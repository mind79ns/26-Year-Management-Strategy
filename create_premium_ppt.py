#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
2026년 경영전략 프리미엄 PPT 생성 스크립트
제조1팀 - 고급 디자인 및 시각화 최적화 버전
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_gradient_background(slide, color1, color2):
    """그라데이션 배경 추가"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 90.0
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def create_rounded_box(slide, left, top, width, height, fill_color, line_color=None, line_width=1):
    """둥근 모서리 박스 생성"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_title_with_underline(slide, text, top, color):
    """언더라인 있는 제목 추가"""
    # 제목 텍스트
    title_box = slide.shapes.add_textbox(Inches(1), top, Inches(8), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = color
    title_para.alignment = PP_ALIGN.LEFT

    # 언더라인
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), top + Inches(0.45), Inches(2), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

    return top + Inches(0.7)

def create_premium_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 색상 팔레트 정의
    NAVY = RGBColor(0, 32, 96)
    BLUE = RGBColor(0, 102, 204)
    LIGHT_BLUE = RGBColor(173, 216, 230)
    ORANGE = RGBColor(255, 127, 39)
    LIGHT_ORANGE = RGBColor(255, 218, 185)
    GREEN = RGBColor(46, 125, 50)
    LIGHT_GREEN = RGBColor(200, 230, 201)
    RED = RGBColor(211, 47, 47)
    GRAY = RGBColor(97, 97, 97)
    LIGHT_GRAY = RGBColor(245, 245, 245)
    WHITE = RGBColor(255, 255, 255)

    # ========== 슬라이드 1: 표지 ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide1, NAVY, BLUE)

    # 메인 타이틀
    main_title = slide1.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(1.2))
    main_frame = main_title.text_frame
    main_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    title_text = main_frame.paragraphs[0]
    title_text.text = "2026년 경영전략"
    title_text.font.size = Pt(66)
    title_text.font.bold = True
    title_text.font.color.rgb = WHITE
    title_text.alignment = PP_ALIGN.CENTER

    # 서브 타이틀
    subtitle = slide1.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(7), Inches(1))
    sub_frame = subtitle.text_frame
    sub_frame.word_wrap = True

    sub_text = sub_frame.paragraphs[0]
    sub_text.text = "자동화 제조라인 스마트화를 통한"
    sub_text.font.size = Pt(26)
    sub_text.font.color.rgb = LIGHT_BLUE
    sub_text.alignment = PP_ALIGN.CENTER

    sub_text2 = sub_frame.add_paragraph()
    sub_text2.text = "가공비 절감 및 품질 개선"
    sub_text2.font.size = Pt(26)
    sub_text2.font.color.rgb = LIGHT_BLUE
    sub_text2.alignment = PP_ALIGN.CENTER

    # 핵심 키워드 박스
    keyword_y = Inches(5)
    keywords = [
        ("가공비", "-10%"),
        ("품질불량", "-10%"),
        ("손실시간", "-5%")
    ]

    box_width = Inches(2.2)
    spacing = Inches(0.3)
    total_width = (box_width * 3) + (spacing * 2)
    start_x = (Inches(10) - total_width) / 2

    for i, (label, value) in enumerate(keywords):
        x_pos = start_x + (i * (box_width + spacing))

        # 박스
        box = create_rounded_box(
            slide1, x_pos, keyword_y, box_width, Inches(0.8),
            WHITE, None
        )

        # 텍스트
        box_frame = box.text_frame
        box_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        label_para = box_frame.paragraphs[0]
        label_para.text = label
        label_para.font.size = Pt(18)
        label_para.font.color.rgb = NAVY
        label_para.alignment = PP_ALIGN.CENTER

        value_para = box_frame.add_paragraph()
        value_para.text = value
        value_para.font.size = Pt(28)
        value_para.font.bold = True
        value_para.font.color.rgb = RED
        value_para.alignment = PP_ALIGN.CENTER

    # 하단 팀 정보
    team_box = slide1.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.4))
    team_frame = team_box.text_frame
    team_para = team_frame.paragraphs[0]
    team_para.text = "제조1팀  |  Manufacturing Team 1"
    team_para.font.size = Pt(16)
    team_para.font.color.rgb = LIGHT_BLUE
    team_para.alignment = PP_ALIGN.CENTER

    # ========== 슬라이드 2: 목표 및 현황 분석 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = WHITE

    # 제목
    current_y = add_title_with_underline(slide2, "목표 및 현황 분석", Inches(0.4), NAVY)

    # 왼쪽: 현재 문제점
    problem_box = create_rounded_box(
        slide2, Inches(0.5), current_y, Inches(4.4), Inches(5.2),
        RGBColor(255, 245, 245), RED, 2
    )

    prob_frame = problem_box.text_frame
    prob_frame.margin_left = Inches(0.3)
    prob_frame.margin_right = Inches(0.3)
    prob_frame.margin_top = Inches(0.25)

    prob_title = prob_frame.paragraphs[0]
    prob_title.text = "⚠️ 현재 문제점"
    prob_title.font.size = Pt(24)
    prob_title.font.bold = True
    prob_title.font.color.rgb = RED
    prob_title.space_after = Pt(15)

    problems = [
        ("1. 느린 대응", "• 설비 이상 발견 → DATA 집계 → 대응", "• 기회 손실 지속 발생", ""),
        ("2. 반복 작업", "• 수동적 원인 분석", "• 업무 비효율 심화", ""),
        ("3. 품질 문제", "• 불량 발생 원인 대응 지연", "• 재발 방지 체계 미흡", "• 현장 체감 저하")
    ]

    for title, *details in problems:
        # 문제 제목
        para = prob_frame.add_paragraph()
        para.text = title
        para.font.size = Pt(16)
        para.font.bold = True
        para.font.color.rgb = GRAY
        para.space_before = Pt(12)
        para.space_after = Pt(6)

        # 세부사항
        for detail in details:
            if detail:
                para = prob_frame.add_paragraph()
                para.text = detail
                para.font.size = Pt(13)
                para.font.color.rgb = GRAY
                para.space_after = Pt(3)
                para.level = 1

    # 오른쪽 상단: 핵심 목표
    goal_box = create_rounded_box(
        slide2, Inches(5.1), current_y, Inches(4.4), Inches(2.4),
        BLUE, None
    )

    goal_frame = goal_box.text_frame
    goal_frame.margin_left = Inches(0.3)
    goal_frame.margin_right = Inches(0.3)
    goal_frame.margin_top = Inches(0.25)
    goal_frame.vertical_anchor = MSO_ANCHOR.TOP

    goal_title = goal_frame.paragraphs[0]
    goal_title.text = "🎯 핵심 목표 (2026)"
    goal_title.font.size = Pt(24)
    goal_title.font.bold = True
    goal_title.font.color.rgb = WHITE
    goal_title.space_after = Pt(12)

    goals = [
        "✓ 가공비 10% 절감",
        "✓ 품질 불량 10% 감소",
        "✓ 손실 시간 5% 감소",
        "✓ OH(간접비) 달성"
    ]

    for goal in goals:
        para = goal_frame.add_paragraph()
        para.text = goal
        para.font.size = Pt(16)
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.space_after = Pt(8)

    # 오른쪽 하단: 기대 효과
    effect_box = create_rounded_box(
        slide2, Inches(5.1), current_y + Inches(2.8), Inches(4.4), Inches(2.4),
        GREEN, None
    )

    effect_frame = effect_box.text_frame
    effect_frame.margin_left = Inches(0.3)
    effect_frame.margin_right = Inches(0.3)
    effect_frame.margin_top = Inches(0.25)

    effect_title = effect_frame.paragraphs[0]
    effect_title.text = "💡 기대 효과"
    effect_title.font.size = Pt(24)
    effect_title.font.bold = True
    effect_title.font.color.rgb = WHITE
    effect_title.space_after = Pt(12)

    effects = [
        "→ 기회손실 최소화",
        "→ 설비 CAPA 증가",
        "→ OVERTIME 감소",
        "→ 점당 가공비 직접 절감"
    ]

    for effect in effects:
        para = effect_frame.add_paragraph()
        para.text = effect
        para.font.size = Pt(16)
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.space_after = Pt(8)

    # 페이지 번호
    page_box = slide2.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page_frame = page_box.text_frame
    page_frame.text = "2"
    page_frame.paragraphs[0].font.size = Pt(14)
    page_frame.paragraphs[0].font.color.rgb = GRAY
    page_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ========== 슬라이드 3: 전략 1 - 손실 시간 제로화 ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = WHITE

    # 헤더 배너
    header = create_rounded_box(
        slide3, Inches(0.5), Inches(0.4), Inches(9), Inches(0.7),
        BLUE, None
    )
    header_frame = header.text_frame
    header_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    header_para = header_frame.paragraphs[0]
    header_para.text = "전략 1️⃣   손실 시간 제로화 프로젝트"
    header_para.font.size = Pt(32)
    header_para.font.bold = True
    header_para.font.color.rgb = WHITE
    header_para.alignment = PP_ALIGN.CENTER

    # 전략 개요
    overview_y = Inches(1.4)
    overview_box = create_rounded_box(
        slide3, Inches(0.5), overview_y, Inches(9), Inches(1.1),
        RGBColor(230, 240, 255), BLUE, 2
    )

    over_frame = overview_box.text_frame
    over_frame.margin_left = Inches(0.3)
    over_frame.margin_top = Inches(0.2)

    over_para = over_frame.paragraphs[0]
    over_para.text = "💡 전략 개요"
    over_para.font.size = Pt(20)
    over_para.font.bold = True
    over_para.font.color.rgb = BLUE
    over_para.space_after = Pt(8)

    over_desc = over_frame.add_paragraph()
    over_desc.text = "MES DATA 활용 자동분석 TOOL 구축 → 순간유실 가시화 → 초 단위 손실 기록 → 우선순위 타겟 집중 개선"
    over_desc.font.size = Pt(15)
    over_desc.font.color.rgb = GRAY

    # 3개 액션 박스
    action_y = Inches(2.8)
    actions = [
        {
            "num": "①",
            "title": "자동분석 TOOL 제작",
            "items": [
                "• MES DATA 활용 이상 감지",
                "• C/T 변화 모니터링 활성화",
                "• 실시간 알람 시스템 구축"
            ]
        },
        {
            "num": "②",
            "title": "순간유실 가시화",
            "items": [
                "• 초 단위 손실 기록",
                "• 모바일 앱 원터치 입력",
                "• LINE별 일/주/월 분석"
            ]
        },
        {
            "num": "③",
            "title": "TOP 10 집중 개선",
            "items": [
                "• 손실 항목 순위화",
                "• WORST 품목 집중 타격",
                "• 주간 리포트 자동화"
            ]
        }
    ]

    action_width = Inches(2.8)
    action_spacing = Inches(0.3)
    action_start = Inches(0.5)

    for i, action in enumerate(actions):
        x_pos = action_start + (i * (action_width + action_spacing))

        # 액션 박스
        action_box = create_rounded_box(
            slide3, x_pos, action_y, action_width, Inches(2.5),
            LIGHT_BLUE, BLUE, 2
        )

        act_frame = action_box.text_frame
        act_frame.margin_left = Inches(0.2)
        act_frame.margin_right = Inches(0.2)
        act_frame.margin_top = Inches(0.2)

        # 번호
        num_para = act_frame.paragraphs[0]
        num_para.text = action["num"]
        num_para.font.size = Pt(32)
        num_para.font.bold = True
        num_para.font.color.rgb = BLUE
        num_para.alignment = PP_ALIGN.CENTER
        num_para.space_after = Pt(8)

        # 제목
        title_para = act_frame.add_paragraph()
        title_para.text = action["title"]
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = NAVY
        title_para.alignment = PP_ALIGN.CENTER
        title_para.space_after = Pt(12)

        # 항목들
        for item in action["items"]:
            item_para = act_frame.add_paragraph()
            item_para.text = item
            item_para.font.size = Pt(12)
            item_para.font.color.rgb = GRAY
            item_para.space_after = Pt(4)

    # KPI 박스
    kpi_box = create_rounded_box(
        slide3, Inches(0.5), Inches(5.6), Inches(9), Inches(1.1),
        RED, None
    )

    kpi_frame = kpi_box.text_frame
    kpi_frame.margin_top = Inches(0.15)
    kpi_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    kpi_title = kpi_frame.paragraphs[0]
    kpi_title.text = "📊 핵심 KPI"
    kpi_title.font.size = Pt(22)
    kpi_title.font.bold = True
    kpi_title.font.color.rgb = WHITE
    kpi_title.alignment = PP_ALIGN.CENTER
    kpi_title.space_after = Pt(8)

    kpi_detail = kpi_frame.add_paragraph()
    kpi_detail.text = "손실시간 5% 감소  |  기회손실 최소화  |  점당 가공비 직접 절감 효과"
    kpi_detail.font.size = Pt(18)
    kpi_detail.font.color.rgb = WHITE
    kpi_detail.alignment = PP_ALIGN.CENTER

    # 페이지 번호
    page_box = slide3.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page_frame = page_box.text_frame
    page_frame.text = "3"
    page_frame.paragraphs[0].font.size = Pt(14)
    page_frame.paragraphs[0].font.color.rgb = GRAY
    page_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ========== 슬라이드 4: 전략 2 - 불량 재발 Zero ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    slide4.background.fill.solid()
    slide4.background.fill.fore_color.rgb = WHITE

    # 헤더 배너
    header = create_rounded_box(
        slide4, Inches(0.5), Inches(0.4), Inches(9), Inches(0.7),
        ORANGE, None
    )
    header_frame = header.text_frame
    header_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    header_para = header_frame.paragraphs[0]
    header_para.text = "전략 2️⃣   불량 재발 Zero 챌린지"
    header_para.font.size = Pt(32)
    header_para.font.bold = True
    header_para.font.color.rgb = WHITE
    header_para.alignment = PP_ALIGN.CENTER

    # 전략 개요
    overview_box = create_rounded_box(
        slide4, Inches(0.5), Inches(1.4), Inches(9), Inches(1.1),
        LIGHT_ORANGE, ORANGE, 2
    )

    over_frame = overview_box.text_frame
    over_frame.margin_left = Inches(0.3)
    over_frame.margin_top = Inches(0.2)

    over_para = over_frame.paragraphs[0]
    over_para.text = "💡 전략 개요"
    over_para.font.size = Pt(20)
    over_para.font.bold = True
    over_para.font.color.rgb = ORANGE
    over_para.space_after = Pt(8)

    over_desc = over_frame.add_paragraph()
    over_desc.text = "불량 사진 즉시 공유 시스템 구축 → 전 조 자동 알람 → 재발 불량 추적 강화 → 조치사항 DB 구축 → 재발 Zero 달성"
    over_desc.font.size = Pt(15)
    over_desc.font.color.rgb = GRAY

    # 3개 액션 박스
    actions = [
        {
            "num": "①",
            "title": "품질 즉시 FEEDBACK",
            "items": [
                "• 불량 사진 + MES 이력 연동",
                "• 모델/일자/LINE/담당자 자동 기록",
                "• 스마트폰 즉시 업로드"
            ]
        },
        {
            "num": "②",
            "title": "전 조 자동 알람",
            "items": [
                "• 공정/설비 태그 자동 분류",
                "• 조치내용 한 줄 메모 공유",
                "• 주간조 → 야간조 자동 전달"
            ]
        },
        {
            "num": "③",
            "title": "재발 불량 추적",
            "items": [
                "• 동일 불량 자동 '재발' 표시",
                "• 월별 재발 불량 추적",
                "• 재발 Zero KPI화"
            ]
        }
    ]

    for i, action in enumerate(actions):
        x_pos = action_start + (i * (action_width + action_spacing))

        action_box = create_rounded_box(
            slide4, x_pos, Inches(2.8), action_width, Inches(2.5),
            LIGHT_ORANGE, ORANGE, 2
        )

        act_frame = action_box.text_frame
        act_frame.margin_left = Inches(0.2)
        act_frame.margin_right = Inches(0.2)
        act_frame.margin_top = Inches(0.2)

        num_para = act_frame.paragraphs[0]
        num_para.text = action["num"]
        num_para.font.size = Pt(32)
        num_para.font.bold = True
        num_para.font.color.rgb = ORANGE
        num_para.alignment = PP_ALIGN.CENTER
        num_para.space_after = Pt(8)

        title_para = act_frame.add_paragraph()
        title_para.text = action["title"]
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(139, 69, 19)
        title_para.alignment = PP_ALIGN.CENTER
        title_para.space_after = Pt(12)

        for item in action["items"]:
            item_para = act_frame.add_paragraph()
            item_para.text = item
            item_para.font.size = Pt(12)
            item_para.font.color.rgb = GRAY
            item_para.space_after = Pt(4)

    # KPI 박스
    kpi_box = create_rounded_box(
        slide4, Inches(0.5), Inches(5.6), Inches(9), Inches(1.1),
        RED, None
    )

    kpi_frame = kpi_box.text_frame
    kpi_frame.margin_top = Inches(0.15)
    kpi_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    kpi_title = kpi_frame.paragraphs[0]
    kpi_title.text = "📊 핵심 KPI"
    kpi_title.font.size = Pt(22)
    kpi_title.font.bold = True
    kpi_title.font.color.rgb = WHITE
    kpi_title.alignment = PP_ALIGN.CENTER
    kpi_title.space_after = Pt(8)

    kpi_detail = kpi_frame.add_paragraph()
    kpi_detail.text = "재발 불량 30% 감소  |  품질 불량 10% 감소  |  조치사항 DB 100% 구축"
    kpi_detail.font.size = Pt(18)
    kpi_detail.font.color.rgb = WHITE
    kpi_detail.alignment = PP_ALIGN.CENTER

    # 페이지 번호
    page_box = slide4.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page_frame = page_box.text_frame
    page_frame.text = "4"
    page_frame.paragraphs[0].font.size = Pt(14)
    page_frame.paragraphs[0].font.color.rgb = GRAY
    page_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ========== 슬라이드 5: 전략 3 - 설비 CAPA 증가 ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    slide5.background.fill.solid()
    slide5.background.fill.fore_color.rgb = WHITE

    # 헤더 배너
    header = create_rounded_box(
        slide5, Inches(0.5), Inches(0.4), Inches(9), Inches(0.7),
        GREEN, None
    )
    header_frame = header.text_frame
    header_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    header_para = header_frame.paragraphs[0]
    header_para.text = "전략 3️⃣   설비 CAPA 증가 및 공정 최적화"
    header_para.font.size = Pt(32)
    header_para.font.bold = True
    header_para.font.color.rgb = WHITE
    header_para.alignment = PP_ALIGN.CENTER

    # 전략 개요
    overview_box = create_rounded_box(
        slide5, Inches(0.5), Inches(1.4), Inches(9), Inches(1.1),
        LIGHT_GREEN, GREEN, 2
    )

    over_frame = overview_box.text_frame
    over_frame.margin_left = Inches(0.3)
    over_frame.margin_top = Inches(0.2)

    over_para = over_frame.paragraphs[0]
    over_para.text = "💡 전략 개요"
    over_para.font.size = Pt(20)
    over_para.font.bold = True
    over_para.font.color.rgb = GREEN
    over_para.space_after = Pt(8)

    over_desc = over_frame.add_paragraph()
    over_desc.text = "C/T 단축 DEEP 분석 → 단 1초라도 줄이기 위한 활동 → 설비 효율 성능/PM 지속 개선 → 공정 최적화"
    over_desc.font.size = Pt(15)
    over_desc.font.color.rgb = GRAY

    # 3개 액션 박스
    actions = [
        {
            "num": "①",
            "title": "C/T 단축 개선",
            "items": [
                "• 현 P/G 운영 DEEP 분석",
                "• 최단거리 프로세스 설계",
                "• 단 1초라도 줄이기 활동",
                "• 병목 공정 표적 개선"
            ]
        },
        {
            "num": "②",
            "title": "기본 BASE 강화",
            "items": [
                "• 설비 효율 유지 및 향상",
                "• 성능 관련 모든 활동 반복",
                "• PM(예방정비) 체계 강화",
                "• 지속적 모니터링"
            ]
        },
        {
            "num": "③",
            "title": "공정 최적화",
            "items": [
                "• RADIAL2 수삽설비 안정화",
                "• SMD 공정 LAY OUT 개선",
                "• 최적 운영 방안 검토",
                "• 설비 배치 효율화"
            ]
        }
    ]

    for i, action in enumerate(actions):
        x_pos = action_start + (i * (action_width + action_spacing))

        action_box = create_rounded_box(
            slide5, x_pos, Inches(2.8), action_width, Inches(2.5),
            LIGHT_GREEN, GREEN, 2
        )

        act_frame = action_box.text_frame
        act_frame.margin_left = Inches(0.2)
        act_frame.margin_right = Inches(0.2)
        act_frame.margin_top = Inches(0.2)

        num_para = act_frame.paragraphs[0]
        num_para.text = action["num"]
        num_para.font.size = Pt(32)
        num_para.font.bold = True
        num_para.font.color.rgb = GREEN
        num_para.alignment = PP_ALIGN.CENTER
        num_para.space_after = Pt(8)

        title_para = act_frame.add_paragraph()
        title_para.text = action["title"]
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(27, 94, 32)
        title_para.alignment = PP_ALIGN.CENTER
        title_para.space_after = Pt(12)

        for item in action["items"]:
            item_para = act_frame.add_paragraph()
            item_para.text = item
            item_para.font.size = Pt(11)
            item_para.font.color.rgb = GRAY
            item_para.space_after = Pt(3)

    # KPI 박스
    kpi_box = create_rounded_box(
        slide5, Inches(0.5), Inches(5.6), Inches(9), Inches(1.1),
        RED, None
    )

    kpi_frame = kpi_box.text_frame
    kpi_frame.margin_top = Inches(0.15)
    kpi_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    kpi_title = kpi_frame.paragraphs[0]
    kpi_title.text = "📊 핵심 KPI"
    kpi_title.font.size = Pt(22)
    kpi_title.font.bold = True
    kpi_title.font.color.rgb = WHITE
    kpi_title.alignment = PP_ALIGN.CENTER
    kpi_title.space_after = Pt(8)

    kpi_detail = kpi_frame.add_paragraph()
    kpi_detail.text = "설비 CAPA 증가  |  OVERTIME 감소  |  OH(간접비) 감소 달성"
    kpi_detail.font.size = Pt(18)
    kpi_detail.font.color.rgb = WHITE
    kpi_detail.alignment = PP_ALIGN.CENTER

    # 페이지 번호
    page_box = slide5.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page_frame = page_box.text_frame
    page_frame.text = "5"
    page_frame.paragraphs[0].font.size = Pt(14)
    page_frame.paragraphs[0].font.color.rgb = GRAY
    page_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ========== 슬라이드 6: 실행 로드맵 및 종합 ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    slide6.background.fill.solid()
    slide6.background.fill.fore_color.rgb = WHITE

    # 제목
    current_y = add_title_with_underline(slide6, "실행 로드맵 및 종합 KPI", Inches(0.4), NAVY)

    # 타임라인 박스들
    timeline_y = Inches(1.3)
    quarters = [
        {"q": "Q1", "color": RGBColor(100, 181, 246), "tasks": [
            "• 자동분석 TOOL 개발",
            "• 모바일 앱 구축",
            "• 불량 공유 시스템 개발"
        ]},
        {"q": "Q2", "color": RGBColor(66, 165, 245), "tasks": [
            "• 시범 운영 및 피드백",
            "• 전 조 확대 적용",
            "• C/T 분석 시작"
        ]},
        {"q": "Q3", "color": RGBColor(42, 149, 224), "tasks": [
            "• 전사 확대 운영",
            "• 재발 불량 추적 강화",
            "• 공정 최적화 실행"
        ]},
        {"q": "Q4", "color": RGBColor(25, 118, 210), "tasks": [
            "• 목표 달성 점검",
            "• 우수 사례 공유",
            "• 차년도 계획 수립"
        ]}
    ]

    quarter_width = Inches(2.1)
    quarter_spacing = Inches(0.2)
    quarter_start = Inches(0.5)

    for i, quarter in enumerate(quarters):
        x_pos = quarter_start + (i * (quarter_width + quarter_spacing))

        # 분기 박스
        q_box = create_rounded_box(
            slide6, x_pos, timeline_y, quarter_width, Inches(2.2),
            quarter["color"], None
        )

        q_frame = q_box.text_frame
        q_frame.margin_left = Inches(0.15)
        q_frame.margin_right = Inches(0.15)
        q_frame.margin_top = Inches(0.15)

        # 분기 제목
        q_title = q_frame.paragraphs[0]
        q_title.text = quarter["q"]
        q_title.font.size = Pt(28)
        q_title.font.bold = True
        q_title.font.color.rgb = WHITE
        q_title.alignment = PP_ALIGN.CENTER
        q_title.space_after = Pt(10)

        # 과제들
        for task in quarter["tasks"]:
            task_para = q_frame.add_paragraph()
            task_para.text = task
            task_para.font.size = Pt(11)
            task_para.font.color.rgb = WHITE
            task_para.space_after = Pt(5)

    # 종합 KPI 영역
    kpi_area_y = Inches(3.8)

    # KPI 제목
    kpi_title_box = slide6.shapes.add_textbox(Inches(0.5), kpi_area_y, Inches(9), Inches(0.4))
    kpi_title_frame = kpi_title_box.text_frame
    kpi_title_para = kpi_title_frame.paragraphs[0]
    kpi_title_para.text = "📊 종합 KPI 대시보드"
    kpi_title_para.font.size = Pt(24)
    kpi_title_para.font.bold = True
    kpi_title_para.font.color.rgb = NAVY
    kpi_title_para.alignment = PP_ALIGN.CENTER

    # 3개 주요 KPI 박스
    kpi_y = kpi_area_y + Inches(0.6)
    kpis = [
        {"label": "가공비 절감", "value": "-10%", "color": BLUE},
        {"label": "품질 개선", "value": "-10%", "color": ORANGE},
        {"label": "손실 시간", "value": "-5%", "color": GREEN}
    ]

    kpi_box_width = Inches(2.8)
    kpi_spacing = Inches(0.3)
    kpi_start = Inches(0.5)

    for i, kpi in enumerate(kpis):
        x_pos = kpi_start + (i * (kpi_box_width + kpi_spacing))

        # KPI 박스
        kpi_box = create_rounded_box(
            slide6, x_pos, kpi_y, kpi_box_width, Inches(1.4),
            kpi["color"], None
        )

        kpi_frame = kpi_box.text_frame
        kpi_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 라벨
        label_para = kpi_frame.paragraphs[0]
        label_para.text = kpi["label"]
        label_para.font.size = Pt(18)
        label_para.font.color.rgb = WHITE
        label_para.alignment = PP_ALIGN.CENTER
        label_para.space_after = Pt(8)

        # 값
        value_para = kpi_frame.add_paragraph()
        value_para.text = kpi["value"]
        value_para.font.size = Pt(42)
        value_para.font.bold = True
        value_para.font.color.rgb = WHITE
        value_para.alignment = PP_ALIGN.CENTER

    # 하단 종합 효과
    effect_y = kpi_y + Inches(1.7)
    effect_box = create_rounded_box(
        slide6, Inches(0.5), effect_y, Inches(9), Inches(1.2),
        RGBColor(240, 240, 240), GRAY, 2
    )

    effect_frame = effect_box.text_frame
    effect_frame.margin_left = Inches(0.3)
    effect_frame.margin_top = Inches(0.2)

    effect_title = effect_frame.paragraphs[0]
    effect_title.text = "🎯 최종 기대 효과"
    effect_title.font.size = Pt(20)
    effect_title.font.bold = True
    effect_title.font.color.rgb = NAVY
    effect_title.space_after = Pt(10)

    effects_list = "기회손실 최소화   |   설비 CAPA 증가   |   OVERTIME 감소   |   점당 가공비 직접 절감   |   OH 달성"
    effect_detail = effect_frame.add_paragraph()
    effect_detail.text = effects_list
    effect_detail.font.size = Pt(16)
    effect_detail.font.color.rgb = GRAY
    effect_detail.alignment = PP_ALIGN.CENTER

    # 페이지 번호
    page_box = slide6.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page_frame = page_box.text_frame
    page_frame.text = "6"
    page_frame.paragraphs[0].font.size = Pt(14)
    page_frame.paragraphs[0].font.color.rgb = GRAY
    page_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 저장
    output_file = '2026년_경영전략_제조1팀_프리미엄.pptx'
    prs.save(output_file)
    print(f"✅ 프리미엄 PPT 파일이 성공적으로 생성되었습니다: {output_file}")
    print(f"📄 총 슬라이드 수: {len(prs.slides)}")
    print(f"🎨 페이지 구성:")
    print(f"   1. 표지 - 임팩트 있는 메인 비주얼")
    print(f"   2. 목표 및 현황 분석 - 문제점/목표/효과 구조화")
    print(f"   3. 전략 1 - 손실 시간 제로화 프로젝트")
    print(f"   4. 전략 2 - 불량 재발 Zero 챌린지")
    print(f"   5. 전략 3 - 설비 CAPA 증가")
    print(f"   6. 실행 로드맵 및 종합 KPI")
    return output_file

if __name__ == "__main__":
    create_premium_presentation()
