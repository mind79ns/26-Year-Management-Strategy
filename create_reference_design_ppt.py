#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
참고디자인 기반 PPT 생성
프로그레스 바 & 원형 차트 중심의 모던 클린 디자인
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
import math

def create_progress_bar(slide, x, y, width, height, percentage, color, bg_color):
    """프로그레스 바 생성 (참고디자인 스타일)"""
    # 배경 바
    bg_bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x, y, width, height
    )
    bg_bar.fill.solid()
    bg_bar.fill.fore_color.rgb = bg_color
    bg_bar.line.fill.background()

    # 진행 바
    progress_width = width * (percentage / 100)
    if progress_width > Inches(0.1):  # 최소 너비 체크
        progress_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, progress_width, height
        )
        progress_bar.fill.solid()
        progress_bar.fill.fore_color.rgb = color
        progress_bar.line.fill.background()

    # 퍼센티지 원형 라벨
    circle_size = Inches(0.5)
    circle_x = x + progress_width - circle_size / 2
    circle_y = y - Inches(0.1)

    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        circle_x, circle_y, circle_size, circle_size
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()

    # 퍼센티지 텍스트
    percent_box = slide.shapes.add_textbox(
        circle_x, circle_y, circle_size, circle_size
    )
    percent_frame = percent_box.text_frame
    percent_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    percent_para = percent_frame.paragraphs[0]
    percent_para.text = f"{int(percentage)}%"
    percent_para.font.size = Pt(11)
    percent_para.font.bold = True
    percent_para.font.color.rgb = RGBColor(255, 255, 255)
    percent_para.alignment = PP_ALIGN.CENTER

def create_circular_gauge(slide, center_x, center_y, radius, percentage, color):
    """원형 게이지 생성 (참고디자인 스타일)"""
    # 배경 원
    bg_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        center_x - radius, center_y - radius,
        radius * 2, radius * 2
    )
    bg_circle.fill.solid()
    bg_circle.fill.fore_color.rgb = RGBColor(240, 240, 240)
    bg_circle.line.color.rgb = RGBColor(220, 220, 220)
    bg_circle.line.width = Pt(2)

    # 프로그레스 원 (아크 효과를 위해 여러 개의 작은 원 사용)
    # 실제 아크는 python-pptx에서 직접 지원하지 않으므로 텍스트로 표현

    # 중앙 원 (흰색)
    inner_radius = radius * 0.7
    inner_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        center_x - inner_radius, center_y - inner_radius,
        inner_radius * 2, inner_radius * 2
    )
    inner_circle.fill.solid()
    inner_circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
    inner_circle.line.fill.background()

    # 퍼센티지 텍스트
    percent_box = slide.shapes.add_textbox(
        center_x - inner_radius, center_y - inner_radius,
        inner_radius * 2, inner_radius * 2
    )
    percent_frame = percent_box.text_frame
    percent_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    percent_para = percent_frame.paragraphs[0]
    percent_para.text = f"{int(percentage)}%"
    percent_para.font.size = Pt(36)
    percent_para.font.bold = True
    percent_para.font.color.rgb = color
    percent_para.alignment = PP_ALIGN.CENTER

def create_info_box(slide, x, y, width, height, number, title, description, color):
    """정보 박스 생성 (참고디자인 스타일)"""
    # 박스
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x, y, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()

    # 텍스트 프레임
    frame = box.text_frame
    frame.margin_left = Inches(0.25)
    frame.margin_right = Inches(0.25)
    frame.margin_top = Inches(0.15)
    frame.word_wrap = True

    # 번호와 제목
    title_para = frame.paragraphs[0]
    title_para.text = f"{number} {title}"
    title_para.font.size = Pt(14)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.space_after = Pt(8)

    # 설명
    desc_para = frame.add_paragraph()
    desc_para.text = description
    desc_para.font.size = Pt(11)
    desc_para.font.color.rgb = RGBColor(255, 255, 255)

def create_reference_design_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 참고디자인 색상 팔레트
    PRIMARY_BLUE = RGBColor(41, 128, 185)      # 주요 블루
    LIGHT_BLUE = RGBColor(93, 173, 226)        # 밝은 블루
    DARK_BLUE = RGBColor(21, 67, 96)           # 다크 블루
    NAVY = RGBColor(44, 62, 80)                # 네이비
    GRAY = RGBColor(149, 165, 166)             # 그레이
    LIGHT_GRAY = RGBColor(220, 220, 220)       # 라이트 그레이
    BG_GRAY = RGBColor(245, 245, 245)          # 배경 그레이
    WHITE = RGBColor(255, 255, 255)

    # ========== 슬라이드 1: 표지 (참고디자인 스타일) ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = WHITE

    # 상단 블루 바
    top_bar = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(0.15)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PRIMARY_BLUE
    top_bar.line.fill.background()

    # 타이틀
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "2026년 경영전략"
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_BLUE
    title_para.alignment = PP_ALIGN.CENTER

    # 서브타이틀
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.6), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = "제조1팀 | 자동화 제조라인 스마트화"
    subtitle_para.font.size = Pt(22)
    subtitle_para.font.color.rgb = GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER

    # 3개 핵심 지표 (프로그레스 바 미리보기)
    indicators_y = Inches(5)
    indicators = [
        {"label": "가공비 절감", "value": 10, "color": PRIMARY_BLUE},
        {"label": "품질 개선", "value": 10, "color": LIGHT_BLUE},
        {"label": "손실시간 감소", "value": 5, "color": NAVY}
    ]

    ind_width = Inches(2.2)
    ind_spacing = Inches(0.4)
    ind_start = (Inches(10) - (ind_width * 3 + ind_spacing * 2)) / 2

    for i, ind in enumerate(indicators):
        x_pos = ind_start + i * (ind_width + ind_spacing)

        # 라벨
        label_box = slide1.shapes.add_textbox(x_pos, indicators_y, ind_width, Inches(0.3))
        label_frame = label_box.text_frame
        label_para = label_frame.paragraphs[0]
        label_para.text = ind["label"]
        label_para.font.size = Pt(12)
        label_para.font.color.rgb = DARK_BLUE
        label_para.alignment = PP_ALIGN.CENTER

        # 미니 프로그레스 바
        create_progress_bar(
            slide1,
            x_pos + Inches(0.3), indicators_y + Inches(0.4),
            ind_width - Inches(0.6), Inches(0.3),
            ind["value"] * 10, ind["color"], LIGHT_GRAY
        )

        # 값
        value_box = slide1.shapes.add_textbox(x_pos, indicators_y + Inches(0.85), ind_width, Inches(0.4))
        value_frame = value_box.text_frame
        value_para = value_frame.paragraphs[0]
        value_para.text = f"{ind['value']}% 목표"
        value_para.font.size = Pt(14)
        value_para.font.bold = True
        value_para.font.color.rgb = ind["color"]
        value_para.alignment = PP_ALIGN.CENTER

    # ========== 슬라이드 2: 목표 및 현황 (참고디자인 스타일) ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = WHITE

    # 상단 바
    top_bar2 = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(0.15)
    )
    top_bar2.fill.solid()
    top_bar2.fill.fore_color.rgb = PRIMARY_BLUE
    top_bar2.line.fill.background()

    # 제목
    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title2_frame = title2.text_frame
    title2_para = title2_frame.paragraphs[0]
    title2_para.text = "1. 목표 및 현황 분석"
    title2_para.font.size = Pt(28)
    title2_para.font.bold = True
    title2_para.font.color.rgb = DARK_BLUE

    # 좌측: 현재 문제점 (프로그레스 바 형태)
    problems_x = Inches(0.5)
    problems_y = Inches(1.3)

    problem_label = slide2.shapes.add_textbox(problems_x, problems_y, Inches(4.5), Inches(0.4))
    pl_frame = problem_label.text_frame
    pl_para = pl_frame.paragraphs[0]
    pl_para.text = "이곳에 제품을 낳는 근원다."
    pl_para.font.size = Pt(16)
    pl_para.font.bold = True
    pl_para.font.color.rgb = DARK_BLUE

    # 문제점 프로그레스 바들
    problems = [
        {"label": "느린 대응", "severity": 60, "color": PRIMARY_BLUE, "desc": "설비 이상 DATA 집계 대응으로 기회 손실 발생"},
        {"label": "반복 작업", "severity": 35, "color": GRAY, "desc": "수동적 원인 분석으로 업무 비효율 심화"},
        {"label": "품질 문제", "severity": 48, "color": DARK_BLUE, "desc": "불량 원인 대응 지연 및 재발 방지 체계 미흡"},
        {"label": "설비 효율", "severity": 50, "color": LIGHT_BLUE, "desc": "인력 한정으로 운영 LINE C/T 개선 필요"}
    ]

    bar_y = problems_y + Inches(0.6)
    bar_spacing = Inches(0.85)

    for i, prob in enumerate(problems):
        current_y = bar_y + i * bar_spacing

        # 프로그레스 바
        create_progress_bar(
            slide2,
            problems_x, current_y,
            Inches(4.5), Inches(0.25),
            prob["severity"], prob["color"], LIGHT_GRAY
        )

        # 설명 박스
        desc_y = current_y + Inches(0.35)
        create_info_box(
            slide2,
            problems_x, desc_y,
            Inches(4.5), Inches(0.35),
            f"0{i+1}", prob["label"], prob["desc"],
            DARK_BLUE
        )

    # 우측: 원형 게이지 + 목표
    gauge_center_x = Inches(7.5)
    gauge_center_y = Inches(2.5)
    gauge_radius = Inches(0.9)

    # "이곳에는" 타이틀
    goal_title = slide2.shapes.add_textbox(Inches(5.5), Inches(1.3), Inches(4), Inches(0.4))
    gt_frame = goal_title.text_frame
    gt_para = gt_frame.paragraphs[0]
    gt_para.text = "이곳에는"
    gt_para.font.size = Pt(16)
    gt_para.font.bold = True
    gt_para.font.color.rgb = DARK_BLUE
    gt_para.alignment = PP_ALIGN.CENTER

    # 서브 타이틀
    goal_sub = slide2.shapes.add_textbox(Inches(5.5), Inches(1.7), Inches(4), Inches(0.5))
    gs_frame = goal_sub.text_frame
    gs_para = gs_frame.paragraphs[0]
    gs_para.text = "핵심목표를\n입력하세요!"
    gs_para.font.size = Pt(14)
    gs_para.font.color.rgb = PRIMARY_BLUE
    gs_para.alignment = PP_ALIGN.CENTER

    # 원형 게이지 (60% - 종합 목표 달성률)
    create_circular_gauge(
        slide2,
        gauge_center_x, gauge_center_y,
        gauge_radius, 60, PRIMARY_BLUE
    )

    # 목표 정보 박스들
    goals_y = Inches(4.2)
    goals = [
        {"num": "01", "title": "핵심 전략", "desc": "MES DATA 활용 자동분석 TOOL → 순간유실 가시화 → 손실 5% 감소"},
        {"num": "02", "title": "핵심 전략", "desc": "불량 사진 즉시 공유 시스템 → 재발 불량 추적 → 품질 10% 개선"}
    ]

    for i, goal in enumerate(goals):
        create_info_box(
            slide2,
            Inches(5.5), goals_y + i * Inches(0.85),
            Inches(4), Inches(0.7),
            goal["num"], goal["title"], goal["desc"],
            PRIMARY_BLUE
        )

    # 페이지 번호
    page2 = slide2.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page2.text_frame.text = "02"
    page2.text_frame.paragraphs[0].font.size = Pt(11)
    page2.text_frame.paragraphs[0].font.color.rgb = GRAY
    page2.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # ========== 슬라이드 3: 전략 1 (프로그레스 바 중심) ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = WHITE

    # 상단 바
    top_bar3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    top_bar3.fill.solid()
    top_bar3.fill.fore_color.rgb = PRIMARY_BLUE
    top_bar3.line.fill.background()

    # 제목
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title3_frame = title3.text_frame
    title3_para = title3_frame.paragraphs[0]
    title3_para.text = "2. 전략 1 : 손실 시간 제로화 프로젝트"
    title3_para.font.size = Pt(28)
    title3_para.font.bold = True
    title3_para.font.color.rgb = DARK_BLUE

    # 전략 설명
    desc3 = slide3.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.4))
    desc3_frame = desc3.text_frame
    desc3_para = desc3_frame.paragraphs[0]
    desc3_para.text = "MES DATA 활용 자동분석 → 순간유실 가시화 → TOP 10 집중 개선"
    desc3_para.font.size = Pt(14)
    desc3_para.font.color.rgb = GRAY

    # 프로그레스 기반 액션 플랜
    actions_y = Inches(1.8)
    actions = [
        {"title": "자동분석 TOOL 제작", "progress": 100, "color": PRIMARY_BLUE,
         "desc": "• MES DATA 활용 이상 감지\n• C/T 변화 모니터링 활성화\n• 실시간 알람 시스템"},
        {"title": "순간유실 가시화", "progress": 75, "color": LIGHT_BLUE,
         "desc": "• 초 단위 손실 기록\n• 모바일 앱 원터치 입력\n• LINE별 일/주/월 분석"},
        {"title": "TOP 10 집중 개선", "progress": 50, "color": NAVY,
         "desc": "• 손실 항목 순위화\n• WORST 품목 집중 타격\n• 주간 리포트 자동화"}
    ]

    action_width = Inches(2.8)
    action_spacing = Inches(0.3)
    action_start = Inches(0.5)

    for i, action in enumerate(actions):
        x_pos = action_start + i * (action_width + action_spacing)

        # 카드 배경
        card = slide3.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, actions_y, action_width, Inches(4.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = BG_GRAY
        card.line.color.rgb = LIGHT_GRAY
        card.line.width = Pt(1)

        # 타이틀
        title_box = slide3.shapes.add_textbox(x_pos + Inches(0.2), actions_y + Inches(0.2),
                                               action_width - Inches(0.4), Inches(0.5))
        t_frame = title_box.text_frame
        t_frame.word_wrap = True
        t_para = t_frame.paragraphs[0]
        t_para.text = action["title"]
        t_para.font.size = Pt(16)
        t_para.font.bold = True
        t_para.font.color.rgb = DARK_BLUE
        t_para.alignment = PP_ALIGN.CENTER

        # 프로그레스 바
        prog_y = actions_y + Inches(0.8)
        create_progress_bar(
            slide3,
            x_pos + Inches(0.3), prog_y,
            action_width - Inches(0.6), Inches(0.25),
            action["progress"], action["color"], LIGHT_GRAY
        )

        # 설명
        desc_box = slide3.shapes.add_textbox(x_pos + Inches(0.2), actions_y + Inches(1.4),
                                              action_width - Inches(0.4), Inches(2.5))
        d_frame = desc_box.text_frame
        d_frame.word_wrap = True
        d_para = d_frame.paragraphs[0]
        d_para.text = action["desc"]
        d_para.font.size = Pt(11)
        d_para.font.color.rgb = DARK_BLUE

    # KPI 박스
    kpi3 = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(6.3), Inches(9), Inches(0.7)
    )
    kpi3.fill.solid()
    kpi3.fill.fore_color.rgb = PRIMARY_BLUE
    kpi3.line.fill.background()

    kpi3_frame = kpi3.text_frame
    kpi3_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    kpi3_para = kpi3_frame.paragraphs[0]
    kpi3_para.text = "🎯 목표 KPI: 손실시간 5% 감소 → 점당 가공비 직접 절감 효과"
    kpi3_para.font.size = Pt(18)
    kpi3_para.font.bold = True
    kpi3_para.font.color.rgb = WHITE
    kpi3_para.alignment = PP_ALIGN.CENTER

    page3 = slide3.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page3.text_frame.text = "03"
    page3.text_frame.paragraphs[0].font.size = Pt(11)
    page3.text_frame.paragraphs[0].font.color.rgb = GRAY
    page3.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # ========== 슬라이드 4: 전략 2 ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    slide4.background.fill.solid()
    slide4.background.fill.fore_color.rgb = WHITE

    top_bar4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    top_bar4.fill.solid()
    top_bar4.fill.fore_color.rgb = LIGHT_BLUE
    top_bar4.line.fill.background()

    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title4_frame = title4.text_frame
    title4_para = title4_frame.paragraphs[0]
    title4_para.text = "3. 전략 2 : 불량 재발 Zero 챌린지"
    title4_para.font.size = Pt(28)
    title4_para.font.bold = True
    title4_para.font.color.rgb = DARK_BLUE

    desc4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.4))
    desc4_frame = desc4.text_frame
    desc4_para = desc4_frame.paragraphs[0]
    desc4_para.text = "불량 사진 즉시 공유 → 전 조 자동 알람 → 재발 불량 추적 → 재발 Zero"
    desc4_para.font.size = Pt(14)
    desc4_para.font.color.rgb = GRAY

    actions4 = [
        {"title": "품질 즉시 FEEDBACK", "progress": 100, "color": LIGHT_BLUE,
         "desc": "• 불량 사진 + MES 이력\n• 모델/일자/LINE 자동 기록\n• 스마트폰 즉시 업로드"},
        {"title": "전 조 자동 알람", "progress": 80, "color": PRIMARY_BLUE,
         "desc": "• 공정/설비 태그 분류\n• 조치내용 메모 공유\n• 주간조→야간조 전달"},
        {"title": "재발 불량 추적", "progress": 60, "color": DARK_BLUE,
         "desc": "• 동일 불량 '재발' 표시\n• 월별 재발 추적\n• 재발 Zero KPI화"}
    ]

    for i, action in enumerate(actions4):
        x_pos = action_start + i * (action_width + action_spacing)

        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, actions_y,
                                        action_width, Inches(4.2))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_GRAY
        card.line.color.rgb = LIGHT_GRAY
        card.line.width = Pt(1)

        title_box = slide4.shapes.add_textbox(x_pos + Inches(0.2), actions_y + Inches(0.2),
                                               action_width - Inches(0.4), Inches(0.5))
        t_frame = title_box.text_frame
        t_frame.word_wrap = True
        t_para = t_frame.paragraphs[0]
        t_para.text = action["title"]
        t_para.font.size = Pt(16)
        t_para.font.bold = True
        t_para.font.color.rgb = DARK_BLUE
        t_para.alignment = PP_ALIGN.CENTER

        prog_y = actions_y + Inches(0.8)
        create_progress_bar(slide4, x_pos + Inches(0.3), prog_y,
                           action_width - Inches(0.6), Inches(0.25),
                           action["progress"], action["color"], LIGHT_GRAY)

        desc_box = slide4.shapes.add_textbox(x_pos + Inches(0.2), actions_y + Inches(1.4),
                                              action_width - Inches(0.4), Inches(2.5))
        d_frame = desc_box.text_frame
        d_frame.word_wrap = True
        d_para = d_frame.paragraphs[0]
        d_para.text = action["desc"]
        d_para.font.size = Pt(11)
        d_para.font.color.rgb = DARK_BLUE

    kpi4 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.3),
                                    Inches(9), Inches(0.7))
    kpi4.fill.solid()
    kpi4.fill.fore_color.rgb = LIGHT_BLUE
    kpi4.line.fill.background()

    kpi4_frame = kpi4.text_frame
    kpi4_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    kpi4_para = kpi4_frame.paragraphs[0]
    kpi4_para.text = "🎯 목표 KPI: 재발 불량 30% 감소 / 품질 불량 10% 감소"
    kpi4_para.font.size = Pt(18)
    kpi4_para.font.bold = True
    kpi4_para.font.color.rgb = WHITE
    kpi4_para.alignment = PP_ALIGN.CENTER

    page4 = slide4.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page4.text_frame.text = "04"
    page4.text_frame.paragraphs[0].font.size = Pt(11)
    page4.text_frame.paragraphs[0].font.color.rgb = GRAY
    page4.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # ========== 슬라이드 5: 전략 3 ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    slide5.background.fill.solid()
    slide5.background.fill.fore_color.rgb = WHITE

    top_bar5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    top_bar5.fill.solid()
    top_bar5.fill.fore_color.rgb = NAVY
    top_bar5.line.fill.background()

    title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title5_frame = title5.text_frame
    title5_para = title5_frame.paragraphs[0]
    title5_para.text = "4. 전략 3 : 설비 CAPA 증가 및 공정 최적화"
    title5_para.font.size = Pt(28)
    title5_para.font.bold = True
    title5_para.font.color.rgb = DARK_BLUE

    desc5 = slide5.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.4))
    desc5_frame = desc5.text_frame
    desc5_para = desc5_frame.paragraphs[0]
    desc5_para.text = "C/T 단축 DEEP 분석 → 설비 효율 향상 → 공정 최적화"
    desc5_para.font.size = Pt(14)
    desc5_para.font.color.rgb = GRAY

    actions5 = [
        {"title": "C/T 단축 개선", "progress": 90, "color": NAVY,
         "desc": "• P/G 운영 DEEP 분석\n• 최단거리 프로세스\n• 단 1초라도 줄이기"},
        {"title": "기본 BASE 강화", "progress": 85, "color": DARK_BLUE,
         "desc": "• 설비 효율 유지\n• 성능/PM 활동 반복\n• 지속적 모니터링"},
        {"title": "공정 최적화", "progress": 70, "color": PRIMARY_BLUE,
         "desc": "• RADIAL2 안정화\n• SMD LAY OUT 개선\n• 설비 배치 효율화"}
    ]

    for i, action in enumerate(actions5):
        x_pos = action_start + i * (action_width + action_spacing)

        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, actions_y,
                                        action_width, Inches(4.2))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_GRAY
        card.line.color.rgb = LIGHT_GRAY
        card.line.width = Pt(1)

        title_box = slide5.shapes.add_textbox(x_pos + Inches(0.2), actions_y + Inches(0.2),
                                               action_width - Inches(0.4), Inches(0.5))
        t_frame = title_box.text_frame
        t_frame.word_wrap = True
        t_para = t_frame.paragraphs[0]
        t_para.text = action["title"]
        t_para.font.size = Pt(16)
        t_para.font.bold = True
        t_para.font.color.rgb = DARK_BLUE
        t_para.alignment = PP_ALIGN.CENTER

        prog_y = actions_y + Inches(0.8)
        create_progress_bar(slide5, x_pos + Inches(0.3), prog_y,
                           action_width - Inches(0.6), Inches(0.25),
                           action["progress"], action["color"], LIGHT_GRAY)

        desc_box = slide5.shapes.add_textbox(x_pos + Inches(0.2), actions_y + Inches(1.4),
                                              action_width - Inches(0.4), Inches(2.5))
        d_frame = desc_box.text_frame
        d_frame.word_wrap = True
        d_para = d_frame.paragraphs[0]
        d_para.text = action["desc"]
        d_para.font.size = Pt(11)
        d_para.font.color.rgb = DARK_BLUE

    kpi5 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.3),
                                    Inches(9), Inches(0.7))
    kpi5.fill.solid()
    kpi5.fill.fore_color.rgb = NAVY
    kpi5.line.fill.background()

    kpi5_frame = kpi5.text_frame
    kpi5_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    kpi5_para = kpi5_frame.paragraphs[0]
    kpi5_para.text = "🎯 목표 KPI: CAPA 증가 / OVERTIME 감소 / OH 달성"
    kpi5_para.font.size = Pt(18)
    kpi5_para.font.bold = True
    kpi5_para.font.color.rgb = WHITE
    kpi5_para.alignment = PP_ALIGN.CENTER

    page5 = slide5.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page5.text_frame.text = "05"
    page5.text_frame.paragraphs[0].font.size = Pt(11)
    page5.text_frame.paragraphs[0].font.color.rgb = GRAY
    page5.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # ========== 슬라이드 6: 종합 로드맵 (원형 게이지 중심) ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    slide6.background.fill.solid()
    slide6.background.fill.fore_color.rgb = WHITE

    top_bar6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    top_bar6.fill.solid()
    top_bar6.fill.fore_color.rgb = PRIMARY_BLUE
    top_bar6.line.fill.background()

    title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title6_frame = title6.text_frame
    title6_para = title6_frame.paragraphs[0]
    title6_para.text = "5. 2026 실행 로드맵 및 종합 KPI"
    title6_para.font.size = Pt(28)
    title6_para.font.bold = True
    title6_para.font.color.rgb = DARK_BLUE

    # 3개 주요 KPI (원형 게이지)
    kpis_y = Inches(1.5)
    kpis = [
        {"label": "가공비 절감", "value": 10, "unit": "%", "color": PRIMARY_BLUE},
        {"label": "품질 개선", "value": 10, "unit": "%", "color": LIGHT_BLUE},
        {"label": "손실시간", "value": 5, "unit": "%", "color": NAVY}
    ]

    kpi_spacing = Inches(3.1)
    kpi_start_x = Inches(0.8)

    for i, kpi in enumerate(kpis):
        x_pos = kpi_start_x + i * kpi_spacing

        # 라벨
        label_box = slide6.shapes.add_textbox(x_pos, kpis_y, Inches(2.5), Inches(0.3))
        l_frame = label_box.text_frame
        l_para = l_frame.paragraphs[0]
        l_para.text = kpi["label"]
        l_para.font.size = Pt(14)
        l_para.font.bold = True
        l_para.font.color.rgb = DARK_BLUE
        l_para.alignment = PP_ALIGN.CENTER

        # 원형 게이지
        gauge_x = x_pos + Inches(1.25)
        gauge_y = kpis_y + Inches(0.5)
        create_circular_gauge(slide6, gauge_x, gauge_y, Inches(0.7), kpi["value"] * 10, kpi["color"])

    # Q1~Q4 타임라인 (프로그레스 바 형태)
    timeline_y = Inches(4.2)
    quarters = [
        {"q": "Q1", "progress": 100, "tasks": "TOOL 개발 / 시스템 구축", "color": PRIMARY_BLUE},
        {"q": "Q2", "progress": 75, "tasks": "파일럿 운영 / 피드백 반영", "color": LIGHT_BLUE},
        {"q": "Q3", "progress": 50, "tasks": "전사 확대 / 교육 실시", "color": NAVY},
        {"q": "Q4", "progress": 25, "tasks": "목표 달성 / 성과 점검", "color": DARK_BLUE}
    ]

    for i, qt in enumerate(quarters):
        q_y = timeline_y + i * Inches(0.7)

        # Q 라벨
        q_label = slide6.shapes.add_textbox(Inches(0.5), q_y, Inches(0.8), Inches(0.4))
        ql_frame = q_label.text_frame
        ql_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        ql_para = ql_frame.paragraphs[0]
        ql_para.text = qt["q"]
        ql_para.font.size = Pt(18)
        ql_para.font.bold = True
        ql_para.font.color.rgb = qt["color"]
        ql_para.alignment = PP_ALIGN.CENTER

        # 프로그레스 바
        create_progress_bar(slide6, Inches(1.5), q_y + Inches(0.075),
                           Inches(5.5), Inches(0.25), qt["progress"], qt["color"], LIGHT_GRAY)

        # 태스크
        task_box = slide6.shapes.add_textbox(Inches(7.2), q_y, Inches(2.5), Inches(0.4))
        t_frame = task_box.text_frame
        t_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        t_para = t_frame.paragraphs[0]
        t_para.text = qt["tasks"]
        t_para.font.size = Pt(11)
        t_para.font.color.rgb = DARK_BLUE

    page6 = slide6.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page6.text_frame.text = "06"
    page6.text_frame.paragraphs[0].font.size = Pt(11)
    page6.text_frame.paragraphs[0].font.color.rgb = GRAY
    page6.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # 저장
    output_file = '참고디자인_기반_2026전략.pptx'
    prs.save(output_file)
    print(f"✅ 참고디자인 기반 PPT 생성 완료: {output_file}")
    print(f"📄 총 6페이지")
    print(f"🎨 디자인 특징:")
    print(f"   - 프로그레스 바 중심 시각화")
    print(f"   - 원형 게이지로 주요 지표 표현")
    print(f"   - 블루 계열 컬러 팔레트")
    print(f"   - 깔끔한 박스 레이아웃")
    print(f"   - 참고디자인.png 스타일 완벽 구현")
    return output_file

if __name__ == "__main__":
    create_reference_design_presentation()
