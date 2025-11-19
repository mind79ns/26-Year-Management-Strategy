#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
제조1팀 전략 이력 PPT 생성 (2021-2025)
- 과거 전략 분석 결과 시각화
- 연도별 핵심 과제 및 KPI
- 2026년 전략과의 연결성
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# 색상 팔레트
NAVY = RGBColor(25, 42, 86)
GOLD = RGBColor(212, 175, 55)
LIGHT_BLUE = RGBColor(52, 152, 219)
GREEN = RGBColor(46, 204, 113)
ORANGE = RGBColor(230, 126, 34)
RED = RGBColor(231, 76, 60)
PURPLE = RGBColor(155, 89, 182)
GRAY = RGBColor(127, 140, 141)
LIGHT_GRAY = RGBColor(236, 240, 241)
WHITE = RGBColor(255, 255, 255)

def create_title_slide(prs):
    """페이지 1: 커버"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "제조1팀 경영전략 이력"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = GOLD
    title_para.alignment = PP_ALIGN.CENTER

    # 부제목
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4), Inches(8), Inches(0.6)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "2021-2025 전략 분석 및 2026 방향"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = WHITE
    subtitle_para.alignment = PP_ALIGN.CENTER

    # 하단 정보
    info_box = slide.shapes.add_textbox(
        Inches(1), Inches(6), Inches(8), Inches(0.8)
    )
    info_frame = info_box.text_frame
    info_frame.text = "분석 대상: 5개 전략 PPT, 72개 슬라이드\n추출 데이터: 53개 과제, 67개 KPI, 754개 키워드"
    for para in info_frame.paragraphs:
        para.font.size = Pt(14)
        para.font.color.rgb = RGBColor(200, 200, 200)
        para.alignment = PP_ALIGN.CENTER

def create_overview(prs):
    """페이지 2: 분석 개요"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "분석 개요 및 대상"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # 분석 대상 파일 (좌측)
    files_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.3)
    )
    files_box.fill.solid()
    files_box.fill.fore_color.rgb = RGBColor(245, 248, 250)
    files_box.line.color.rgb = LIGHT_BLUE
    files_box.line.width = Pt(2)

    files_text = files_box.text_frame
    files_text.text = "📂 분석 대상 파일"
    files_text.paragraphs[0].font.size = Pt(18)
    files_text.paragraphs[0].font.bold = True
    files_text.paragraphs[0].font.color.rgb = LIGHT_BLUE
    files_text.paragraphs[0].alignment = PP_ALIGN.CENTER

    files_list = [
        ("21년 SMD 전략", "8", "슬라이드"),
        ("21년 경영전략 (상반기)", "9", "슬라이드"),
        ("21년 경영전략 (하반기)", "22", "슬라이드"),
        ("22년 제조1 경영전략", "19", "슬라이드"),
        ("22년 경영전략 (하반기)", "14", "슬라이드")
    ]

    for i, (name, count, unit) in enumerate(files_list):
        p = files_text.add_paragraph()
        p.text = f"• {name}"
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY
        p.font.bold = True
        p.space_before = Pt(10)

        p2 = files_text.add_paragraph()
        p2.text = f"  {count} {unit}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY
        p2.level = 1

    # 통계 (우측)
    stats_y = 1.2
    stats = [
        {"label": "총 슬라이드", "value": "72", "color": LIGHT_BLUE},
        {"label": "추출 과제", "value": "53", "color": GREEN},
        {"label": "추출 KPI", "value": "67", "color": ORANGE},
        {"label": "고유 키워드", "value": "754", "color": PURPLE}
    ]

    for i, stat in enumerate(stats):
        y = stats_y + i * 1.3

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.3), Inches(y), Inches(4.2), Inches(1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = stat["color"]
        box.line.width = Pt(3)

        # 라벨
        label_box = slide.shapes.add_textbox(
            Inches(5.5), Inches(y + 0.15), Inches(1.5), Inches(0.3)
        )
        label_frame = label_box.text_frame
        label_frame.text = stat["label"]
        label_frame.paragraphs[0].font.size = Pt(14)
        label_frame.paragraphs[0].font.color.rgb = GRAY

        # 값
        value_box = slide.shapes.add_textbox(
            Inches(7.2), Inches(y + 0.05), Inches(2), Inches(0.9)
        )
        value_frame = value_box.text_frame
        value_frame.text = stat["value"]
        value_frame.paragraphs[0].font.size = Pt(36)
        value_frame.paragraphs[0].font.bold = True
        value_frame.paragraphs[0].font.color.rgb = stat["color"]
        value_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        value_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

def create_theme_analysis(prs):
    """페이지 3: 주제 분포 분석"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "주제 분포 분석 (Top 10)"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # 주제 데이터
    themes = [
        {"name": "생산성", "count": 63, "pct": 87.5, "color": LIGHT_BLUE},
        {"name": "실행계획", "count": 48, "pct": 66.7, "color": GREEN},
        {"name": "유실시간", "count": 48, "pct": 66.7, "color": ORANGE},
        {"name": "설비", "count": 47, "pct": 65.3, "color": PURPLE},
        {"name": "목표", "count": 39, "pct": 54.2, "color": RED}
    ]

    start_y = 1.3
    max_width = 7

    for i, theme in enumerate(themes):
        y = start_y + i * 1

        # 순위 원
        rank_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.5), Inches(y), Inches(0.5), Inches(0.5)
        )
        rank_circle.fill.solid()
        rank_circle.fill.fore_color.rgb = theme["color"]
        rank_circle.line.fill.background()

        rank_text = rank_circle.text_frame
        rank_text.text = str(i + 1)
        rank_text.paragraphs[0].font.size = Pt(20)
        rank_text.paragraphs[0].font.bold = True
        rank_text.paragraphs[0].font.color.rgb = WHITE
        rank_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        rank_text.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 주제명
        name_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(y + 0.05), Inches(1.2), Inches(0.4)
        )
        name_frame = name_box.text_frame
        name_frame.text = theme["name"]
        name_frame.paragraphs[0].font.size = Pt(16)
        name_frame.paragraphs[0].font.bold = True
        name_frame.paragraphs[0].font.color.rgb = NAVY
        name_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 프로그레스 바
        bar_width = max_width * (theme["pct"] / 100)

        bg_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.5), Inches(y + 0.1), Inches(max_width), Inches(0.3)
        )
        bg_bar.fill.solid()
        bg_bar.fill.fore_color.rgb = LIGHT_GRAY
        bg_bar.line.fill.background()

        progress_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.5), Inches(y + 0.1), Inches(bar_width), Inches(0.3)
        )
        progress_bar.fill.solid()
        progress_bar.fill.fore_color.rgb = theme["color"]
        progress_bar.line.fill.background()

        # 통계
        stats_box = slide.shapes.add_textbox(
            Inches(2.5 + bar_width + 0.1), Inches(y + 0.05), Inches(1.5), Inches(0.4)
        )
        stats_frame = stats_box.text_frame
        stats_frame.text = f"{theme['count']}회 ({theme['pct']}%)"
        stats_frame.paragraphs[0].font.size = Pt(13)
        stats_frame.paragraphs[0].font.bold = True
        stats_frame.paragraphs[0].font.color.rgb = theme["color"]
        stats_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 인사이트 박스
    insight_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(6.3), Inches(9), Inches(0.9)
    )
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(255, 250, 230)
    insight_box.line.color.rgb = GOLD
    insight_box.line.width = Pt(2)

    insight_text = insight_box.text_frame
    insight_text.text = "💡 핵심 인사이트"
    insight_text.paragraphs[0].font.size = Pt(16)
    insight_text.paragraphs[0].font.bold = True
    insight_text.paragraphs[0].font.color.rgb = GOLD

    p2 = insight_text.add_paragraph()
    p2.text = "생산성 향상(87.5%)이 압도적 1순위. 유실시간(66.7%)과 설비(65.3%)가 모든 전략의 기반"
    p2.font.size = Pt(14)
    p2.font.color.rgb = NAVY

def create_keyword_analysis(prs):
    """페이지 4: 고빈도 키워드 분석"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "고빈도 키워드 분석 (Top 20)"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # 키워드 데이터
    keywords = [
        ("개선", 134, LIGHT_BLUE),
        ("설비", 75, PURPLE),
        ("유실", 69, ORANGE),
        ("관리", 55, GREEN),
        ("과제", 53, RED),
        ("향상", 51, LIGHT_BLUE),
        ("운영", 46, PURPLE),
        ("진행", 44, GREEN),
        ("효율", 41, ORANGE),
        ("자동화", 34, LIGHT_BLUE)
    ]

    # 워드 클라우드 스타일 배치
    positions = [
        (1.5, 1.5, 2, 0.8),
        (4, 1.3, 1.8, 0.7),
        (6.5, 1.6, 1.6, 0.65),
        (1, 2.6, 1.7, 0.6),
        (3.2, 2.5, 1.5, 0.55),
        (5.2, 2.7, 1.8, 0.65),
        (7.5, 2.8, 1.4, 0.5),
        (1.8, 3.8, 1.6, 0.6),
        (4.2, 3.7, 1.5, 0.55),
        (6.8, 3.9, 1.3, 0.5)
    ]

    for i, (keyword, count, color) in enumerate(keywords):
        if i < len(positions):
            x, y, w, h = positions[i]

            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(w), Inches(h)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.fill.background()

            # 투명도 효과 (순위별)
            box.fill.transparency = 0.2 + (i * 0.03)

            text_frame = box.text_frame
            text_frame.text = keyword
            text_frame.paragraphs[0].font.size = Pt(24 - i)
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].font.color.rgb = WHITE
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            p2 = text_frame.add_paragraph()
            p2.text = f"{count}회"
            p2.font.size = Pt(12)
            p2.font.color.rgb = WHITE
            p2.alignment = PP_ALIGN.CENTER

    # 하단 분석
    analysis_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(5), Inches(9), Inches(1.8)
    )
    analysis_box.fill.solid()
    analysis_box.fill.fore_color.rgb = RGBColor(245, 248, 250)
    analysis_box.line.color.rgb = NAVY
    analysis_box.line.width = Pt(2)

    analysis_text = analysis_box.text_frame
    analysis_text.text = "🔍 키워드 트렌드 분석"
    analysis_text.paragraphs[0].font.size = Pt(18)
    analysis_text.paragraphs[0].font.bold = True
    analysis_text.paragraphs[0].font.color.rgb = NAVY
    analysis_text.paragraphs[0].alignment = PP_ALIGN.CENTER

    categories = [
        "• 개선 중심: 개선(134) + 향상(51) = 185회 → 지속적 개선이 핵심 DNA",
        "• 설비 관련: 설비(75) + 효율(41) = 116회 → 설비 중심 제조 혁신",
        "• 유실시간: 유실(69) + 관리(55) = 124회 → 유실시간 감소가 최우선",
        "• 실행 중심: 과제(53) + 진행(44) = 97회 → 체계적 실행력 강조"
    ]

    for cat in categories:
        p = analysis_text.add_paragraph()
        p.text = cat
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY
        p.space_before = Pt(8)

def create_2021_summary(prs):
    """페이지 5: 2021년 핵심 과제"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "2021년 핵심 과제 및 KPI"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # 4개 과제 박스
    tasks = [
        {
            "title": "가공비 절감",
            "icon": "💰",
            "kpis": ["CAPA 23% ↓", "가동율 24% ↓", "인원 18% ↓"],
            "actions": ["설비능력 향상", "노무비 절감", "월별 비용 관리"],
            "color": LIGHT_BLUE,
            "x": 0.5, "y": 1.2
        },
        {
            "title": "유실시간 개선",
            "icon": "⏱️",
            "kpis": ["순간정지 개선", "흡착율 관리", "PM 활동 강화"],
            "actions": ["WORST LINE 개선", "편차 분석", "습관화 정착"],
            "color": GREEN,
            "x": 5.2, "y": 1.2
        },
        {
            "title": "설비 관리",
            "icon": "🔧",
            "kpis": ["IMT 95% 비용", "고장 89% 집중", "예방보전 강화"],
            "actions": ["SPARE PART 관리", "성능 복원", "이력 분석"],
            "color": ORANGE,
            "x": 0.5, "y": 4
        },
        {
            "title": "생산성 향상",
            "icon": "📈",
            "kpis": ["SMD 효율 개선", "Line별 최적화", "CAPA 증대"],
            "actions": ["POINT BY POINT", "편차 개선", "자동화 확대"],
            "color": PURPLE,
            "x": 5.2, "y": 4
        }
    ]

    for task in tasks:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(task["x"]), Inches(task["y"]),
            Inches(4.3), Inches(2.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = task["color"]
        box.line.width = Pt(3)

        text_frame = box.text_frame
        text_frame.text = f"{task['icon']} {task['title']}"
        text_frame.paragraphs[0].font.size = Pt(18)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = task["color"]
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # KPI
        p_kpi = text_frame.add_paragraph()
        p_kpi.text = "\nKPI:"
        p_kpi.font.size = Pt(13)
        p_kpi.font.bold = True
        p_kpi.font.color.rgb = NAVY

        for kpi in task["kpis"]:
            p = text_frame.add_paragraph()
            p.text = f"• {kpi}"
            p.font.size = Pt(11)
            p.font.color.rgb = GRAY
            p.space_before = Pt(3)

        # Actions
        p_act = text_frame.add_paragraph()
        p_act.text = "\n실행:"
        p_act.font.size = Pt(13)
        p_act.font.bold = True
        p_act.font.color.rgb = NAVY

        for action in task["actions"]:
            p = text_frame.add_paragraph()
            p.text = f"→ {action}"
            p.font.size = Pt(11)
            p.font.color.rgb = task["color"]
            p.space_before = Pt(3)

def create_2022_summary(prs):
    """페이지 6: 2022년 전략 진화"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "2022년 전략 진화 및 신규 과제"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # 진화 다이어그램
    evolution = [
        {
            "from": "수동 분석",
            "to": "실시간 모니터링",
            "color": LIGHT_BLUE,
            "y": 1.5
        },
        {
            "from": "예방 보전",
            "to": "예측 보전",
            "color": GREEN,
            "y": 2.5
        },
        {
            "from": "비용 절감",
            "to": "효율 향상",
            "color": ORANGE,
            "y": 3.5
        },
        {
            "from": "개별 시스템",
            "to": "MES 통합",
            "color": PURPLE,
            "y": 4.5
        }
    ]

    for evo in evolution:
        # From 박스
        from_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(evo["y"]), Inches(2.2), Inches(0.7)
        )
        from_box.fill.solid()
        from_box.fill.fore_color.rgb = LIGHT_GRAY
        from_box.line.fill.background()

        from_text = from_box.text_frame
        from_text.text = evo["from"]
        from_text.paragraphs[0].font.size = Pt(16)
        from_text.paragraphs[0].font.color.rgb = GRAY
        from_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        from_text.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 화살표
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(3.2), Inches(evo["y"] + 0.2), Inches(1), Inches(0.3)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = evo["color"]
        arrow.line.fill.background()

        # To 박스
        to_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(4.5), Inches(evo["y"]), Inches(2.2), Inches(0.7)
        )
        to_box.fill.solid()
        to_box.fill.fore_color.rgb = evo["color"]
        to_box.line.fill.background()

        to_text = to_box.text_frame
        to_text.text = evo["to"]
        to_text.paragraphs[0].font.size = Pt(16)
        to_text.paragraphs[0].font.bold = True
        to_text.paragraphs[0].font.color.rgb = WHITE
        to_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        to_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 신규 과제
    new_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.2), Inches(1.5), Inches(2.3), Inches(4.2)
    )
    new_box.fill.solid()
    new_box.fill.fore_color.rgb = RGBColor(255, 250, 230)
    new_box.line.color.rgb = GOLD
    new_box.line.width = Pt(3)

    new_text = new_box.text_frame
    new_text.text = "⭐ 2022 신규"
    new_text.paragraphs[0].font.size = Pt(18)
    new_text.paragraphs[0].font.bold = True
    new_text.paragraphs[0].font.color.rgb = GOLD
    new_text.paragraphs[0].alignment = PP_ALIGN.CENTER

    new_items = [
        "디지털 전환",
        "MES 본격 도입",
        "IoT 센서 활용",
        "AI 예측 분석",
        "품질 자동화",
        "데이터 기반 결정"
    ]

    for item in new_items:
        p = new_text.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY
        p.space_before = Pt(10)

    # 하단 성과
    result_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(5.8), Inches(6.5), Inches(0.9)
    )
    result_box.fill.solid()
    result_box.fill.fore_color.rgb = RGBColor(230, 245, 255)
    result_box.line.color.rgb = LIGHT_BLUE
    result_box.line.width = Pt(2)

    result_text = result_box.text_frame
    result_text.text = "📊 2022 예상 성과: 가공비 5-10% 절감 | 유실시간 10-15% 감소 | 설비 가동률 5-7% 향상"
    result_text.paragraphs[0].font.size = Pt(15)
    result_text.paragraphs[0].font.bold = True
    result_text.paragraphs[0].font.color.rgb = LIGHT_BLUE
    result_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    result_text.vertical_anchor = MSO_ANCHOR.MIDDLE

def create_continuity_analysis(prs):
    """페이지 7: 전략 연속성 분석"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "전략 연속성 및 2026 방향"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # 타임라인
    timeline_y = 1.2
    years = ["2021", "2022", "2023-2025", "2026"]
    colors = [GRAY, LIGHT_BLUE, GREEN, GOLD]

    for i, (year, color) in enumerate(zip(years, colors)):
        x = 0.8 + i * 2.2

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(timeline_y), Inches(0.8), Inches(0.8)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.color.rgb = WHITE
        circle.line.width = Pt(3)

        text = circle.text_frame
        text.text = year
        text.paragraphs[0].font.size = Pt(14)
        text.paragraphs[0].font.bold = True
        text.paragraphs[0].font.color.rgb = WHITE
        text.paragraphs[0].alignment = PP_ALIGN.CENTER
        text.vertical_anchor = MSO_ANCHOR.MIDDLE

        if i < len(years) - 1:
            line = slide.shapes.add_connector(
                1,
                Inches(x + 0.8), Inches(timeline_y + 0.4),
                Inches(x + 2.2), Inches(timeline_y + 0.4)
            )
            line.line.color.rgb = GRAY
            line.line.width = Pt(2)

    # 연속성 매트릭스
    continuity = [
        {"area": "유실시간", "21": "수동 분석", "26": "실시간 감지", "status": "✓ 강화"},
        {"area": "설비관리", "21": "예방 보전", "26": "AI 예측", "status": "✓ 혁신"},
        {"area": "가공비", "21": "비용 절감", "26": "효율 극대화", "status": "✓ 전환"},
        {"area": "불량", "21": "사후 대응", "26": "재발 Zero", "status": "★ 신규 강화"},
        {"area": "자동화", "21": "개별 도입", "26": "MES 통합", "status": "✓ 확대"}
    ]

    # 테이블 헤더
    headers = ["영역", "2021-2022", "2026", "연속성"]
    header_x = [0.5, 2.5, 5.5, 7.8]
    header_w = [1.8, 2.8, 2, 1.7]

    for i, (header, x, w) in enumerate(zip(headers, header_x, header_w)):
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(2.5), Inches(w), Inches(0.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY
        box.line.fill.background()

        text = box.text_frame
        text.text = header
        text.paragraphs[0].font.size = Pt(14)
        text.paragraphs[0].font.bold = True
        text.paragraphs[0].font.color.rgb = WHITE
        text.paragraphs[0].alignment = PP_ALIGN.CENTER
        text.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 테이블 내용
    for i, cont in enumerate(continuity):
        y = 3.1 + i * 0.6

        contents = [cont["area"], cont["21"], cont["26"], cont["status"]]

        for j, (content, x, w) in enumerate(zip(contents, header_x, header_w)):
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(w), Inches(0.5)
            )
            box.fill.solid()
            if j == 3:
                box.fill.fore_color.rgb = RGBColor(255, 250, 230) if "신규" in content else RGBColor(245, 255, 245)
            else:
                box.fill.fore_color.rgb = WHITE
            box.line.color.rgb = LIGHT_GRAY

            text = box.text_frame
            text.text = content
            text.paragraphs[0].font.size = Pt(12)
            if j == 0:
                text.paragraphs[0].font.bold = True
            text.paragraphs[0].font.color.rgb = NAVY
            text.paragraphs[0].alignment = PP_ALIGN.CENTER
            text.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 하단 인사이트
    insight_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(6.2), Inches(9), Inches(0.8)
    )
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(255, 245, 235)
    insight_box.line.color.rgb = ORANGE
    insight_box.line.width = Pt(2)

    insight_text = insight_box.text_frame
    insight_text.text = "🎯 2026 차별화 전략"
    insight_text.paragraphs[0].font.size = Pt(16)
    insight_text.paragraphs[0].font.bold = True
    insight_text.paragraphs[0].font.color.rgb = ORANGE

    p2 = insight_text.add_paragraph()
    p2.text = "과거 기반 유지·강화 + 불량 재발 Zero 신규 강화 + 설비관리 혁신으로 도약"
    p2.font.size = Pt(14)
    p2.font.color.rgb = NAVY

def create_2026_connection(prs):
    """페이지 8: 2026 전략 연결"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "2026 전략: 과거 경험 + 기술 혁신"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # 4대 전략 연결
    strategies = [
        {
            "name": "순간유실 Zero",
            "past": "2021: 순간정지 개선\n2022: 실시간 모니터링",
            "new": "2026: AI 실시간 감지\n자동 대응 시스템",
            "improvement": "1000배 빠른 대응",
            "color": LIGHT_BLUE,
            "y": 1.2
        },
        {
            "name": "불량재발 Zero",
            "past": "2021: 사후 분석\n2022: 재발 방지 시작",
            "new": "2026: 비전검사 + AI\n불량 DB 패턴 분석",
            "improvement": "재발률 70% 감소",
            "color": GREEN,
            "y": 3
        },
        {
            "name": "설비 CAPA 증대",
            "past": "2021: CAPA 향상 활동\n2022: 병목 개선",
            "new": "2026: 고속화 + 자동화\nLine 증설",
            "improvement": "생산능력 15% 향상",
            "color": ORANGE,
            "y": 4.8
        }
    ]

    for strat in strategies:
        # 과거 박스
        past_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(strat["y"]), Inches(2.8), Inches(1.5)
        )
        past_box.fill.solid()
        past_box.fill.fore_color.rgb = LIGHT_GRAY
        past_box.line.fill.background()

        past_text = past_box.text_frame
        past_text.text = f"📋 {strat['name']}\n과거"
        past_text.paragraphs[0].font.size = Pt(14)
        past_text.paragraphs[0].font.bold = True
        past_text.paragraphs[0].font.color.rgb = NAVY
        past_text.paragraphs[0].alignment = PP_ALIGN.CENTER

        p = past_text.add_paragraph()
        p.text = f"\n{strat['past']}"
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.LEFT

        # 화살표
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(3.5), Inches(strat["y"] + 0.5), Inches(1.2), Inches(0.5)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = strat["color"]
        arrow.line.fill.background()

        arrow_text = arrow.text_frame
        arrow_text.text = "진화"
        arrow_text.paragraphs[0].font.size = Pt(12)
        arrow_text.paragraphs[0].font.bold = True
        arrow_text.paragraphs[0].font.color.rgb = WHITE
        arrow_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        arrow_text.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 2026 박스
        new_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5), Inches(strat["y"]), Inches(2.8), Inches(1.5)
        )
        new_box.fill.solid()
        new_box.fill.fore_color.rgb = strat["color"]
        new_box.line.fill.background()

        new_text = new_box.text_frame
        new_text.text = f"⭐ 2026"
        new_text.paragraphs[0].font.size = Pt(14)
        new_text.paragraphs[0].font.bold = True
        new_text.paragraphs[0].font.color.rgb = WHITE
        new_text.paragraphs[0].alignment = PP_ALIGN.CENTER

        p = new_text.add_paragraph()
        p.text = f"\n{strat['new']}"
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT

        # 개선 효과
        effect_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(8.1), Inches(strat["y"] + 0.3), Inches(1.4), Inches(0.9)
        )
        effect_box.fill.solid()
        effect_box.fill.fore_color.rgb = WHITE
        effect_box.line.color.rgb = strat["color"]
        effect_box.line.width = Pt(2)

        effect_text = effect_box.text_frame
        effect_text.text = strat["improvement"]
        effect_text.paragraphs[0].font.size = Pt(11)
        effect_text.paragraphs[0].font.bold = True
        effect_text.paragraphs[0].font.color.rgb = strat["color"]
        effect_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        effect_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 신규 전략 (설비관리 혁신)
    new_strategy_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(6.5), Inches(9), Inches(0.8)
    )
    new_strategy_box.fill.solid()
    new_strategy_box.fill.fore_color.rgb = RGBColor(245, 235, 255)
    new_strategy_box.line.color.rgb = PURPLE
    new_strategy_box.line.width = Pt(3)

    new_text = new_strategy_box.text_frame
    new_text.text = "★ 전략4: 설비관리 혁신 (신규)"
    new_text.paragraphs[0].font.size = Pt(16)
    new_text.paragraphs[0].font.bold = True
    new_text.paragraphs[0].font.color.rgb = PURPLE

    p2 = new_text.add_paragraph()
    p2.text = "과거 예방보전 → 2026 IoT 센서 + AI 수명 예측 + 긴급 대응 체계 (설비 고장 50% 감소)"
    p2.font.size = Pt(13)
    p2.font.color.rgb = NAVY

def create_conclusion(prs):
    """페이지 9: 결론 및 기대효과"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "결론: 2021-2026 전략의 완성"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # 3단계 진화
    stages = [
        {
            "stage": "2021-2022\n기반 구축",
            "content": "• 문제 인식\n• 개선 활동 정착\n• 데이터 수집",
            "color": GRAY,
            "x": 0.8
        },
        {
            "stage": "2023-2025\n체계화",
            "content": "• 시스템 고도화\n• MES 연동\n• 자동화 확대",
            "color": LIGHT_BLUE,
            "x": 3.8
        },
        {
            "stage": "2026\n혁신 도약",
            "content": "• AI/IoT 활용\n• 실시간 대응\n• 통합 플랫폼",
            "color": GOLD,
            "x": 6.8
        }
    ]

    for stage in stages:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(stage["x"]), Inches(1.2), Inches(2.5), Inches(2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = stage["color"]
        box.line.fill.background()

        text = box.text_frame
        text.text = stage["stage"]
        text.paragraphs[0].font.size = Pt(16)
        text.paragraphs[0].font.bold = True
        text.paragraphs[0].font.color.rgb = WHITE
        text.paragraphs[0].alignment = PP_ALIGN.CENTER

        p = text.add_paragraph()
        p.text = f"\n{stage['content']}"
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT

    # 화살표
    for i in range(2):
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(3.5 + i * 3), Inches(2), Inches(0.5), Inches(0.4)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ORANGE
        arrow.line.fill.background()

    # 기대효과
    effects_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(3.8), Inches(9), Inches(2.5)
    )
    effects_box.fill.solid()
    effects_box.fill.fore_color.rgb = RGBColor(245, 250, 255)
    effects_box.line.color.rgb = NAVY
    effects_box.line.width = Pt(3)

    effects_text = effects_box.text_frame
    effects_text.text = "🎯 2026 기대효과"
    effects_text.paragraphs[0].font.size = Pt(24)
    effects_text.paragraphs[0].font.bold = True
    effects_text.paragraphs[0].font.color.rgb = NAVY
    effects_text.paragraphs[0].alignment = PP_ALIGN.CENTER

    effect_items = [
        ("💰 가공비 10% 절감", "연간 수억원 비용 절감", LIGHT_BLUE),
        ("📈 품질 10% 개선", "고객 만족도 대폭 향상", GREEN),
        ("⏱️ 유실시간 60% 감소", "생산성 5% 증가", ORANGE),
        ("🔧 설비 수명 20% 연장", "투자비 장기 절감", PURPLE),
        ("📊 평가가동 효율", "SMD 91%, RADIAL 85%, AXIAL 85%", RED)
    ]

    for item, detail, color in effect_items:
        p = effects_text.add_paragraph()
        p.text = f"\n{item}"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_before = Pt(10)

        p2 = effects_text.add_paragraph()
        p2.text = f"  → {detail}"
        p2.font.size = Pt(12)
        p2.font.color.rgb = GRAY
        p2.level = 1

    # 핵심 메시지
    message_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(6.5), Inches(9), Inches(0.8)
    )
    message_box.fill.solid()
    message_box.fill.fore_color.rgb = GOLD
    message_box.line.fill.background()

    message_text = message_box.text_frame
    message_text.text = "과거 5년의 경험과 노하우 + 최신 기술 = 2026년 스마트 제조 혁신 완성"
    message_text.paragraphs[0].font.size = Pt(20)
    message_text.paragraphs[0].font.bold = True
    message_text.paragraphs[0].font.color.rgb = WHITE
    message_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    message_text.vertical_anchor = MSO_ANCHOR.MIDDLE

def main():
    """메인 실행 함수"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("=" * 70)
    print("제조1팀 전략 이력 PPT 생성 시작")
    print("=" * 70)

    create_title_slide(prs)
    print("✓ 페이지 1: 커버")

    create_overview(prs)
    print("✓ 페이지 2: 분석 개요")

    create_theme_analysis(prs)
    print("✓ 페이지 3: 주제 분포")

    create_keyword_analysis(prs)
    print("✓ 페이지 4: 키워드 분석")

    create_2021_summary(prs)
    print("✓ 페이지 5: 2021년 과제")

    create_2022_summary(prs)
    print("✓ 페이지 6: 2022년 진화")

    create_continuity_analysis(prs)
    print("✓ 페이지 7: 전략 연속성")

    create_2026_connection(prs)
    print("✓ 페이지 8: 2026 연결")

    create_conclusion(prs)
    print("✓ 페이지 9: 결론")

    output_file = '제조1팀_전략이력_2021-2026.pptx'
    prs.save(output_file)

    print("\n" + "=" * 70)
    print(f"✅ PPT 생성 완료: {output_file}")
    print(f"📄 총 9페이지")
    print(f"🎨 특징:")
    print(f"   - 2021-2022년 전략 심층 분석 시각화")
    print(f"   - 주제 분포 및 키워드 트렌드 분석")
    print(f"   - 연도별 핵심 과제 및 KPI 정리")
    print(f"   - 전략 연속성 및 진화 과정 표현")
    print(f"   - 2026년 전략과의 명확한 연결고리")
    print("=" * 70)

if __name__ == "__main__":
    main()
