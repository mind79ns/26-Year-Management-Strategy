#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
옵션 3: 타임라인/프로세스 중심 디자인
플로우차트 스타일, 화살표, 단계별 프로세스
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_arrow(slide, x1, y1, x2, y2, color):
    """화살표 추가"""
    # 라인
    connector = slide.shapes.add_connector(
        1, x1, y1, x2, y2  # STRAIGHT connector
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(3)

def create_process_box(slide, left, top, width, height, number, title, color, text_color=None):
    """프로세스 박스 생성"""
    if text_color is None:
        text_color = RGBColor(255, 255, 255)

    # 박스
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = color
    box.line.width = Pt(2)

    frame = box.text_frame
    frame.margin_left = Inches(0.2)
    frame.margin_right = Inches(0.2)
    frame.margin_top = Inches(0.15)
    frame.vertical_anchor = MSO_ANCHOR.TOP

    # 번호
    num_para = frame.paragraphs[0]
    num_para.text = number
    num_para.font.size = Pt(32)
    num_para.font.bold = True
    num_para.font.color.rgb = text_color
    num_para.alignment = PP_ALIGN.CENTER
    num_para.space_after = Pt(8)

    # 제목
    title_para = frame.add_paragraph()
    title_para.text = title
    title_para.font.size = Pt(16)
    title_para.font.bold = True
    title_para.font.color.rgb = text_color
    title_para.alignment = PP_ALIGN.CENTER

    return box

def create_timeline_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 색상 팔레트
    NAVY = RGBColor(25, 55, 109)
    BLUE = RGBColor(52, 152, 219)
    LIGHT_BLUE = RGBColor(174, 214, 241)
    ORANGE = RGBColor(230, 126, 34)
    LIGHT_ORANGE = RGBColor(245, 203, 167)
    GREEN = RGBColor(39, 174, 96)
    LIGHT_GREEN = RGBColor(169, 223, 191)
    RED = RGBColor(231, 76, 60)
    YELLOW = RGBColor(241, 196, 15)
    PURPLE = RGBColor(142, 68, 173)
    GRAY = RGBColor(127, 140, 141)
    DARK_GRAY = RGBColor(52, 73, 94)
    WHITE = RGBColor(255, 255, 255)

    # ========== 슬라이드 1: 표지 - 플로우 스타일 ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = NAVY

    # 중앙 메인 박스
    main_box = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(2.5), Inches(6), Inches(2.5)
    )
    main_box.fill.solid()
    main_box.fill.fore_color.rgb = WHITE
    main_box.line.color.rgb = BLUE
    main_box.line.width = Pt(5)

    main_frame = main_box.text_frame
    main_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 년도
    year_para = main_frame.paragraphs[0]
    year_para.text = "2026"
    year_para.font.size = Pt(28)
    year_para.font.bold = True
    year_para.font.color.rgb = BLUE
    year_para.alignment = PP_ALIGN.CENTER
    year_para.space_after = Pt(10)

    # 타이틀
    title_para = main_frame.add_paragraph()
    title_para.text = "경영전략"
    title_para.font.size = Pt(56)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY
    title_para.alignment = PP_ALIGN.CENTER
    title_para.space_after = Pt(15)

    # 서브타이틀
    sub_para = main_frame.add_paragraph()
    sub_para.text = "제조1팀 스마트화 로드맵"
    sub_para.font.size = Pt(18)
    sub_para.font.color.rgb = GRAY
    sub_para.alignment = PP_ALIGN.CENTER

    # 4개 코너 박스
    corner_boxes = [
        {"text": "가공비\n-10%", "x": Inches(0.5), "y": Inches(0.8), "color": BLUE},
        {"text": "품질\n-10%", "x": Inches(8.3), "y": Inches(0.8), "color": GREEN},
        {"text": "손실시간\n-5%", "x": Inches(0.5), "y": Inches(6), "color": ORANGE},
        {"text": "OH\n달성", "x": Inches(8.3), "y": Inches(6), "color": RED}
    ]

    for cb in corner_boxes:
        box = slide1.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            cb["x"], cb["y"], Inches(1.2), Inches(1.1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = cb["color"]
        box.line.fill.background()

        frame = box.text_frame
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = frame.paragraphs[0]
        para.text = cb["text"]
        para.font.size = Pt(16)
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER

    # ========== 슬라이드 2: 문제 → 전략 플로우 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = WHITE

    # 제목
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "현황 → 전략 프로세스 맵"
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY
    title_para.alignment = PP_ALIGN.CENTER

    # 좌측: 문제점 (3개 박스)
    problems_x = Inches(0.5)
    problems = [
        {"y": Inches(1.5), "title": "느린 대응", "color": RED},
        {"y": Inches(3), "title": "반복 작업", "color": ORANGE},
        {"y": Inches(4.5), "title": "품질 문제", "color": PURPLE}
    ]

    for prob in problems:
        box = slide2.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            problems_x, prob["y"], Inches(2.2), Inches(1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = prob["color"]
        box.line.fill.background()

        frame = box.text_frame
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = frame.paragraphs[0]
        para.text = prob["title"]
        para.font.size = Pt(18)
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER

        # 화살표 (오른쪽으로)
        arrow = slide2.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            problems_x + Inches(2.3), prob["y"] + Inches(0.3),
            Inches(0.8), Inches(0.4)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRAY
        arrow.line.fill.background()

    # 중앙: 전략 허브
    hub_box = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.5), Inches(2.5), Inches(3), Inches(1.5)
    )
    hub_box.fill.solid()
    hub_box.fill.fore_color.rgb = NAVY
    hub_box.line.color.rgb = BLUE
    hub_box.line.width = Pt(4)

    hub_frame = hub_box.text_frame
    hub_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    hub_para = hub_frame.paragraphs[0]
    hub_para.text = "3대 전략"
    hub_para.font.size = Pt(32)
    hub_para.font.bold = True
    hub_para.font.color.rgb = WHITE
    hub_para.alignment = PP_ALIGN.CENTER
    hub_para.space_after = Pt(8)

    hub_sub = hub_frame.add_paragraph()
    hub_sub.text = "스마트화 솔루션"
    hub_sub.font.size = Pt(14)
    hub_sub.font.color.rgb = LIGHT_BLUE
    hub_sub.alignment = PP_ALIGN.CENTER

    # 우측: 결과 (3개 박스)
    results_x = Inches(7.2)
    results = [
        {"y": Inches(1.5), "title": "손실 시간\n제로화", "color": BLUE},
        {"y": Inches(3), "title": "불량 재발\nZero", "color": GREEN},
        {"y": Inches(4.5), "title": "설비 CAPA\n증가", "color": ORANGE}
    ]

    for res in results:
        # 화살표
        arrow = slide2.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(6.6), res["y"] + Inches(0.3),
            Inches(0.5), Inches(0.4)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = res["color"]
        arrow.line.fill.background()

        # 결과 박스
        box = slide2.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            results_x, res["y"], Inches(2.3), Inches(1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = res["color"]
        box.line.fill.background()

        frame = box.text_frame
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = frame.paragraphs[0]
        para.text = res["title"]
        para.font.size = Pt(16)
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER

    # 페이지 번호
    page_box = slide2.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page_frame = page_box.text_frame
    page_frame.text = "02"
    page_frame.paragraphs[0].font.size = Pt(14)
    page_frame.paragraphs[0].font.color.rgb = GRAY

    # ========== 슬라이드 3: 전략 1 프로세스 플로우 ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = RGBColor(240, 248, 255)

    # 헤더
    header = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE
    header.line.fill.background()

    header_frame = header.text_frame
    header_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    header_para = header_frame.paragraphs[0]
    header_para.text = "전략 1  →  손실 시간 제로화 프로젝트"
    header_para.font.size = Pt(28)
    header_para.font.bold = True
    header_para.font.color.rgb = WHITE
    header_para.alignment = PP_ALIGN.CENTER

    # 프로세스 플로우 (좌 → 우)
    process_y = Inches(2)
    processes = [
        {"num": "1", "title": "DATA 수집", "desc": "MES 연동\n이상 감지", "color": BLUE},
        {"num": "2", "title": "분석", "desc": "손실 항목\n순위화", "color": BLUE},
        {"num": "3", "title": "실행", "desc": "TOP 10\n집중 개선", "color": GREEN},
        {"num": "4", "title": "결과", "desc": "5% 감소\n달성", "color": GREEN}
    ]

    process_width = Inches(1.8)
    process_spacing = Inches(0.4)
    process_start = Inches(1)

    for i, proc in enumerate(processes):
        x_pos = process_start + (i * (process_width + process_spacing))

        # 프로세스 박스
        box = create_process_box(
            slide3, x_pos, process_y, process_width, Inches(1.8),
            proc["num"], proc["title"], proc["color"]
        )

        # 설명 추가
        frame = box.text_frame
        desc_para = frame.add_paragraph()
        desc_para.text = proc["desc"]
        desc_para.font.size = Pt(12)
        desc_para.font.color.rgb = WHITE
        desc_para.alignment = PP_ALIGN.CENTER
        desc_para.space_before = Pt(10)

        # 화살표 (마지막 제외)
        if i < 3:
            arrow = slide3.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x_pos + process_width + Inches(0.05), process_y + Inches(0.7),
                Inches(0.3), Inches(0.4)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = DARK_GRAY
            arrow.line.fill.background()

    # 하단: 세부 액션
    action_y = Inches(4.3)

    action_title = slide3.shapes.add_textbox(Inches(1), action_y, Inches(8), Inches(0.3))
    action_title_frame = action_title.text_frame
    action_title_para = action_title_frame.paragraphs[0]
    action_title_para.text = "핵심 액션"
    action_title_para.font.size = Pt(20)
    action_title_para.font.bold = True
    action_title_para.font.color.rgb = NAVY

    actions = [
        "✓ 자동분석 TOOL 제작 → MES DATA 활용",
        "✓ 순간유실 가시화 → 모바일 앱 원터치 입력",
        "✓ LINE별 일/주/월 분석 → 우선순위 타겟팅"
    ]

    action_box = slide3.shapes.add_textbox(Inches(1), action_y + Inches(0.5), Inches(8), Inches(1.5))
    action_frame = action_box.text_frame
    action_frame.word_wrap = True

    for i, action in enumerate(actions):
        if i == 0:
            para = action_frame.paragraphs[0]
        else:
            para = action_frame.add_paragraph()
        para.text = action
        para.font.size = Pt(16)
        para.font.color.rgb = DARK_GRAY
        para.space_after = Pt(12)

    # KPI 박스
    kpi_box = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(6.3), Inches(8), Inches(0.8)
    )
    kpi_box.fill.solid()
    kpi_box.fill.fore_color.rgb = RED
    kpi_box.line.fill.background()

    kpi_frame = kpi_box.text_frame
    kpi_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    kpi_para = kpi_frame.paragraphs[0]
    kpi_para.text = "🎯 KPI: 손실시간 5% 감소 → 점당 가공비 직접 절감"
    kpi_para.font.size = Pt(20)
    kpi_para.font.bold = True
    kpi_para.font.color.rgb = WHITE
    kpi_para.alignment = PP_ALIGN.CENTER

    # 페이지 번호
    page_box = slide3.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page_frame = page_box.text_frame
    page_frame.text = "03"
    page_frame.paragraphs[0].font.size = Pt(14)
    page_frame.paragraphs[0].font.color.rgb = GRAY

    # ========== 슬라이드 4: 전략 2 프로세스 플로우 ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    slide4.background.fill.solid()
    slide4.background.fill.fore_color.rgb = RGBColor(255, 250, 240)

    # 헤더
    header = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = GREEN
    header.line.fill.background()

    header_frame = header.text_frame
    header_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    header_para = header_frame.paragraphs[0]
    header_para.text = "전략 2  →  불량 재발 Zero 챌린지"
    header_para.font.size = Pt(28)
    header_para.font.bold = True
    header_para.font.color.rgb = WHITE
    header_para.alignment = PP_ALIGN.CENTER

    # 순환 프로세스 (원형 배치)
    processes_circ = [
        {"num": "1", "title": "불량 발생", "desc": "사진 촬영", "x": Inches(1.5), "y": Inches(2), "color": RED},
        {"num": "2", "title": "즉시 공유", "desc": "MES 연동", "x": Inches(4.2), "y": Inches(1.5), "color": ORANGE},
        {"num": "3", "title": "전 조 알람", "desc": "자동 전달", "x": Inches(6.9), "y": Inches(2), "color": YELLOW},
        {"num": "4", "title": "조치 입력", "desc": "DB 저장", "x": Inches(6.9), "y": Inches(4), "color": GREEN},
        {"num": "5", "title": "재발 추적", "desc": "Zero 목표", "x": Inches(4.2), "y": Inches(4.5), "color": BLUE},
        {"num": "6", "title": "재발 방지", "desc": "완료", "x": Inches(1.5), "y": Inches(4), "color": PURPLE}
    ]

    box_size = Inches(1.5)
    for proc in processes_circ:
        box = create_process_box(
            slide4, proc["x"], proc["y"], box_size, Inches(1.2),
            proc["num"], proc["title"], proc["color"]
        )

        frame = box.text_frame
        desc_para = frame.add_paragraph()
        desc_para.text = proc["desc"]
        desc_para.font.size = Pt(10)
        desc_para.font.color.rgb = WHITE
        desc_para.alignment = PP_ALIGN.CENTER
        desc_para.space_before = Pt(5)

    # 중앙 텍스트
    center_box = slide4.shapes.add_textbox(Inches(3.5), Inches(2.8), Inches(3), Inches(0.8))
    center_frame = center_box.text_frame
    center_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    center_para = center_frame.paragraphs[0]
    center_para.text = "순환\n프로세스"
    center_para.font.size = Pt(22)
    center_para.font.bold = True
    center_para.font.color.rgb = NAVY
    center_para.alignment = PP_ALIGN.CENTER

    # KPI 박스
    kpi_box = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(6.3), Inches(8), Inches(0.8)
    )
    kpi_box.fill.solid()
    kpi_box.fill.fore_color.rgb = RED
    kpi_box.line.fill.background()

    kpi_frame = kpi_box.text_frame
    kpi_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    kpi_para = kpi_frame.paragraphs[0]
    kpi_para.text = "🎯 KPI: 재발 불량 30% 감소 / 품질 불량 10% 감소"
    kpi_para.font.size = Pt(20)
    kpi_para.font.bold = True
    kpi_para.font.color.rgb = WHITE
    kpi_para.alignment = PP_ALIGN.CENTER

    # 페이지 번호
    page_box = slide4.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page_frame = page_box.text_frame
    page_frame.text = "04"
    page_frame.paragraphs[0].font.size = Pt(14)
    page_frame.paragraphs[0].font.color.rgb = GRAY

    # ========== 슬라이드 5: 전략 3 프로세스 플로우 ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    slide5.background.fill.solid()
    slide5.background.fill.fore_color.rgb = RGBColor(255, 245, 240)

    # 헤더
    header = slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = ORANGE
    header.line.fill.background()

    header_frame = header.text_frame
    header_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    header_para = header_frame.paragraphs[0]
    header_para.text = "전략 3  →  설비 CAPA 증가 및 공정 최적화"
    header_para.font.size = Pt(28)
    header_para.font.bold = True
    header_para.font.color.rgb = WHITE
    header_para.alignment = PP_ALIGN.CENTER

    # 3단계 레이어 (상→하)
    layers = [
        {
            "title": "① C/T 단축 개선",
            "items": ["P/G 운영 DEEP 분석", "최단거리 프로세스", "단 1초라도 줄이기"],
            "y": Inches(1.5),
            "color": ORANGE
        },
        {
            "title": "② 기본 BASE 강화",
            "items": ["설비 효율 유지", "성능/PM 관련 활동", "지속적 모니터링"],
            "y": Inches(3.3),
            "color": BLUE
        },
        {
            "title": "③ 공정 최적화",
            "items": ["RADIAL2 안정화", "SMD LAY OUT 개선", "설비 배치 효율화"],
            "y": Inches(5.1),
            "color": GREEN
        }
    ]

    for layer in layers:
        # 레이어 박스
        layer_box = slide5.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), layer["y"], Inches(8), Inches(1.5)
        )
        layer_box.fill.solid()
        layer_box.fill.fore_color.rgb = layer["color"]
        layer_box.line.fill.background()

        frame = layer_box.text_frame
        frame.margin_left = Inches(0.3)
        frame.margin_top = Inches(0.2)

        # 제목
        title_para = frame.paragraphs[0]
        title_para.text = layer["title"]
        title_para.font.size = Pt(20)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        title_para.space_after = Pt(10)

        # 항목들
        items_para = frame.add_paragraph()
        items_para.text = "  →  ".join(layer["items"])
        items_para.font.size = Pt(14)
        items_para.font.color.rgb = WHITE

    # KPI 박스
    kpi_box = slide5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(6.8), Inches(8), Inches(0.5)
    )
    kpi_box.fill.solid()
    kpi_box.fill.fore_color.rgb = RED
    kpi_box.line.fill.background()

    kpi_frame = kpi_box.text_frame
    kpi_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    kpi_para = kpi_frame.paragraphs[0]
    kpi_para.text = "🎯 KPI: CAPA 증가 / OVERTIME 감소 / OH 달성"
    kpi_para.font.size = Pt(18)
    kpi_para.font.bold = True
    kpi_para.font.color.rgb = WHITE
    kpi_para.alignment = PP_ALIGN.CENTER

    # 페이지 번호
    page_box = slide5.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page_frame = page_box.text_frame
    page_frame.text = "05"
    page_frame.paragraphs[0].font.size = Pt(14)
    page_frame.paragraphs[0].font.color.rgb = GRAY

    # ========== 슬라이드 6: 통합 타임라인 로드맵 ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    slide6.background.fill.solid()
    slide6.background.fill.fore_color.rgb = WHITE

    # 제목
    title_box = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "2026 통합 실행 로드맵"
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY
    title_para.alignment = PP_ALIGN.CENTER

    # 타임라인 (좌→우)
    timeline_y = Inches(1.5)
    quarters_timeline = [
        {"q": "Q1", "color": BLUE, "tasks": ["TOOL 개발", "시스템 구축", "앱 제작"]},
        {"q": "Q2", "color": GREEN, "tasks": ["파일럿 운영", "피드백", "개선"]},
        {"q": "Q3", "color": ORANGE, "tasks": ["전사 확대", "교육", "모니터링"]},
        {"q": "Q4", "color": RED, "tasks": ["목표 달성", "성과 점검", "2027 계획"]}
    ]

    q_width = Inches(2)
    q_spacing = Inches(0.25)
    q_start = Inches(0.5)

    for i, qt in enumerate(quarters_timeline):
        x_pos = q_start + (i * (q_width + q_spacing))

        # 분기 박스
        q_box = slide6.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, timeline_y, q_width, Inches(2.5)
        )
        q_box.fill.solid()
        q_box.fill.fore_color.rgb = qt["color"]
        q_box.line.fill.background()

        frame = q_box.text_frame
        frame.margin_left = Inches(0.2)
        frame.margin_top = Inches(0.2)

        # 분기명
        q_para = frame.paragraphs[0]
        q_para.text = qt["q"]
        q_para.font.size = Pt(36)
        q_para.font.bold = True
        q_para.font.color.rgb = WHITE
        q_para.alignment = PP_ALIGN.CENTER
        q_para.space_after = Pt(15)

        # 태스크
        for task in qt["tasks"]:
            task_para = frame.add_paragraph()
            task_para.text = f"• {task}"
            task_para.font.size = Pt(13)
            task_para.font.color.rgb = WHITE
            task_para.space_after = Pt(8)

        # 화살표 (마지막 제외)
        if i < 3:
            arrow = slide6.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x_pos + q_width + Inches(0.05), timeline_y + Inches(1),
                Inches(0.15), Inches(0.5)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = DARK_GRAY
            arrow.line.fill.background()

    # 하단: 최종 목표 플로우
    goal_y = Inches(4.5)

    goal_title = slide6.shapes.add_textbox(Inches(0.5), goal_y, Inches(9), Inches(0.4))
    goal_title_frame = goal_title.text_frame
    goal_title_para = goal_title_frame.paragraphs[0]
    goal_title_para.text = "최종 목표 달성 플로우"
    goal_title_para.font.size = Pt(24)
    goal_title_para.font.bold = True
    goal_title_para.font.color.rgb = NAVY
    goal_title_para.alignment = PP_ALIGN.CENTER

    # 목표 박스들 (좌→우)
    goals_flow = [
        {"text": "기회손실\n최소화", "color": BLUE},
        {"text": "가공비\n10% 절감", "color": GREEN},
        {"text": "품질\n10% 개선", "color": ORANGE},
        {"text": "OH\n달성", "color": RED}
    ]

    goal_width = Inches(1.8)
    goal_spacing = Inches(0.3)
    goal_start = Inches(1)
    goal_box_y = goal_y + Inches(0.7)

    for i, goal in enumerate(goals_flow):
        x_pos = goal_start + (i * (goal_width + goal_spacing))

        box = slide6.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, goal_box_y, goal_width, Inches(1.3)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = goal["color"]
        box.line.fill.background()

        frame = box.text_frame
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = frame.paragraphs[0]
        para.text = goal["text"]
        para.font.size = Pt(18)
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER

        # 화살표 (마지막 제외)
        if i < 3:
            arrow = slide6.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x_pos + goal_width + Inches(0.05), goal_box_y + Inches(0.5),
                Inches(0.2), Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = DARK_GRAY
            arrow.line.fill.background()

    # 페이지 번호
    page_box = slide6.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page_frame = page_box.text_frame
    page_frame.text = "06"
    page_frame.paragraphs[0].font.size = Pt(14)
    page_frame.paragraphs[0].font.color.rgb = GRAY

    # 저장
    output_file = '옵션3_타임라인중심_2026전략.pptx'
    prs.save(output_file)
    print(f"✅ 옵션 3 (타임라인/프로세스 중심) PPT 생성 완료: {output_file}")
    return output_file

if __name__ == "__main__":
    create_timeline_presentation()
