#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
2026년 경영전략 PPT 생성 스크립트
제조1팀 - 자동화 제조라인 스마트화 전략
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_strategy_presentation():
    # 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== 슬라이드 1: 전략 개요 ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃

    # 배경색 설정
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 제목
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "2026년 경영전략"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER

    # 부제목
    subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "자동화 제조라인 스마트화를 통한 가공비 절감 및 품질 개선"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = RGBColor(102, 102, 102)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # 핵심 목표 박스 (상단 강조)
    goal_box = slide1.shapes.add_shape(
        1,  # Rectangle
        Inches(1), Inches(1.7), Inches(8), Inches(1.2)
    )
    goal_box.fill.solid()
    goal_box.fill.fore_color.rgb = RGBColor(0, 102, 204)
    goal_box.line.color.rgb = RGBColor(0, 51, 102)
    goal_box.line.width = Pt(2)

    goal_frame = goal_box.text_frame
    goal_frame.margin_top = Inches(0.1)
    goal_frame.margin_bottom = Inches(0.1)
    goal_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    goal_title = goal_frame.paragraphs[0]
    goal_title.text = "🎯 핵심 목표"
    goal_title.font.size = Pt(24)
    goal_title.font.bold = True
    goal_title.font.color.rgb = RGBColor(255, 255, 255)
    goal_title.alignment = PP_ALIGN.CENTER

    goal_detail = goal_frame.add_paragraph()
    goal_detail.text = "순간유실 개선  |  인력운영비 감소  |  점당가공비 10% 절감  |  품질불량 10% 감소"
    goal_detail.font.size = Pt(16)
    goal_detail.font.color.rgb = RGBColor(255, 255, 255)
    goal_detail.alignment = PP_ALIGN.CENTER
    goal_detail.space_before = Pt(8)

    # 현재 문제점 섹션
    problem_title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(4), Inches(0.4))
    problem_title_frame = problem_title_box.text_frame
    problem_title_frame.text = "📌 현재 문제점"
    problem_title_para = problem_title_frame.paragraphs[0]
    problem_title_para.font.size = Pt(20)
    problem_title_para.font.bold = True
    problem_title_para.font.color.rgb = RGBColor(204, 0, 0)

    # 문제점 상세
    problem_box = slide1.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(4), Inches(2))
    problem_frame = problem_box.text_frame
    problem_frame.word_wrap = True

    problems = [
        "① 느린 대응",
        "  설비 이상 DATA 집계 대응으로 기회손실 발생",
        "",
        "② 반복 작업",
        "  수동적 원인 분석으로 업무 비효율 심화",
        "",
        "③ 품질 문제",
        "  불량 발생 원인 대응 체감 저하"
    ]

    for i, prob in enumerate(problems):
        if i == 0:
            para = problem_frame.paragraphs[0]
        else:
            para = problem_frame.add_paragraph()
        para.text = prob
        if prob.startswith("①") or prob.startswith("②") or prob.startswith("③"):
            para.font.size = Pt(14)
            para.font.bold = True
            para.font.color.rgb = RGBColor(51, 51, 51)
        else:
            para.font.size = Pt(12)
            para.font.color.rgb = RGBColor(102, 102, 102)
        para.space_after = Pt(2)

    # 기대 효과 섹션
    effect_title_box = slide1.shapes.add_textbox(Inches(5.2), Inches(3.2), Inches(4), Inches(0.4))
    effect_title_frame = effect_title_box.text_frame
    effect_title_frame.text = "✨ 기대 효과"
    effect_title_para = effect_title_frame.paragraphs[0]
    effect_title_para.font.size = Pt(20)
    effect_title_para.font.bold = True
    effect_title_para.font.color.rgb = RGBColor(0, 153, 51)

    # 기대 효과 상세
    effect_box = slide1.shapes.add_textbox(Inches(5.2), Inches(3.7), Inches(4), Inches(2))
    effect_frame = effect_box.text_frame
    effect_frame.word_wrap = True

    effects = [
        "✓ 기회손실 최소화",
        "  → 손실시간 5% 감소 목표",
        "",
        "✓ 설비 CAPA 증가",
        "  → 생산성 향상 및 효율 개선",
        "",
        "✓ OVERTIME 감소",
        "  → 인력운영비 절감",
        "",
        "✓ 점당 가공비 10% 절감",
        "  → OH 감소 달성"
    ]

    for i, eff in enumerate(effects):
        if i == 0:
            para = effect_frame.paragraphs[0]
        else:
            para = effect_frame.add_paragraph()
        para.text = eff
        if eff.startswith("✓"):
            para.font.size = Pt(14)
            para.font.bold = True
            para.font.color.rgb = RGBColor(0, 102, 51)
        else:
            para.font.size = Pt(12)
            para.font.color.rgb = RGBColor(102, 102, 102)
        para.space_after = Pt(2)

    # 페이지 번호
    page_num_box = slide1.shapes.add_textbox(Inches(9.2), Inches(7.1), Inches(0.6), Inches(0.3))
    page_num_frame = page_num_box.text_frame
    page_num_frame.text = "1"
    page_num_para = page_num_frame.paragraphs[0]
    page_num_para.font.size = Pt(12)
    page_num_para.font.color.rgb = RGBColor(128, 128, 128)
    page_num_para.alignment = PP_ALIGN.CENTER

    # ========== 슬라이드 2: 전략 과제 및 실행 계획 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경색
    background2 = slide2.background
    fill2 = background2.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(255, 255, 255)

    # 제목
    title2_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title2_frame = title2_box.text_frame
    title2_frame.text = "전략 과제 및 실행 계획"
    title2_para = title2_frame.paragraphs[0]
    title2_para.font.size = Pt(36)
    title2_para.font.bold = True
    title2_para.font.color.rgb = RGBColor(0, 51, 102)
    title2_para.alignment = PP_ALIGN.CENTER

    # 전략 1: 손실 시간 제로화 프로젝트
    strategy1_box = slide2.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(1.2), Inches(9), Inches(1.8)
    )
    strategy1_box.fill.solid()
    strategy1_box.fill.fore_color.rgb = RGBColor(230, 240, 255)
    strategy1_box.line.color.rgb = RGBColor(0, 102, 204)
    strategy1_box.line.width = Pt(3)

    s1_frame = strategy1_box.text_frame
    s1_frame.margin_left = Inches(0.2)
    s1_frame.margin_top = Inches(0.15)
    s1_frame.word_wrap = True

    s1_title = s1_frame.paragraphs[0]
    s1_title.text = "전략 1️⃣  손실 시간 제로화 프로젝트"
    s1_title.font.size = Pt(22)
    s1_title.font.bold = True
    s1_title.font.color.rgb = RGBColor(0, 51, 153)

    s1_content = [
        "▶ 핵심 액션",
        "   • 자동분석 TOOL 제작: MES DATA 활용한 이상 감지, C/T 변화, 모니터링 활성화",
        "   • 순간유실 가시화: 초 단위 손실 기록 → 세부 항목 원터치 이력 저장 (모바일 앱)",
        "   • TOP 10 집중 개선: LINE별 일/주/월 손실 분석 → 우선순위 타겟 개선",
        "",
        "📊 KPI:  손실시간 5% 감소  |  기회손실 최소화 → 점당 가공비 직접 절감 효과"
    ]

    for content in s1_content:
        para = s1_frame.add_paragraph()
        para.text = content
        if content.startswith("▶"):
            para.font.size = Pt(16)
            para.font.bold = True
            para.font.color.rgb = RGBColor(0, 102, 204)
        elif content.startswith("📊"):
            para.font.size = Pt(14)
            para.font.bold = True
            para.font.color.rgb = RGBColor(204, 0, 0)
        else:
            para.font.size = Pt(13)
            para.font.color.rgb = RGBColor(51, 51, 51)
        para.space_after = Pt(3)

    # 전략 2: 불량 재발 Zero 챌린지
    strategy2_box = slide2.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(3.2), Inches(9), Inches(1.8)
    )
    strategy2_box.fill.solid()
    strategy2_box.fill.fore_color.rgb = RGBColor(255, 240, 230)
    strategy2_box.line.color.rgb = RGBColor(255, 102, 0)
    strategy2_box.line.width = Pt(3)

    s2_frame = strategy2_box.text_frame
    s2_frame.margin_left = Inches(0.2)
    s2_frame.margin_top = Inches(0.15)
    s2_frame.word_wrap = True

    s2_title = s2_frame.paragraphs[0]
    s2_title.text = "전략 2️⃣  불량 재발 Zero 챌린지"
    s2_title.font.size = Pt(22)
    s2_title.font.bold = True
    s2_title.font.color.rgb = RGBColor(204, 51, 0)

    s2_content = [
        "▶ 핵심 액션",
        "   • 품질 즉시 FEEDBACK: 불량 사진 + MES 이력정보 (모델/일자/LINE/담당자) 즉시 공유",
        "   • 전 조 자동 알람: 스마트폰 사진 업로드 → 공정/설비 태그 → 조치내용 공유 → 재발 방지",
        "   • 재발 불량 추적: 동일 불량 발생 시 자동 '재발' 표시 → 월별 추적 및 Zero 목표 KPI화",
        "",
        "📊 KPI:  재발 불량 30% 감소  |  품질 불량 10% 감소 목표 달성"
    ]

    for content in s2_content:
        para = s2_frame.add_paragraph()
        para.text = content
        if content.startswith("▶"):
            para.font.size = Pt(16)
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 102, 0)
        elif content.startswith("📊"):
            para.font.size = Pt(14)
            para.font.bold = True
            para.font.color.rgb = RGBColor(204, 0, 0)
        else:
            para.font.size = Pt(13)
            para.font.color.rgb = RGBColor(51, 51, 51)
        para.space_after = Pt(3)

    # 전략 3: 설비 CAPA 증가
    strategy3_box = slide2.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(5.2), Inches(9), Inches(1.6)
    )
    strategy3_box.fill.solid()
    strategy3_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    strategy3_box.line.color.rgb = RGBColor(0, 153, 51)
    strategy3_box.line.width = Pt(3)

    s3_frame = strategy3_box.text_frame
    s3_frame.margin_left = Inches(0.2)
    s3_frame.margin_top = Inches(0.15)
    s3_frame.word_wrap = True

    s3_title = s3_frame.paragraphs[0]
    s3_title.text = "전략 3️⃣  설비 CAPA 증가 및 공정 최적화"
    s3_title.font.size = Pt(22)
    s3_title.font.bold = True
    s3_title.font.color.rgb = RGBColor(0, 102, 51)

    s3_content = [
        "▶ 핵심 액션",
        "   • C/T 단축 개선: 현 P/G 운영 DEEP 분석 → 최단거리 단 1초라도 줄이기 위한 활동",
        "   • 기본 BASE 유지: 설비 효율 향상 위한 성능/PM 관련 모든 활동 반복 운영",
        "   • 기타 과제: RADIAL2 수삽설비 조기 안정화 / SMD 공정 LAY OUT 최적화",
        "",
        "📊 KPI:  CAPA 증가  |  OVERTIME 감소  |  OH 감소 달성"
    ]

    for content in s3_content:
        para = s3_frame.add_paragraph()
        para.text = content
        if content.startswith("▶"):
            para.font.size = Pt(16)
            para.font.bold = True
            para.font.color.rgb = RGBColor(0, 153, 51)
        elif content.startswith("📊"):
            para.font.size = Pt(14)
            para.font.bold = True
            para.font.color.rgb = RGBColor(204, 0, 0)
        else:
            para.font.size = Pt(13)
            para.font.color.rgb = RGBColor(51, 51, 51)
        para.space_after = Pt(3)

    # 하단 팀 정보
    team_box = slide2.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(8.5), Inches(0.3))
    team_frame = team_box.text_frame
    team_frame.text = "제조1팀  |  2026년 경영전략"
    team_para = team_frame.paragraphs[0]
    team_para.font.size = Pt(12)
    team_para.font.color.rgb = RGBColor(128, 128, 128)
    team_para.alignment = PP_ALIGN.LEFT

    # 페이지 번호
    page2_num_box = slide2.shapes.add_textbox(Inches(9.2), Inches(7.1), Inches(0.6), Inches(0.3))
    page2_num_frame = page2_num_box.text_frame
    page2_num_frame.text = "2"
    page2_num_para = page2_num_frame.paragraphs[0]
    page2_num_para.font.size = Pt(12)
    page2_num_para.font.color.rgb = RGBColor(128, 128, 128)
    page2_num_para.alignment = PP_ALIGN.CENTER

    # 저장
    output_file = '2026년_경영전략_제조1팀.pptx'
    prs.save(output_file)
    print(f"✅ PPT 파일이 성공적으로 생성되었습니다: {output_file}")
    print(f"📄 총 슬라이드 수: {len(prs.slides)}")
    return output_file

if __name__ == "__main__":
    create_strategy_presentation()
