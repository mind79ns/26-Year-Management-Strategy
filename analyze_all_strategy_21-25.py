#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
전략참고자료 폴더의 모든 PPT 심층 분석 (21-25년 전체)
- MAIN 브랜치의 11개 전략 PPT 파일 분석
- 모든 슬라이드의 모든 텍스트 추출
- 페이지별 주제 분류
- 핵심 KPI 및 목표 추출
- 과제 및 실행 계획 추출
- 연도별 트렌드 분석
"""

from pptx import Presentation
import os
import json
import re
from collections import defaultdict, Counter

def extract_detailed_content(ppt_path):
    """PPT 파일에서 상세 내용 추출"""
    try:
        prs = Presentation(ppt_path)
        file_info = {
            "file_name": os.path.basename(ppt_path),
            "file_path": ppt_path,
            "total_slides": len(prs.slides),
            "slides": []
        }

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_info = {
                "slide_number": slide_num,
                "texts": [],
                "all_text": "",
                "bullets": [],
                "numbers": [],
                "keywords": []
            }

            # 모든 텍스트 수집
            all_texts = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    all_texts.append(text)
                    slide_info["texts"].append(text)

                    # 텍스트 프레임이 있는 경우 단락별 분석
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = paragraph.text.strip()
                            if para_text:
                                # 불릿 포인트 감지
                                if para_text.startswith(('•', '-', '·', '※', '→', '▶', '□', '■', '○', '●', '☑')):
                                    slide_info["bullets"].append(para_text)

            # 전체 텍스트 통합
            slide_info["all_text"] = " ".join(all_texts)

            # 숫자 패턴 추출 (목표, KPI 등)
            numbers = re.findall(r'\d+(?:\.\d+)?(?:%|개|명|시간|분|초|건|회|대|억|만|천)', slide_info["all_text"])
            slide_info["numbers"] = numbers

            # 키워드 추출 (한글 2글자 이상)
            korean_words = re.findall(r'[가-힣]{2,}', slide_info["all_text"])
            slide_info["keywords"] = korean_words

            file_info["slides"].append(slide_info)

        return file_info

    except Exception as e:
        print(f"Error processing {ppt_path}: {e}")
        return None

def analyze_themes(file_info):
    """슬라이드별 주제 분류"""
    theme_keywords = {
        "목표": ["목표", "Target", "전략", "방향", "비전"],
        "유실시간": ["유실", "Loss", "정지", "비가동", "가동률"],
        "불량": ["불량", "품질", "Defect", "PPM", "재발"],
        "설비": ["설비", "Equipment", "CAPA", "능력", "생산력"],
        "인력": ["인력", "인원", "작업자", "교육", "훈련"],
        "원가": ["원가", "비용", "Cost", "절감", "가공비"],
        "생산성": ["생산성", "효율", "Productivity", "개선"],
        "자동화": ["자동화", "Automation", "스마트", "MES", "시스템"],
        "실행계획": ["실행", "계획", "일정", "로드맵", "추진"],
        "성과": ["성과", "결과", "달성", "실적"],
        "안전": ["안전", "Safety", "재해", "사고"],
        "환경": ["환경", "Environment", "에너지", "탄소"]
    }

    for slide in file_info["slides"]:
        slide["themes"] = []
        text = slide["all_text"].lower()

        for theme, keywords in theme_keywords.items():
            for keyword in keywords:
                if keyword.lower() in text:
                    if theme not in slide["themes"]:
                        slide["themes"].append(theme)
                    break

def extract_tasks_and_kpis(file_info):
    """과제 및 KPI 추출"""
    tasks = []
    kpis = []

    for slide in file_info["slides"]:
        # 과제 패턴 추출
        for bullet in slide["bullets"]:
            # 과제로 보이는 패턴
            if any(keyword in bullet for keyword in ["개선", "구축", "도입", "추진", "실행", "수립", "강화", "확대", "고도화", "혁신"]):
                tasks.append({
                    "slide": slide["slide_number"],
                    "task": bullet,
                    "themes": slide["themes"]
                })

        # KPI 패턴 추출 (숫자 포함)
        for text in slide["texts"]:
            if any(char in text for char in ['%', '→', '▶']) and any(char.isdigit() for char in text):
                kpis.append({
                    "slide": slide["slide_number"],
                    "kpi": text,
                    "numbers": slide["numbers"]
                })

    return tasks, kpis

def extract_year_from_filename(filename):
    """파일명에서 연도 추출"""
    match = re.search(r'(\d{2})년', filename)
    if match:
        year = int(match.group(1))
        return 2000 + year
    return None

def main():
    """메인 분석 실행"""
    ppt_folder = "전략자료참고"

    # PPT 파일 자동 검색
    ppt_files = []
    for file in os.listdir(ppt_folder):
        if file.endswith('.pptx') and not file.startswith('~'):
            ppt_files.append(file)

    # 파일명으로 정렬 (연도순)
    ppt_files.sort()

    all_analysis = {
        "total_files": len(ppt_files),
        "total_slides": 0,
        "files": [],
        "all_tasks": [],
        "all_kpis": [],
        "keyword_frequency": {},
        "theme_distribution": defaultdict(int),
        "yearly_summary": {}
    }

    print("=" * 70)
    print("전략참고자료 심층 분석 시작 (21-25년 전체)")
    print("=" * 70)

    yearly_data = defaultdict(lambda: {
        "files": [],
        "total_slides": 0,
        "tasks": [],
        "kpis": [],
        "keywords": Counter()
    })

    for ppt_file in ppt_files:
        ppt_path = os.path.join(ppt_folder, ppt_file)

        if not os.path.exists(ppt_path):
            print(f"⚠ 파일 없음: {ppt_file}")
            continue

        print(f"\n📄 분석 중: {ppt_file}")

        # 상세 내용 추출
        file_info = extract_detailed_content(ppt_path)

        if file_info:
            # 주제 분석
            analyze_themes(file_info)

            # 과제 및 KPI 추출
            tasks, kpis = extract_tasks_and_kpis(file_info)

            file_info["tasks"] = tasks
            file_info["kpis"] = kpis

            # 연도 추출
            year = extract_year_from_filename(ppt_file)
            file_info["year"] = year

            all_analysis["files"].append(file_info)
            all_analysis["total_slides"] += file_info["total_slides"]
            all_analysis["all_tasks"].extend(tasks)
            all_analysis["all_kpis"].extend(kpis)

            # 연도별 데이터 집계
            if year:
                yearly_data[year]["files"].append(ppt_file)
                yearly_data[year]["total_slides"] += file_info["total_slides"]
                yearly_data[year]["tasks"].extend(tasks)
                yearly_data[year]["kpis"].extend(kpis)

                for slide in file_info["slides"]:
                    for keyword in slide["keywords"]:
                        if len(keyword) >= 2:
                            yearly_data[year]["keywords"][keyword] += 1

            # 주제 분포 집계
            for slide in file_info["slides"]:
                for theme in slide["themes"]:
                    all_analysis["theme_distribution"][theme] += 1

            # 키워드 빈도 집계
            for slide in file_info["slides"]:
                for keyword in slide["keywords"]:
                    if len(keyword) >= 2:
                        all_analysis["keyword_frequency"][keyword] = \
                            all_analysis["keyword_frequency"].get(keyword, 0) + 1

            print(f"   ✓ {file_info['total_slides']}개 슬라이드 분석 완료")
            print(f"   ✓ {len(tasks)}개 과제 추출")
            print(f"   ✓ {len(kpis)}개 KPI 추출")

    # 연도별 요약 정리
    for year in sorted(yearly_data.keys()):
        data = yearly_data[year]
        all_analysis["yearly_summary"][str(year)] = {
            "files": data["files"],
            "total_slides": data["total_slides"],
            "total_tasks": len(data["tasks"]),
            "total_kpis": len(data["kpis"]),
            "top_keywords": data["keywords"].most_common(20)
        }

    # 키워드 빈도 상위 정렬
    sorted_keywords = sorted(
        all_analysis["keyword_frequency"].items(),
        key=lambda x: x[1],
        reverse=True
    )
    all_analysis["top_keywords"] = sorted_keywords[:100]

    # 주제 분포 정렬
    all_analysis["theme_distribution"] = dict(
        sorted(all_analysis["theme_distribution"].items(),
               key=lambda x: x[1],
               reverse=True)
    )

    # JSON 저장
    output_file = "전략자료_전체분석_21-25년.json"
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(all_analysis, f, ensure_ascii=False, indent=2)

    print("\n" + "=" * 70)
    print("✅ 분석 완료")
    print("=" * 70)
    print(f"총 파일: {all_analysis['total_files']}개")
    print(f"총 슬라이드: {all_analysis['total_slides']}개")
    print(f"총 과제: {len(all_analysis['all_tasks'])}개")
    print(f"총 KPI: {len(all_analysis['all_kpis'])}개")
    print(f"고유 키워드: {len(all_analysis['keyword_frequency'])}개")
    print(f"\n결과 저장: {output_file}")
    print("=" * 70)

    # 연도별 통계
    print("\n📅 연도별 통계:")
    for year in sorted(yearly_data.keys()):
        data = yearly_data[year]
        print(f"\n{year}년:")
        print(f"   파일: {len(data['files'])}개")
        print(f"   슬라이드: {data['total_slides']}개")
        print(f"   과제: {len(data['tasks'])}개")
        print(f"   KPI: {len(data['kpis'])}개")
        print(f"   상위 키워드: {', '.join([kw for kw, _ in data['keywords'].most_common(5)])}")

    # 주요 통계 출력
    print("\n📊 주제 분포 (상위 10개):")
    for i, (theme, count) in enumerate(list(all_analysis["theme_distribution"].items())[:10], 1):
        print(f"   {i}. {theme}: {count}회")

    print("\n🔑 고빈도 키워드 (상위 20개):")
    for i, (keyword, count) in enumerate(all_analysis["top_keywords"][:20], 1):
        print(f"   {i}. {keyword}: {count}회")

    return all_analysis

if __name__ == "__main__":
    analysis_result = main()
