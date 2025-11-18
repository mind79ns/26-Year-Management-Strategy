#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Ultra Premium Professional Design
최신 고급 비즈니스 프레젠테이션 스타일
- 다크 네이비 & 골드 컬러
- 모던 그라데이션
- 세련된 타이포그래피
- 비대칭 레이아웃
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_premium_gradient(shape, color1, color2, angle=90):
    """프리미엄 그라데이션 추가"""
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = angle
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def create_premium_shape(slide, left, top, width, height, shape_type, color, shadow=True):
    """그림자가 있는 프리미엄 도형"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    if shadow:
        shape.shadow.inherit = False
        shape.shadow.visible = True
        shape.shadow.distance = Pt(3)
        shape.shadow.angle = 45
        shape.shadow.blur_radius = Pt(8)

    return shape

def create_ultra_premium_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 프리미엄 색상 팔레트
    DARK_NAVY = RGBColor(15, 32, 56)      # 다크 네이비
    NAVY = RGBColor(25, 55, 109)          # 네이비
    GOLD = RGBColor(212, 175, 55)         # 골드
    LIGHT_GOLD = RGBColor(255, 223, 128)  # 라이트 골드
    SILVER = RGBColor(192, 192, 192)      # 실버
    WHITE = RGBColor(255, 255, 255)       # 화이트
    OFF_WHITE = RGBColor(248, 248, 248)   # 오프화이트
    BLUE_GRAY = RGBColor(96, 125, 139)    # 블루그레이
    ACCENT_BLUE = RGBColor(41, 128, 185)  # 액센트 블루
    ACCENT_GREEN = RGBColor(39, 174, 96)  # 액센트 그린
    DARK_TEXT = RGBColor(33, 33, 33)      # 다크 텍스트

    # ========== 슬라이드 1: 프리미엄 표지 ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    # 다크 네이비 배경
    bg = slide1.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_NAVY

    # 좌측 골드 액센트 바 (세로)
    accent_bar = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.15), Inches(7.5)
    )
    add_premium_gradient(accent_bar, GOLD, LIGHT_GOLD, 0)
    accent_bar.line.fill.background()

    # 대형 년도 (고급 타이포그래피)
    year_box = slide1.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(1.5))
    year_frame = year_box.text_frame
    year_para = year_frame.paragraphs[0]
    year_para.text = "2026"
    year_para.font.size = Pt(140)
    year_para.font.bold = True
    year_para.font.color.rgb = RGBColor(35, 62, 96)  # 약간 밝은 네이비 (워터마크 효과)
    year_para.alignment = PP_ALIGN.LEFT

    # 메인 타이틀
    title_box = slide1.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(7), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "경영전략"
    title_para.font.size = Pt(68)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.LEFT

    # 서브타이틀 (골드 라인과 함께)
    gold_line = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(3.9), Inches(4), Inches(0.03)
    )
    gold_line.fill.solid()
    gold_line.fill.fore_color.rgb = GOLD
    gold_line.line.fill.background()

    subtitle_box = slide1.shapes.add_textbox(Inches(1.5), Inches(4.1), Inches(7), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = "제조1팀  |  자동화 제조라인 스마트화"
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = SILVER
    subtitle_para.alignment = PP_ALIGN.LEFT

    # 우하단: 핵심 지표 카드 (3개)
    kpi_y = Inches(5.3)
    kpis = [
        {"label": "가공비", "value": "10%", "icon": "↓"},
        {"label": "품질", "value": "10%", "icon": "↑"},
        {"label": "효율", "value": "5%", "icon": "↑"}
    ]

    card_width = Inches(2)
    card_height = Inches(1.4)
    card_spacing = Inches(0.25)
    card_start = Inches(1.5)

    for i, kpi in enumerate(kpis):
        x_pos = card_start + (i * (card_width + card_spacing))

        # 카드 박스 (그림자 있음)
        card = create_premium_shape(
            slide1, x_pos, kpi_y, card_width, card_height,
            MSO_SHAPE.ROUNDED_RECTANGLE, RGBColor(30, 55, 85), shadow=True
        )

        # 라벨
        label_box = slide1.shapes.add_textbox(x_pos + Inches(0.2), kpi_y + Inches(0.2), card_width - Inches(0.4), Inches(0.3))
        label_frame = label_box.text_frame
        label_para = label_frame.paragraphs[0]
        label_para.text = kpi["label"]
        label_para.font.size = Pt(14)
        label_para.font.color.rgb = SILVER
        label_para.alignment = PP_ALIGN.LEFT

        # 값
        value_box = slide1.shapes.add_textbox(x_pos + Inches(0.2), kpi_y + Inches(0.55), card_width - Inches(0.4), Inches(0.6))
        value_frame = value_box.text_frame
        value_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        value_para = value_frame.paragraphs[0]
        value_para.text = f"{kpi['icon']} {kpi['value']}"
        value_para.font.size = Pt(36)
        value_para.font.bold = True
        value_para.font.color.rgb = GOLD
        value_para.alignment = PP_ALIGN.LEFT

    # ========== 슬라이드 2: 현황 분석 (프리미엄) ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = OFF_WHITE

    # 좌측 골드 액센트
    accent = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.08), Inches(7.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    # 헤더 영역
    header_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    header_frame = header_box.text_frame

    # 페이지 번호 (작게)
    page_para = header_frame.paragraphs[0]
    page_para.text = "01"
    page_para.font.size = Pt(14)
    page_para.font.color.rgb = BLUE_GRAY
    page_para.alignment = PP_ALIGN.LEFT
    page_para.space_after = Pt(5)

    # 제목
    title_para = header_frame.add_paragraph()
    title_para.text = "현황 분석 및 전략 방향"
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_NAVY
    title_para.alignment = PP_ALIGN.LEFT

    # 언더라인
    underline = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.4), Inches(2.5), Inches(0.04)
    )
    add_premium_gradient(underline, GOLD, LIGHT_GOLD, 90)
    underline.line.fill.background()

    # 좌측: 문제점 (다크 카드)
    problem_card = create_premium_shape(
        slide2, Inches(0.5), Inches(2), Inches(4.3), Inches(4.8),
        MSO_SHAPE.ROUNDED_RECTANGLE, DARK_NAVY, shadow=True
    )

    prob_frame = problem_card.text_frame
    prob_frame.margin_left = Inches(0.35)
    prob_frame.margin_right = Inches(0.35)
    prob_frame.margin_top = Inches(0.3)
    prob_frame.margin_bottom = Inches(0.3)

    # 타이틀
    prob_title = prob_frame.paragraphs[0]
    prob_title.text = "Current Challenges"
    prob_title.font.size = Pt(11)
    prob_title.font.color.rgb = SILVER
    prob_title.alignment = PP_ALIGN.LEFT
    prob_title.space_after = Pt(5)

    prob_title_ko = prob_frame.add_paragraph()
    prob_title_ko.text = "현재 문제점"
    prob_title_ko.font.size = Pt(24)
    prob_title_ko.font.bold = True
    prob_title_ko.font.color.rgb = WHITE
    prob_title_ko.space_after = Pt(20)

    # 문제점 리스트
    problems = [
        {"icon": "⚠", "title": "느린 대응", "desc": "설비 이상 발견 → DATA 집계 → 대응\n기회 손실 지속 발생"},
        {"icon": "↻", "title": "반복 작업", "desc": "수동적 원인 분석\n업무 비효율 심화"},
        {"icon": "!", "title": "품질 문제", "desc": "불량 원인 대응 지연\n재발 방지 체계 미흡"}
    ]

    for prob in problems:
        # 아이콘
        icon_para = prob_frame.add_paragraph()
        icon_para.text = prob["icon"]
        icon_para.font.size = Pt(20)
        icon_para.font.color.rgb = GOLD
        icon_para.space_after = Pt(5)

        # 제목
        title_para = prob_frame.add_paragraph()
        title_para.text = prob["title"]
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        title_para.space_after = Pt(6)

        # 설명
        desc_para = prob_frame.add_paragraph()
        desc_para.text = prob["desc"]
        desc_para.font.size = Pt(12)
        desc_para.font.color.rgb = SILVER
        desc_para.space_after = Pt(18)

    # 우측 상단: 목표 (화이트 카드)
    goal_card = create_premium_shape(
        slide2, Inches(5.1), Inches(2), Inches(4.4), Inches(2.2),
        MSO_SHAPE.ROUNDED_RECTANGLE, WHITE, shadow=True
    )

    goal_frame = goal_card.text_frame
    goal_frame.margin_left = Inches(0.35)
    goal_frame.margin_top = Inches(0.25)

    goal_label = goal_frame.paragraphs[0]
    goal_label.text = "2026 Target"
    goal_label.font.size = Pt(11)
    goal_label.font.color.rgb = BLUE_GRAY
    goal_label.space_after = Pt(5)

    goal_title = goal_frame.add_paragraph()
    goal_title.text = "핵심 목표"
    goal_title.font.size = Pt(22)
    goal_title.font.bold = True
    goal_title.font.color.rgb = DARK_NAVY
    goal_title.space_after = Pt(15)

    goals = ["가공비 10% 절감", "품질 불량 10% 감소", "손실 시간 5% 감소"]
    for goal in goals:
        para = goal_frame.add_paragraph()
        para.text = f"▪ {goal}"
        para.font.size = Pt(14)
        para.font.color.rgb = DARK_TEXT
        para.space_after = Pt(8)

    # 우측 하단: 기대효과 (골드 액센트 카드)
    effect_card = create_premium_shape(
        slide2, Inches(5.1), Inches(4.5), Inches(4.4), Inches(2.3),
        MSO_SHAPE.ROUNDED_RECTANGLE, WHITE, shadow=True
    )

    # 좌측 골드 바
    gold_bar = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.1), Inches(4.5), Inches(0.08), Inches(2.3)
    )
    gold_bar.fill.solid()
    gold_bar.fill.fore_color.rgb = GOLD
    gold_bar.line.fill.background()

    effect_frame = effect_card.text_frame
    effect_frame.margin_left = Inches(0.45)
    effect_frame.margin_top = Inches(0.25)

    effect_label = effect_frame.paragraphs[0]
    effect_label.text = "Expected Impact"
    effect_label.font.size = Pt(11)
    effect_label.font.color.rgb = BLUE_GRAY
    effect_label.space_after = Pt(5)

    effect_title = effect_frame.add_paragraph()
    effect_title.text = "기대 효과"
    effect_title.font.size = Pt(22)
    effect_title.font.bold = True
    effect_title.font.color.rgb = DARK_NAVY
    effect_title.space_after = Pt(15)

    effects = ["기회손실 최소화", "설비 CAPA 증가", "OVERTIME 감소", "점당 가공비 직접 절감"]
    for eff in effects:
        para = effect_frame.add_paragraph()
        para.text = f"→ {eff}"
        para.font.size = Pt(13)
        para.font.color.rgb = DARK_TEXT
        para.space_after = Pt(8)

    # 페이지 번호
    page_num = slide2.shapes.add_textbox(Inches(9.2), Inches(7.1), Inches(0.6), Inches(0.3))
    page_num.text_frame.text = "02"
    page_num.text_frame.paragraphs[0].font.size = Pt(11)
    page_num.text_frame.paragraphs[0].font.color.rgb = BLUE_GRAY
    page_num.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # ========== 슬라이드 3: 전략 1 (프리미엄) ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = OFF_WHITE

    # 액센트 바
    accent = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_BLUE
    accent.line.fill.background()

    # 헤더
    header = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    h_frame = header.text_frame

    h_page = h_frame.paragraphs[0]
    h_page.text = "02"
    h_page.font.size = Pt(14)
    h_page.font.color.rgb = BLUE_GRAY
    h_page.space_after = Pt(5)

    h_title = h_frame.add_paragraph()
    h_title.text = "전략 1   손실 시간 제로화 프로젝트"
    h_title.font.size = Pt(32)
    h_title.font.bold = True
    h_title.font.color.rgb = DARK_NAVY

    # 언더라인
    underline = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(2.5), Inches(0.04))
    add_premium_gradient(underline, ACCENT_BLUE, RGBColor(100, 181, 246), 90)
    underline.line.fill.background()

    # 전략 개요 카드
    overview = create_premium_shape(
        slide3, Inches(0.5), Inches(1.8), Inches(9), Inches(1),
        MSO_SHAPE.ROUNDED_RECTANGLE, WHITE, shadow=True
    )

    ov_frame = overview.text_frame
    ov_frame.margin_left = Inches(0.3)
    ov_frame.margin_top = Inches(0.2)

    ov_para = ov_frame.paragraphs[0]
    ov_para.text = "전략 개요  "
    ov_para.font.size = Pt(16)
    ov_para.font.bold = True
    ov_para.font.color.rgb = DARK_NAVY

    ov_desc = ov_frame.add_paragraph()
    ov_desc.text = "MES DATA 활용 자동분석 TOOL 구축 → 순간유실 가시화 → 초 단위 손실 기록 → 우선순위 집중 개선"
    ov_desc.font.size = Pt(14)
    ov_desc.font.color.rgb = DARK_TEXT
    ov_desc.space_before = Pt(5)

    # 3개 액션 카드 (세로 배치)
    actions = [
        {
            "num": "01", "title": "자동분석 TOOL 제작",
            "items": ["MES DATA 활용 이상 감지", "C/T 변화 모니터링 활성화", "실시간 알람 시스템 구축"]
        },
        {
            "num": "02", "title": "순간유실 가시화",
            "items": ["초 단위 손실 기록", "모바일 앱 원터치 입력", "LINE별 일/주/월 분석"]
        },
        {
            "num": "03", "title": "TOP 10 집중 개선",
            "items": ["손실 항목 순위화", "WORST 품목 집중 타격", "주간 리포트 자동화"]
        }
    ]

    action_width = Inches(2.8)
    action_height = Inches(3.4)
    action_spacing = Inches(0.3)
    action_start = Inches(0.5)
    action_y = Inches(3.1)

    for i, action in enumerate(actions):
        x_pos = action_start + (i * (action_width + action_spacing))

        # 액션 카드
        card = create_premium_shape(
            slide3, x_pos, action_y, action_width, action_height,
            MSO_SHAPE.ROUNDED_RECTANGLE, WHITE, shadow=True
        )

        # 상단 액센트 (블루)
        top_accent = slide3.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, action_y, action_width, Inches(0.6)
        )
        top_accent.fill.solid()
        top_accent.fill.fore_color.rgb = ACCENT_BLUE
        top_accent.line.fill.background()

        # 번호 (흰색)
        num_box = slide3.shapes.add_textbox(x_pos, action_y, action_width, Inches(0.6))
        num_frame = num_box.text_frame
        num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        num_para = num_frame.paragraphs[0]
        num_para.text = action["num"]
        num_para.font.size = Pt(28)
        num_para.font.bold = True
        num_para.font.color.rgb = WHITE
        num_para.alignment = PP_ALIGN.CENTER

        # 제목
        title_box = slide3.shapes.add_textbox(x_pos + Inches(0.2), action_y + Inches(0.8), action_width - Inches(0.4), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = action["title"]
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = DARK_NAVY
        title_para.alignment = PP_ALIGN.CENTER
        title_para.space_after = Pt(10)

        # 항목들
        items_box = slide3.shapes.add_textbox(x_pos + Inches(0.2), action_y + Inches(1.5), action_width - Inches(0.4), Inches(1.6))
        items_frame = items_box.text_frame
        items_frame.word_wrap = True

        for j, item in enumerate(action["items"]):
            if j == 0:
                para = items_frame.paragraphs[0]
            else:
                para = items_frame.add_paragraph()
            para.text = f"• {item}"
            para.font.size = Pt(11)
            para.font.color.rgb = DARK_TEXT
            para.space_after = Pt(6)

    # KPI 박스
    kpi_box = create_premium_shape(
        slide3, Inches(0.5), Inches(6.7), Inches(9), Inches(0.6),
        MSO_SHAPE.ROUNDED_RECTANGLE, DARK_NAVY, shadow=False
    )

    kpi_frame = kpi_box.text_frame
    kpi_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    kpi_para = kpi_frame.paragraphs[0]
    kpi_para.text = "KPI   손실시간 5% 감소 → 점당 가공비 직접 절감 효과"
    kpi_para.font.size = Pt(18)
    kpi_para.font.bold = True
    kpi_para.font.color.rgb = WHITE
    kpi_para.alignment = PP_ALIGN.CENTER

    # 골드 라인 추가
    gold_accent = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(6.7), Inches(9), Inches(0.05)
    )
    gold_accent.fill.solid()
    gold_accent.fill.fore_color.rgb = GOLD
    gold_accent.line.fill.background()

    # 페이지 번호
    page_num = slide3.shapes.add_textbox(Inches(9.2), Inches(7.1), Inches(0.6), Inches(0.3))
    page_num.text_frame.text = "03"
    page_num.text_frame.paragraphs[0].font.size = Pt(11)
    page_num.text_frame.paragraphs[0].font.color.rgb = BLUE_GRAY
    page_num.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # ========== 슬라이드 4-6: 동일한 프리미엄 스타일로 전략 2, 3, 로드맵 생성 ==========
    # (간결성을 위해 전략 2, 3은 유사한 구조로 생성)

    # 전략 2
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    slide4.background.fill.solid()
    slide4.background.fill.fore_color.rgb = OFF_WHITE

    accent4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
    accent4.fill.solid()
    accent4.fill.fore_color.rgb = ACCENT_GREEN
    accent4.line.fill.background()

    header4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    h4_frame = header4.text_frame
    h4_page = h4_frame.paragraphs[0]
    h4_page.text = "03"
    h4_page.font.size = Pt(14)
    h4_page.font.color.rgb = BLUE_GRAY
    h4_page.space_after = Pt(5)

    h4_title = h4_frame.add_paragraph()
    h4_title.text = "전략 2   불량 재발 Zero 챌린지"
    h4_title.font.size = Pt(32)
    h4_title.font.bold = True
    h4_title.font.color.rgb = DARK_NAVY

    underline4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(2.5), Inches(0.04))
    add_premium_gradient(underline4, ACCENT_GREEN, RGBColor(169, 223, 191), 90)
    underline4.line.fill.background()

    overview4 = create_premium_shape(
        slide4, Inches(0.5), Inches(1.8), Inches(9), Inches(1),
        MSO_SHAPE.ROUNDED_RECTANGLE, WHITE, shadow=True
    )

    ov4_frame = overview4.text_frame
    ov4_frame.margin_left = Inches(0.3)
    ov4_frame.margin_top = Inches(0.2)

    ov4_para = ov4_frame.paragraphs[0]
    ov4_para.text = "전략 개요  "
    ov4_para.font.size = Pt(16)
    ov4_para.font.bold = True
    ov4_para.font.color.rgb = DARK_NAVY

    ov4_desc = ov4_frame.add_paragraph()
    ov4_desc.text = "불량 사진 즉시 공유 시스템 → 전 조 자동 알람 → 재발 불량 추적 강화 → 조치사항 DB 구축 → 재발 Zero"
    ov4_desc.font.size = Pt(14)
    ov4_desc.font.color.rgb = DARK_TEXT
    ov4_desc.space_before = Pt(5)

    actions4 = [
        {"num": "01", "title": "품질 즉시 FEEDBACK", "items": ["불량 사진 + MES 이력 연동", "모델/일자/LINE/담당자 자동 기록", "스마트폰 즉시 업로드"]},
        {"num": "02", "title": "전 조 자동 알람", "items": ["공정/설비 태그 자동 분류", "조치내용 한 줄 메모 공유", "주간조 → 야간조 자동 전달"]},
        {"num": "03", "title": "재발 불량 추적", "items": ["동일 불량 자동 '재발' 표시", "월별 재발 불량 추적", "재발 Zero KPI화"]}
    ]

    for i, action in enumerate(actions4):
        x_pos = action_start + (i * (action_width + action_spacing))

        card = create_premium_shape(slide4, x_pos, action_y, action_width, action_height,
                                    MSO_SHAPE.ROUNDED_RECTANGLE, WHITE, shadow=True)

        top_accent = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, action_y, action_width, Inches(0.6))
        top_accent.fill.solid()
        top_accent.fill.fore_color.rgb = ACCENT_GREEN
        top_accent.line.fill.background()

        num_box = slide4.shapes.add_textbox(x_pos, action_y, action_width, Inches(0.6))
        num_frame = num_box.text_frame
        num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        num_para = num_frame.paragraphs[0]
        num_para.text = action["num"]
        num_para.font.size = Pt(28)
        num_para.font.bold = True
        num_para.font.color.rgb = WHITE
        num_para.alignment = PP_ALIGN.CENTER

        title_box = slide4.shapes.add_textbox(x_pos + Inches(0.2), action_y + Inches(0.8), action_width - Inches(0.4), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = action["title"]
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = DARK_NAVY
        title_para.alignment = PP_ALIGN.CENTER

        items_box = slide4.shapes.add_textbox(x_pos + Inches(0.2), action_y + Inches(1.5), action_width - Inches(0.4), Inches(1.6))
        items_frame = items_box.text_frame
        items_frame.word_wrap = True

        for j, item in enumerate(action["items"]):
            if j == 0:
                para = items_frame.paragraphs[0]
            else:
                para = items_frame.add_paragraph()
            para.text = f"• {item}"
            para.font.size = Pt(11)
            para.font.color.rgb = DARK_TEXT
            para.space_after = Pt(6)

    kpi_box4 = create_premium_shape(slide4, Inches(0.5), Inches(6.7), Inches(9), Inches(0.6),
                                     MSO_SHAPE.ROUNDED_RECTANGLE, DARK_NAVY, shadow=False)

    kpi4_frame = kpi_box4.text_frame
    kpi4_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    kpi4_para = kpi4_frame.paragraphs[0]
    kpi4_para.text = "KPI   재발 불량 30% 감소 / 품질 불량 10% 감소 목표 달성"
    kpi4_para.font.size = Pt(18)
    kpi4_para.font.bold = True
    kpi4_para.font.color.rgb = WHITE
    kpi4_para.alignment = PP_ALIGN.CENTER

    gold4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.7), Inches(9), Inches(0.05))
    gold4.fill.solid()
    gold4.fill.fore_color.rgb = GOLD
    gold4.line.fill.background()

    page_num4 = slide4.shapes.add_textbox(Inches(9.2), Inches(7.1), Inches(0.6), Inches(0.3))
    page_num4.text_frame.text = "04"
    page_num4.text_frame.paragraphs[0].font.size = Pt(11)
    page_num4.text_frame.paragraphs[0].font.color.rgb = BLUE_GRAY
    page_num4.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # 전략 3
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    slide5.background.fill.solid()
    slide5.background.fill.fore_color.rgb = OFF_WHITE

    accent5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
    accent5.fill.solid()
    accent5.fill.fore_color.rgb = RGBColor(230, 126, 34)  # Orange
    accent5.line.fill.background()

    header5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    h5_frame = header5.text_frame
    h5_page = h5_frame.paragraphs[0]
    h5_page.text = "04"
    h5_page.font.size = Pt(14)
    h5_page.font.color.rgb = BLUE_GRAY
    h5_page.space_after = Pt(5)

    h5_title = h5_frame.add_paragraph()
    h5_title.text = "전략 3   설비 CAPA 증가 및 공정 최적화"
    h5_title.font.size = Pt(32)
    h5_title.font.bold = True
    h5_title.font.color.rgb = DARK_NAVY

    underline5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(2.5), Inches(0.04))
    add_premium_gradient(underline5, RGBColor(230, 126, 34), RGBColor(245, 203, 167), 90)
    underline5.line.fill.background()

    overview5 = create_premium_shape(slide5, Inches(0.5), Inches(1.8), Inches(9), Inches(1),
                                     MSO_SHAPE.ROUNDED_RECTANGLE, WHITE, shadow=True)

    ov5_frame = overview5.text_frame
    ov5_frame.margin_left = Inches(0.3)
    ov5_frame.margin_top = Inches(0.2)

    ov5_para = ov5_frame.paragraphs[0]
    ov5_para.text = "전략 개요  "
    ov5_para.font.size = Pt(16)
    ov5_para.font.bold = True
    ov5_para.font.color.rgb = DARK_NAVY

    ov5_desc = ov5_frame.add_paragraph()
    ov5_desc.text = "C/T 단축 DEEP 분석 → 단 1초라도 줄이기 → 설비 효율 성능/PM 지속 개선 → 공정 최적화"
    ov5_desc.font.size = Pt(14)
    ov5_desc.font.color.rgb = DARK_TEXT
    ov5_desc.space_before = Pt(5)

    actions5 = [
        {"num": "01", "title": "C/T 단축 개선", "items": ["현 P/G 운영 DEEP 분석", "최단거리 프로세스 설계", "단 1초라도 줄이기 활동", "병목 공정 표적 개선"]},
        {"num": "02", "title": "기본 BASE 강화", "items": ["설비 효율 유지 및 향상", "성능 관련 모든 활동 반복", "PM(예방정비) 체계 강화", "지속적 모니터링"]},
        {"num": "03", "title": "공정 최적화", "items": ["RADIAL2 수삽설비 안정화", "SMD 공정 LAY OUT 개선", "최적 운영 방안 검토"]}
    ]

    for i, action in enumerate(actions5):
        x_pos = action_start + (i * (action_width + action_spacing))

        card = create_premium_shape(slide5, x_pos, action_y, action_width, action_height,
                                    MSO_SHAPE.ROUNDED_RECTANGLE, WHITE, shadow=True)

        top_accent = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, action_y, action_width, Inches(0.6))
        top_accent.fill.solid()
        top_accent.fill.fore_color.rgb = RGBColor(230, 126, 34)
        top_accent.line.fill.background()

        num_box = slide5.shapes.add_textbox(x_pos, action_y, action_width, Inches(0.6))
        num_frame = num_box.text_frame
        num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        num_para = num_frame.paragraphs[0]
        num_para.text = action["num"]
        num_para.font.size = Pt(28)
        num_para.font.bold = True
        num_para.font.color.rgb = WHITE
        num_para.alignment = PP_ALIGN.CENTER

        title_box = slide5.shapes.add_textbox(x_pos + Inches(0.2), action_y + Inches(0.8), action_width - Inches(0.4), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = action["title"]
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = DARK_NAVY
        title_para.alignment = PP_ALIGN.CENTER

        items_box = slide5.shapes.add_textbox(x_pos + Inches(0.2), action_y + Inches(1.5), action_width - Inches(0.4), Inches(1.6))
        items_frame = items_box.text_frame
        items_frame.word_wrap = True

        for j, item in enumerate(action["items"]):
            if j == 0:
                para = items_frame.paragraphs[0]
            else:
                para = items_frame.add_paragraph()
            para.text = f"• {item}"
            para.font.size = Pt(11)
            para.font.color.rgb = DARK_TEXT
            para.space_after = Pt(6)

    kpi_box5 = create_premium_shape(slide5, Inches(0.5), Inches(6.7), Inches(9), Inches(0.6),
                                     MSO_SHAPE.ROUNDED_RECTANGLE, DARK_NAVY, shadow=False)

    kpi5_frame = kpi_box5.text_frame
    kpi5_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    kpi5_para = kpi5_frame.paragraphs[0]
    kpi5_para.text = "KPI   설비 CAPA 증가 / OVERTIME 감소 / OH(간접비) 감소 달성"
    kpi5_para.font.size = Pt(18)
    kpi5_para.font.bold = True
    kpi5_para.font.color.rgb = WHITE
    kpi5_para.alignment = PP_ALIGN.CENTER

    gold5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.7), Inches(9), Inches(0.05))
    gold5.fill.solid()
    gold5.fill.fore_color.rgb = GOLD
    gold5.line.fill.background()

    page_num5 = slide5.shapes.add_textbox(Inches(9.2), Inches(7.1), Inches(0.6), Inches(0.3))
    page_num5.text_frame.text = "05"
    page_num5.text_frame.paragraphs[0].font.size = Pt(11)
    page_num5.text_frame.paragraphs[0].font.color.rgb = BLUE_GRAY
    page_num5.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # ========== 슬라이드 6: 로드맵 (프리미엄) ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    slide6.background.fill.solid()
    slide6.background.fill.fore_color.rgb = OFF_WHITE

    accent6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
    accent6.fill.solid()
    accent6.fill.fore_color.rgb = GOLD
    accent6.line.fill.background()

    header6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    h6_frame = header6.text_frame
    h6_page = h6_frame.paragraphs[0]
    h6_page.text = "05"
    h6_page.font.size = Pt(14)
    h6_page.font.color.rgb = BLUE_GRAY
    h6_page.space_after = Pt(5)

    h6_title = h6_frame.add_paragraph()
    h6_title.text = "2026 실행 로드맵 및 종합 KPI"
    h6_title.font.size = Pt(32)
    h6_title.font.bold = True
    h6_title.font.color.rgb = DARK_NAVY

    underline6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(2.5), Inches(0.04))
    add_premium_gradient(underline6, GOLD, LIGHT_GOLD, 90)
    underline6.line.fill.background()

    # 타임라인
    quarters = [
        {"q": "Q1", "color": ACCENT_BLUE, "tasks": ["TOOL 개발", "시스템 구축", "앱 제작"]},
        {"q": "Q2", "color": ACCENT_GREEN, "tasks": ["파일럿 운영", "피드백 반영", "개선"]},
        {"q": "Q3", "color": RGBColor(230, 126, 34), "tasks": ["전사 확대", "교육 실시", "모니터링"]},
        {"q": "Q4", "color": RGBColor(211, 47, 47), "tasks": ["목표 달성", "성과 점검", "2027 계획"]}
    ]

    q_width = Inches(2)
    q_spacing = Inches(0.25)
    q_start = Inches(0.5)
    q_y = Inches(2)

    for i, qt in enumerate(quarters):
        x_pos = q_start + (i * (q_width + q_spacing))

        card = create_premium_shape(slide6, x_pos, q_y, q_width, Inches(2.2),
                                    MSO_SHAPE.ROUNDED_RECTANGLE, WHITE, shadow=True)

        # 상단 컬러 헤더
        q_header = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, q_y, q_width, Inches(0.5))
        q_header.fill.solid()
        q_header.fill.fore_color.rgb = qt["color"]
        q_header.line.fill.background()

        # 분기
        q_box = slide6.shapes.add_textbox(x_pos, q_y, q_width, Inches(0.5))
        q_frame = q_box.text_frame
        q_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        q_para = q_frame.paragraphs[0]
        q_para.text = qt["q"]
        q_para.font.size = Pt(24)
        q_para.font.bold = True
        q_para.font.color.rgb = WHITE
        q_para.alignment = PP_ALIGN.CENTER

        # 태스크
        task_box = slide6.shapes.add_textbox(x_pos + Inches(0.2), q_y + Inches(0.7), q_width - Inches(0.4), Inches(1.3))
        task_frame = task_box.text_frame
        task_frame.word_wrap = True

        for j, task in enumerate(qt["tasks"]):
            if j == 0:
                para = task_frame.paragraphs[0]
            else:
                para = task_frame.add_paragraph()
            para.text = f"▪ {task}"
            para.font.size = Pt(12)
            para.font.color.rgb = DARK_TEXT
            para.space_after = Pt(6)

    # 하단 KPI 카드
    kpi_title = slide6.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.4))
    kpi_t_frame = kpi_title.text_frame
    kpi_t_para = kpi_t_frame.paragraphs[0]
    kpi_t_para.text = "종합 KPI Dashboard"
    kpi_t_para.font.size = Pt(22)
    kpi_t_para.font.bold = True
    kpi_t_para.font.color.rgb = DARK_NAVY

    final_kpis = [
        {"label": "가공비 절감", "value": "10%", "color": ACCENT_BLUE},
        {"label": "품질 개선", "value": "10%", "color": ACCENT_GREEN},
        {"label": "손실시간", "value": "5%", "color": RGBColor(230, 126, 34)}
    ]

    kpi_w = Inches(2.8)
    kpi_sp = Inches(0.3)
    kpi_st = Inches(0.5)
    kpi_y_pos = Inches(5.1)

    for i, kpi in enumerate(final_kpis):
        x_pos = kpi_st + (i * (kpi_w + kpi_sp))

        card = create_premium_shape(slide6, x_pos, kpi_y_pos, kpi_w, Inches(1.6),
                                    MSO_SHAPE.ROUNDED_RECTANGLE, kpi["color"], shadow=True)

        # 라벨
        label_box = slide6.shapes.add_textbox(x_pos, kpi_y_pos + Inches(0.2), kpi_w, Inches(0.3))
        label_frame = label_box.text_frame
        label_para = label_frame.paragraphs[0]
        label_para.text = kpi["label"]
        label_para.font.size = Pt(14)
        label_para.font.color.rgb = WHITE
        label_para.alignment = PP_ALIGN.CENTER

        # 값
        value_box = slide6.shapes.add_textbox(x_pos, kpi_y_pos + Inches(0.55), kpi_w, Inches(0.8))
        value_frame = value_box.text_frame
        value_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        value_para = value_frame.paragraphs[0]
        value_para.text = f"-{kpi['value']}"
        value_para.font.size = Pt(56)
        value_para.font.bold = True
        value_para.font.color.rgb = WHITE
        value_para.alignment = PP_ALIGN.CENTER

    page_num6 = slide6.shapes.add_textbox(Inches(9.2), Inches(7.1), Inches(0.6), Inches(0.3))
    page_num6.text_frame.text = "06"
    page_num6.text_frame.paragraphs[0].font.size = Pt(11)
    page_num6.text_frame.paragraphs[0].font.color.rgb = BLUE_GRAY
    page_num6.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # 저장
    output_file = '최종_Ultra_Premium_2026전략.pptx'
    prs.save(output_file)
    print(f"✅ Ultra Premium PPT 생성 완료: {output_file}")
    print(f"📄 총 6페이지")
    print(f"🎨 디자인 특징:")
    print(f"   - 다크 네이비 & 골드 컬러 팔레트")
    print(f"   - 프리미엄 그라데이션 효과")
    print(f"   - 그림자가 있는 3D 카드")
    print(f"   - 세련된 타이포그래피")
    print(f"   - 전문 비즈니스 프레젠테이션 스타일")
    return output_file

if __name__ == "__main__":
    create_ultra_premium_presentation()
