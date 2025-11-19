#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
2026년 제조1팀 경영전략 PPT - 다양한 디자인 통합 버전
- 설비관리 혁신방안 추가
- 과거 이력 연속성 강화
- 평가가동 효율 목표 (SMD 91%, RADIAL 85%, AXIAL 85%)
- 페이지당 다양한 시각화 적용
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import os

# 색상 팔레트
NAVY = RGBColor(25, 42, 86)
GOLD = RGBColor(212, 175, 55)
LIGHT_BLUE = RGBColor(52, 152, 219)
GREEN = RGBColor(46, 204, 113)
ORANGE = RGBColor(230, 126, 34)
RED = RGBColor(231, 76, 60)
PURPLE = RGBColor(155, 89, 182)
GRAY = RGBColor(127, 140, 141)
LIGHT_GRAY = RGBColor(236, 240, 241)
WHITE = RGBColor(255, 255, 255)

def create_title_slide(prs):
    """페이지 1: 프리미엄 커버"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경 그라데이션 효과 (네이비 박스)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2), Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "2026년 경영전략"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = GOLD
    title_para.alignment = PP_ALIGN.CENTER

    # 부제목
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.5), Inches(8), Inches(0.6)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "자동화 제조라인 스마트화를 통한 가공비 절감 및 품질 개선"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = WHITE
    subtitle_para.alignment = PP_ALIGN.CENTER

    # 하단 팀명
    team_box = slide.shapes.add_textbox(
        Inches(1), Inches(6.5), Inches(8), Inches(0.5)
    )
    team_frame = team_box.text_frame
    team_frame.text = "제조1팀"
    team_para = team_frame.paragraphs[0]
    team_para.font.size = Pt(24)
    team_para.font.bold = True
    team_para.font.color.rgb = WHITE
    team_para.alignment = PP_ALIGN.CENTER

def create_history_timeline(prs):
    """페이지 2: 과거 전략 회고 및 연속성 - 타임라인 스타일"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "전략 연속성: 과거 성과 기반 2026 방향"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # 타임라인 화살표
    arrow_y = Inches(1.5)
    arrow_width = Inches(8.5)

    # 배경 화살표
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.7), arrow_y, arrow_width, Inches(0.3)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = LIGHT_GRAY
    arrow.line.fill.background()

    # 시간대별 박스
    periods = [
        {"year": "2021-2022", "color": GRAY, "x": 0.7},
        {"year": "2023-2024", "color": LIGHT_BLUE, "x": 3.2},
        {"year": "2025", "color": GREEN, "x": 5.7},
        {"year": "2026", "color": GOLD, "x": 8.2}
    ]

    for period in periods:
        # 동그라미
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(period["x"]), arrow_y - Inches(0.15), Inches(0.6), Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = period["color"]
        circle.line.color.rgb = WHITE
        circle.line.width = Pt(3)

        # 연도 텍스트
        year_box = slide.shapes.add_textbox(
            Inches(period["x"] - 0.3), arrow_y + Inches(0.6), Inches(1.2), Inches(0.4)
        )
        year_frame = year_box.text_frame
        year_frame.text = period["year"]
        yp = year_frame.paragraphs[0]
        yp.font.size = Pt(14)
        yp.font.bold = True
        yp.font.color.rgb = period["color"]
        yp.alignment = PP_ALIGN.CENTER

    # 과거 핵심 과제 (2021-2024)
    past_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(2.8), Inches(4), Inches(2.5)
    )
    past_box.fill.solid()
    past_box.fill.fore_color.rgb = LIGHT_GRAY
    past_box.line.color.rgb = GRAY

    past_text = past_box.text_frame
    past_text.text = "과거 핵심 과제 (유지·강화)"
    past_text.paragraphs[0].font.size = Pt(18)
    past_text.paragraphs[0].font.bold = True
    past_text.paragraphs[0].font.color.rgb = NAVY
    past_text.paragraphs[0].alignment = PP_ALIGN.CENTER

    tasks = [
        "✓ 유실시간 개선 (연속 추진)",
        "✓ 설비 CAPA 관리 (강화)",
        "✓ 불량률 감소 활동",
        "✓ 공정 표준화"
    ]

    for task in tasks:
        p = past_text.add_paragraph()
        p.text = task
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.space_before = Pt(8)
        p.level = 0

    # 2026 신규 과제
    new_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.5), Inches(2.8), Inches(4), Inches(2.5)
    )
    new_box.fill.solid()
    new_box.fill.fore_color.rgb = RGBColor(255, 250, 230)
    new_box.line.color.rgb = GOLD
    new_box.line.width = Pt(3)

    new_text = new_box.text_frame
    new_text.text = "2026 신규·고도화"
    new_text.paragraphs[0].font.size = Pt(18)
    new_text.paragraphs[0].font.bold = True
    new_text.paragraphs[0].font.color.rgb = NAVY
    new_text.paragraphs[0].alignment = PP_ALIGN.CENTER

    new_tasks = [
        "★ 순간유실 실시간 가시화 (신규)",
        "★ 불량 재발 Zero 시스템 (고도화)",
        "★ 설비관리 혁신방안 (신규)",
        "★ MES 연동 자동화"
    ]

    for task in new_tasks:
        p = new_text.add_paragraph()
        p.text = task
        p.font.size = Pt(14)
        p.font.color.rgb = NAVY
        p.space_before = Pt(8)
        p.level = 0

    # 연결 화살표
    arrow_connect = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(4.6), Inches(4), Inches(0.8), Inches(0.4)
    )
    arrow_connect.fill.solid()
    arrow_connect.fill.fore_color.rgb = GOLD
    arrow_connect.line.fill.background()

def create_strategy_overview(prs):
    """페이지 3: 2026 전략 방향성 - 플로우차트"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "2026 전략 방향: 4대 핵심 과제"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # 중앙 목표
    center_box = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(3.5), Inches(1.5), Inches(3), Inches(1.5)
    )
    center_box.fill.solid()
    center_box.fill.fore_color.rgb = GOLD
    center_box.line.fill.background()

    center_text = center_box.text_frame
    center_text.text = "스마트 제조혁신"
    center_text.paragraphs[0].font.size = Pt(24)
    center_text.paragraphs[0].font.bold = True
    center_text.paragraphs[0].font.color.rgb = WHITE
    center_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    center_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    p2 = center_text.add_paragraph()
    p2.text = "가공비 10% 절감\n품질 10% 개선"
    p2.font.size = Pt(14)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

    # 4개 전략 박스 (사방에 배치)
    strategies = [
        {"name": "순간유실\nZero", "color": LIGHT_BLUE, "x": 1, "y": 3.8},
        {"name": "불량재발\nZero", "color": GREEN, "x": 7.5, "y": 3.8},
        {"name": "설비CAPA\n증대", "color": ORANGE, "x": 1, "y": 5.8},
        {"name": "설비관리\n혁신", "color": PURPLE, "x": 7.5, "y": 5.8}
    ]

    for strat in strategies:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(strat["x"]), Inches(strat["y"]), Inches(1.8), Inches(1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = strat["color"]
        box.line.fill.background()

        text = box.text_frame
        text.text = strat["name"]
        text.paragraphs[0].font.size = Pt(16)
        text.paragraphs[0].font.bold = True
        text.paragraphs[0].font.color.rgb = WHITE
        text.paragraphs[0].alignment = PP_ALIGN.CENTER
        text.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 연결선
        line = slide.shapes.add_connector(
            1, Inches(5), Inches(2.5), Inches(strat["x"] + 0.9), Inches(strat["y"])
        )
        line.line.color.rgb = GRAY
        line.line.width = Pt(2)

def create_strategy1_progress(prs):
    """페이지 4: 전략1 - 순간유실 Zero (프로그레스바 + 아이콘)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "전략 1: 순간유실 Zero 프로젝트"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE

    # 목표 박스
    goal_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.2), Inches(4), Inches(1)
    )
    goal_box.fill.solid()
    goal_box.fill.fore_color.rgb = RGBColor(230, 240, 250)
    goal_box.line.color.rgb = LIGHT_BLUE

    goal_text = goal_box.text_frame
    goal_text.text = "목표"
    goal_text.paragraphs[0].font.size = Pt(16)
    goal_text.paragraphs[0].font.bold = True
    goal_text.paragraphs[0].font.color.rgb = LIGHT_BLUE

    p = goal_text.add_paragraph()
    p.text = "순간유실 5% → 2% 감소"
    p.font.size = Pt(14)
    p.font.color.rgb = NAVY

    # 세부 과제 (프로그레스바)
    tasks = [
        {"name": "실시간 모니터링 시스템 구축", "progress": 40},
        {"name": "순간정지 알람 시스템", "progress": 60},
        {"name": "유실 원인 자동 분석", "progress": 30},
        {"name": "MES 연동 데이터 수집", "progress": 50}
    ]

    start_y = 2.5
    for i, task in enumerate(tasks):
        y = start_y + i * 1

        # 과제명
        name_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(y), Inches(4), Inches(0.3)
        )
        name_frame = name_box.text_frame
        name_frame.text = f"• {task['name']}"
        name_frame.paragraphs[0].font.size = Pt(13)
        name_frame.paragraphs[0].font.color.rgb = NAVY

        # 프로그레스바 배경
        bg_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(y + 0.35), Inches(4), Inches(0.25)
        )
        bg_bar.fill.solid()
        bg_bar.fill.fore_color.rgb = LIGHT_GRAY
        bg_bar.line.fill.background()

        # 프로그레스바
        progress_width = 4 * (task['progress'] / 100)
        progress_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(y + 0.35), Inches(progress_width), Inches(0.25)
        )
        progress_bar.fill.solid()
        progress_bar.fill.fore_color.rgb = LIGHT_BLUE
        progress_bar.line.fill.background()

        # 퍼센트 원
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.5 + progress_width - 0.2), Inches(y + 0.25),
            Inches(0.45), Inches(0.45)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = WHITE
        circle.line.color.rgb = LIGHT_BLUE
        circle.line.width = Pt(3)

        pct_text = circle.text_frame
        pct_text.text = f"{task['progress']}%"
        pct_text.paragraphs[0].font.size = Pt(11)
        pct_text.paragraphs[0].font.bold = True
        pct_text.paragraphs[0].font.color.rgb = LIGHT_BLUE
        pct_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        pct_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 기대효과
    effect_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.3)
    )
    effect_box.fill.solid()
    effect_box.fill.fore_color.rgb = RGBColor(245, 250, 255)
    effect_box.line.color.rgb = LIGHT_BLUE
    effect_box.line.width = Pt(2)

    effect_text = effect_box.text_frame
    effect_text.text = "기대효과"
    effect_text.paragraphs[0].font.size = Pt(18)
    effect_text.paragraphs[0].font.bold = True
    effect_text.paragraphs[0].font.color.rgb = LIGHT_BLUE
    effect_text.paragraphs[0].alignment = PP_ALIGN.CENTER

    effects = [
        "📊 유실시간 60% 감소",
        "⚡ 대응속도 1000배 향상",
        "💰 가공비 3% 절감",
        "🎯 생산성 5% 증가"
    ]

    for effect in effects:
        p = effect_text.add_paragraph()
        p.text = effect
        p.font.size = Pt(16)
        p.font.color.rgb = NAVY
        p.space_before = Pt(15)
        p.alignment = PP_ALIGN.LEFT

def create_strategy2_table(prs):
    """페이지 5: 전략2 - 불량 재발 Zero (테이블 + 체크리스트)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "전략 2: 불량 재발 Zero 시스템"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = GREEN

    # 목표
    goal_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.2), Inches(4), Inches(0.8)
    )
    goal_box.fill.solid()
    goal_box.fill.fore_color.rgb = RGBColor(230, 250, 230)
    goal_box.line.color.rgb = GREEN

    goal_text = goal_box.text_frame
    goal_text.text = "목표: 불량률 10% → 5% 개선"
    goal_text.paragraphs[0].font.size = Pt(16)
    goal_text.paragraphs[0].font.bold = True
    goal_text.paragraphs[0].font.color.rgb = GREEN
    goal_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    goal_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 3단계 프로세스 (좌측)
    processes = [
        {"step": "1단계", "name": "불량 발생 즉시 감지", "icon": "🔍"},
        {"step": "2단계", "name": "원인 자동 분석", "icon": "🧠"},
        {"step": "3단계", "name": "재발 방지 조치", "icon": "🛡️"}
    ]

    start_y = 2.3
    for i, proc in enumerate(processes):
        y = start_y + i * 1.3

        # 단계 박스
        step_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(y), Inches(4.2), Inches(1)
        )
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = WHITE
        step_box.line.color.rgb = GREEN
        step_box.line.width = Pt(2)

        step_text = step_box.text_frame
        step_text.text = f"{proc['icon']}  {proc['step']}: {proc['name']}"
        step_text.paragraphs[0].font.size = Pt(16)
        step_text.paragraphs[0].font.bold = True
        step_text.paragraphs[0].font.color.rgb = NAVY
        step_text.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 화살표
        if i < 2:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(2.3), Inches(y + 1.05), Inches(0.4), Inches(0.2)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GREEN
            arrow.line.fill.background()

    # 체크리스트 (우측)
    checklist_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.3)
    )
    checklist_box.fill.solid()
    checklist_box.fill.fore_color.rgb = RGBColor(245, 255, 245)
    checklist_box.line.color.rgb = GREEN
    checklist_box.line.width = Pt(2)

    checklist_text = checklist_box.text_frame
    checklist_text.text = "세부 실행 과제"
    checklist_text.paragraphs[0].font.size = Pt(18)
    checklist_text.paragraphs[0].font.bold = True
    checklist_text.paragraphs[0].font.color.rgb = GREEN
    checklist_text.paragraphs[0].alignment = PP_ALIGN.CENTER

    tasks = [
        "☑ 비전검사 시스템 도입",
        "☑ 불량 DB 구축 및 AI 분석",
        "☑ 작업자 실시간 알림",
        "☑ SOP 자동 업데이트",
        "☑ 예방점검 자동화",
        "☑ 품질 트렌드 분석"
    ]

    for task in tasks:
        p = checklist_text.add_paragraph()
        p.text = task
        p.font.size = Pt(15)
        p.font.color.rgb = NAVY
        p.space_before = Pt(12)
        p.level = 0

    # 하단 기대효과
    effect_text = checklist_text.add_paragraph()
    effect_text.text = "\n기대효과"
    effect_text.font.size = Pt(16)
    effect_text.font.bold = True
    effect_text.font.color.rgb = GREEN
    effect_text.space_before = Pt(20)
    effect_text.alignment = PP_ALIGN.CENTER

    p_effect = checklist_text.add_paragraph()
    p_effect.text = "불량비용 40% 절감\n고객 클레임 70% 감소"
    p_effect.font.size = Pt(14)
    p_effect.font.color.rgb = NAVY
    p_effect.alignment = PP_ALIGN.CENTER

def create_strategy3_chart(prs):
    """페이지 6: 전략3 - 설비 CAPA 증대 (막대 그래프)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "전략 3: 설비 CAPA 증대"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ORANGE

    # 목표
    goal_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.2), Inches(4), Inches(0.8)
    )
    goal_box.fill.solid()
    goal_box.fill.fore_color.rgb = RGBColor(255, 245, 230)
    goal_box.line.color.rgb = ORANGE

    goal_text = goal_box.text_frame
    goal_text.text = "목표: 생산능력 15% 향상"
    goal_text.paragraphs[0].font.size = Pt(16)
    goal_text.paragraphs[0].font.bold = True
    goal_text.paragraphs[0].font.color.rgb = ORANGE
    goal_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    goal_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 비교 막대 그래프
    chart_data = [
        {"item": "Tact Time", "before": 12, "after": 10, "unit": "초"},
        {"item": "설비 가동률", "before": 75, "after": 90, "unit": "%"},
        {"item": "일일 생산량", "before": 5000, "after": 5750, "unit": "개"}
    ]

    start_y = 2.5
    max_width = 3.5

    for i, data in enumerate(chart_data):
        y = start_y + i * 1.3

        # 항목명
        item_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(y), Inches(1.8), Inches(0.4)
        )
        item_frame = item_box.text_frame
        item_frame.text = data['item']
        item_frame.paragraphs[0].font.size = Pt(14)
        item_frame.paragraphs[0].font.bold = True
        item_frame.paragraphs[0].font.color.rgb = NAVY
        item_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Before 막대
        max_val = max(data['before'], data['after'])
        before_width = max_width * (data['before'] / max_val)

        before_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.5), Inches(y), Inches(before_width), Inches(0.35)
        )
        before_bar.fill.solid()
        before_bar.fill.fore_color.rgb = LIGHT_GRAY
        before_bar.line.fill.background()

        before_text = before_bar.text_frame
        before_text.text = f"{data['before']}{data['unit']}"
        before_text.paragraphs[0].font.size = Pt(11)
        before_text.paragraphs[0].font.color.rgb = GRAY
        before_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        before_text.vertical_anchor = MSO_ANCHOR.MIDDLE

        # After 막대
        after_width = max_width * (data['after'] / max_val)

        after_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.5), Inches(y + 0.45), Inches(after_width), Inches(0.35)
        )
        after_bar.fill.solid()
        after_bar.fill.fore_color.rgb = ORANGE
        after_bar.line.fill.background()

        after_text = after_bar.text_frame
        after_text.text = f"{data['after']}{data['unit']}"
        after_text.paragraphs[0].font.size = Pt(11)
        after_text.paragraphs[0].font.bold = True
        after_text.paragraphs[0].font.color.rgb = WHITE
        after_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        after_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 범례
    legend_before = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.5), Inches(2.2), Inches(0.3), Inches(0.3)
    )
    legend_before.fill.solid()
    legend_before.fill.fore_color.rgb = LIGHT_GRAY
    legend_before.line.fill.background()

    legend_before_text = slide.shapes.add_textbox(
        Inches(6.9), Inches(2.15), Inches(1), Inches(0.4)
    )
    legend_before_text.text_frame.text = "현재"
    legend_before_text.text_frame.paragraphs[0].font.size = Pt(12)
    legend_before_text.text_frame.paragraphs[0].font.color.rgb = GRAY

    legend_after = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7.8), Inches(2.2), Inches(0.3), Inches(0.3)
    )
    legend_after.fill.solid()
    legend_after.fill.fore_color.rgb = ORANGE
    legend_after.line.fill.background()

    legend_after_text = slide.shapes.add_textbox(
        Inches(8.2), Inches(2.15), Inches(1), Inches(0.4)
    )
    legend_after_text.text_frame.text = "목표"
    legend_after_text.text_frame.paragraphs[0].font.size = Pt(12)
    legend_after_text.text_frame.paragraphs[0].font.color.rgb = ORANGE

    # 실행 계획 (우측 하단)
    plan_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(3.5), Inches(4.3), Inches(3)
    )
    plan_box.fill.solid()
    plan_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
    plan_box.line.color.rgb = ORANGE
    plan_box.line.width = Pt(2)

    plan_text = plan_box.text_frame
    plan_text.text = "실행 계획"
    plan_text.paragraphs[0].font.size = Pt(16)
    plan_text.paragraphs[0].font.bold = True
    plan_text.paragraphs[0].font.color.rgb = ORANGE
    plan_text.paragraphs[0].alignment = PP_ALIGN.CENTER

    plans = [
        "1. 병목공정 분석 및 개선",
        "2. 고속화 설비 개조",
        "3. 자동화 라인 증설",
        "4. 작업 동선 최적화",
        "5. 다기능 작업자 양성"
    ]

    for plan in plans:
        p = plan_text.add_paragraph()
        p.text = plan
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY
        p.space_before = Pt(8)

def create_strategy4_facility(prs):
    """페이지 7: 전략4 - 설비관리 혁신방안 (4분할 매트릭스)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "전략 4: 설비관리 혁신방안 (신규)"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    # 목표
    goal_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.2), Inches(9), Inches(0.6)
    )
    goal_box.fill.solid()
    goal_box.fill.fore_color.rgb = RGBColor(245, 235, 255)
    goal_box.line.color.rgb = PURPLE

    goal_text = goal_box.text_frame
    goal_text.text = "목표: 예방보전 체계 고도화로 설비 고장 50% 감소 및 설비 수명 20% 연장"
    goal_text.paragraphs[0].font.size = Pt(15)
    goal_text.paragraphs[0].font.bold = True
    goal_text.paragraphs[0].font.color.rgb = PURPLE
    goal_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    goal_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 4분할 매트릭스
    boxes = [
        {
            "title": "예방보전 고도화",
            "icon": "🔧",
            "items": ["• 주기 기반 → 상태 기반", "• IoT 센서 실시간 모니터링", "• 이상징후 조기 감지"],
            "x": 0.5, "y": 2.2
        },
        {
            "title": "설비 이력 관리",
            "icon": "📋",
            "items": ["• 설비별 정비 이력 DB화", "• 고장 패턴 분석", "• 부품 교체 주기 최적화"],
            "x": 5.2, "y": 2.2
        },
        {
            "title": "부품 수명 예측",
            "icon": "🎯",
            "items": ["• AI 기반 수명 예측", "• 적기 부품 교체", "• 재고 최적화"],
            "x": 0.5, "y": 4.5
        },
        {
            "title": "긴급 정비 체계",
            "icon": "⚡",
            "items": ["• 24시간 대응 체계", "• 비상부품 확보", "• 협력업체 네트워크"],
            "x": 5.2, "y": 4.5
        }
    ]

    for box_data in boxes:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(box_data["x"]), Inches(box_data["y"]),
            Inches(4.3), Inches(2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = PURPLE
        box.line.width = Pt(2)

        text_frame = box.text_frame
        text_frame.text = f"{box_data['icon']} {box_data['title']}"
        text_frame.paragraphs[0].font.size = Pt(16)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = PURPLE
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        for item in box_data['items']:
            p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(12)
            p.font.color.rgb = NAVY
            p.space_before = Pt(6)
            p.level = 0

def create_efficiency_targets(prs):
    """페이지 8: 평가가동 효율 목표 (비교 차트)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "2026 평가가동 효율 목표"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # 3개 라인 비교
    lines = [
        {"name": "SMD", "target": 91, "color": LIGHT_BLUE, "current": 85},
        {"name": "RADIAL", "target": 85, "color": GREEN, "current": 78},
        {"name": "AXIAL", "target": 85, "color": ORANGE, "current": 80}
    ]

    start_x = 1
    box_width = 2.5
    spacing = 0.3

    for i, line in enumerate(lines):
        x = start_x + i * (box_width + spacing)

        # 메인 박스
        main_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(1.5), Inches(box_width), Inches(3.5)
        )
        main_box.fill.solid()
        main_box.fill.fore_color.rgb = WHITE
        main_box.line.color.rgb = line['color']
        main_box.line.width = Pt(3)

        # 라인명
        name_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(1.7), Inches(box_width - 0.4), Inches(0.5)
        )
        name_frame = name_box.text_frame
        name_frame.text = line['name']
        name_frame.paragraphs[0].font.size = Pt(24)
        name_frame.paragraphs[0].font.bold = True
        name_frame.paragraphs[0].font.color.rgb = line['color']
        name_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 목표값 (대형)
        target_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(2.4), Inches(box_width - 0.4), Inches(1)
        )
        target_frame = target_box.text_frame
        target_frame.text = f"{line['target']}%"
        target_frame.paragraphs[0].font.size = Pt(48)
        target_frame.paragraphs[0].font.bold = True
        target_frame.paragraphs[0].font.color.rgb = line['color']
        target_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 목표 라벨
        label_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(3.4), Inches(box_width - 0.4), Inches(0.3)
        )
        label_frame = label_box.text_frame
        label_frame.text = "2026 목표"
        label_frame.paragraphs[0].font.size = Pt(12)
        label_frame.paragraphs[0].font.color.rgb = GRAY
        label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 현재값
        current_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x + 0.3), Inches(3.9), Inches(box_width - 0.6), Inches(0.5)
        )
        current_box.fill.solid()
        current_box.fill.fore_color.rgb = LIGHT_GRAY
        current_box.line.fill.background()

        current_text = current_box.text_frame
        current_text.text = f"현재: {line['current']}%"
        current_text.paragraphs[0].font.size = Pt(14)
        current_text.paragraphs[0].font.color.rgb = GRAY
        current_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        current_text.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 증가 화살표
        improvement = line['target'] - line['current']
        arrow_box = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(4.5), Inches(box_width - 0.6), Inches(0.4)
        )
        arrow_frame = arrow_box.text_frame
        arrow_frame.text = f"↑ {improvement}%p 향상"
        arrow_frame.paragraphs[0].font.size = Pt(13)
        arrow_frame.paragraphs[0].font.bold = True
        arrow_frame.paragraphs[0].font.color.rgb = line['color']
        arrow_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 하단 전략 요약
    strategy_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(5.5), Inches(9), Inches(1)
    )
    strategy_box.fill.solid()
    strategy_box.fill.fore_color.rgb = RGBColor(250, 250, 250)
    strategy_box.line.color.rgb = NAVY
    strategy_box.line.width = Pt(2)

    strategy_text = strategy_box.text_frame
    strategy_text.text = "핵심 전략"
    strategy_text.paragraphs[0].font.size = Pt(16)
    strategy_text.paragraphs[0].font.bold = True
    strategy_text.paragraphs[0].font.color.rgb = NAVY
    strategy_text.paragraphs[0].alignment = PP_ALIGN.CENTER

    p2 = strategy_text.add_paragraph()
    p2.text = "순간유실 Zero + 불량재발 Zero + 설비CAPA 증대 + 설비관리 혁신 = 평가가동 효율 극대화"
    p2.font.size = Pt(14)
    p2.font.color.rgb = NAVY
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)

def create_roadmap(prs):
    """페이지 9: Q1-Q4 로드맵 (간트 차트 스타일)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "2026 실행 로드맵"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # 분기 헤더
    quarters = ["Q1", "Q2", "Q3", "Q4"]
    header_start_x = 2.5
    quarter_width = 1.7

    for i, q in enumerate(quarters):
        x = header_start_x + i * quarter_width
        header_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(1.2), Inches(quarter_width - 0.1), Inches(0.5)
        )
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = NAVY
        header_box.line.fill.background()

        header_text = header_box.text_frame
        header_text.text = q
        header_text.paragraphs[0].font.size = Pt(18)
        header_text.paragraphs[0].font.bold = True
        header_text.paragraphs[0].font.color.rgb = WHITE
        header_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        header_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 과제별 간트 바
    tasks = [
        {
            "name": "순간유실 Zero",
            "color": LIGHT_BLUE,
            "quarters": [1, 1, 1, 1]  # 전 분기
        },
        {
            "name": "불량 재발 Zero",
            "color": GREEN,
            "quarters": [1, 1, 1, 0]  # Q1-Q3
        },
        {
            "name": "설비 CAPA 증대",
            "color": ORANGE,
            "quarters": [0, 1, 1, 1]  # Q2-Q4
        },
        {
            "name": "설비관리 혁신",
            "color": PURPLE,
            "quarters": [1, 1, 0, 0]  # Q1-Q2
        }
    ]

    start_y = 2
    row_height = 0.9

    for i, task in enumerate(tasks):
        y = start_y + i * row_height

        # 과제명
        name_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(y), Inches(1.8), Inches(0.6)
        )
        name_frame = name_box.text_frame
        name_frame.text = task['name']
        name_frame.paragraphs[0].font.size = Pt(13)
        name_frame.paragraphs[0].font.bold = True
        name_frame.paragraphs[0].font.color.rgb = task['color']
        name_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 간트 바
        for q_idx, active in enumerate(task['quarters']):
            x = header_start_x + q_idx * quarter_width

            if active:
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(x), Inches(y + 0.1), Inches(quarter_width - 0.1), Inches(0.4)
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = task['color']
                bar.line.fill.background()
            else:
                # 비활성 영역 (점선 효과)
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(x), Inches(y + 0.1), Inches(quarter_width - 0.1), Inches(0.4)
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = LIGHT_GRAY
                bar.line.fill.background()

    # 마일스톤
    milestones = [
        {"text": "중간 점검", "q": 1, "y": 5.5},
        {"text": "성과 평가", "q": 3, "y": 5.5}
    ]

    for ms in milestones:
        x = header_start_x + ms['q'] * quarter_width

        # 다이아몬드
        diamond = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND,
            Inches(x + 0.65), Inches(ms['y']), Inches(0.4), Inches(0.4)
        )
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = RED
        diamond.line.fill.background()

        # 텍스트
        ms_text = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(ms['y'] + 0.5), Inches(1.3), Inches(0.3)
        )
        ms_frame = ms_text.text_frame
        ms_frame.text = ms['text']
        ms_frame.paragraphs[0].font.size = Pt(11)
        ms_frame.paragraphs[0].font.color.rgb = RED
        ms_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def create_kpi_dashboard(prs):
    """페이지 10: 종합 KPI 대시보드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "2026 핵심 KPI 대시보드"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # 3개 원형 게이지
    kpis = [
        {"name": "가공비 절감", "target": 10, "color": LIGHT_BLUE, "x": 1},
        {"name": "품질 개선", "target": 10, "color": GREEN, "x": 4},
        {"name": "유실시간 감소", "target": 5, "color": ORANGE, "x": 7}
    ]

    for kpi in kpis:
        # 외부 원 (배경)
        outer_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(kpi['x']), Inches(1.5), Inches(2), Inches(2)
        )
        outer_circle.fill.solid()
        outer_circle.fill.fore_color.rgb = LIGHT_GRAY
        outer_circle.line.fill.background()

        # 내부 원 (진행률)
        inner_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(kpi['x'] + 0.15), Inches(1.65), Inches(1.7), Inches(1.7)
        )
        inner_circle.fill.solid()
        inner_circle.fill.fore_color.rgb = kpi['color']
        inner_circle.line.fill.background()

        # 중앙 흰색 원
        center_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(kpi['x'] + 0.5), Inches(2), Inches(1), Inches(1)
        )
        center_circle.fill.solid()
        center_circle.fill.fore_color.rgb = WHITE
        center_circle.line.fill.background()

        # 퍼센트 텍스트
        pct_box = slide.shapes.add_textbox(
            Inches(kpi['x'] + 0.5), Inches(2.2), Inches(1), Inches(0.6)
        )
        pct_frame = pct_box.text_frame
        pct_frame.text = f"{kpi['target']}%"
        pct_frame.paragraphs[0].font.size = Pt(32)
        pct_frame.paragraphs[0].font.bold = True
        pct_frame.paragraphs[0].font.color.rgb = kpi['color']
        pct_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        pct_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # KPI 이름
        name_box = slide.shapes.add_textbox(
            Inches(kpi['x']), Inches(3.7), Inches(2), Inches(0.4)
        )
        name_frame = name_box.text_frame
        name_frame.text = kpi['name']
        name_frame.paragraphs[0].font.size = Pt(16)
        name_frame.paragraphs[0].font.bold = True
        name_frame.paragraphs[0].font.color.rgb = NAVY
        name_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 하단 세부 지표
    metrics_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(4.5), Inches(9), Inches(2)
    )
    metrics_box.fill.solid()
    metrics_box.fill.fore_color.rgb = RGBColor(245, 248, 250)
    metrics_box.line.color.rgb = NAVY
    metrics_box.line.width = Pt(2)

    metrics_text = metrics_box.text_frame
    metrics_text.text = "세부 성과 지표"
    metrics_text.paragraphs[0].font.size = Pt(18)
    metrics_text.paragraphs[0].font.bold = True
    metrics_text.paragraphs[0].font.color.rgb = NAVY
    metrics_text.paragraphs[0].alignment = PP_ALIGN.CENTER

    details = [
        "• 순간유실시간: 5% → 2% (60% 감소)",
        "• 불량률: 10% → 5% (50% 개선)",
        "• 설비 가동률: 75% → 90% (15%p 향상)",
        "• 평가가동 효율: SMD 91%, RADIAL 85%, AXIAL 85%",
        "• 설비 고장: 50% 감소, 수명 20% 연장"
    ]

    for detail in details:
        p = metrics_text.add_paragraph()
        p.text = detail
        p.font.size = Pt(14)
        p.font.color.rgb = NAVY
        p.space_before = Pt(8)
        p.level = 0

def main():
    """메인 실행 함수"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("PPT 생성 시작...")

    create_title_slide(prs)
    print("✓ 페이지 1: 커버")

    create_history_timeline(prs)
    print("✓ 페이지 2: 과거 전략 회고 (타임라인)")

    create_strategy_overview(prs)
    print("✓ 페이지 3: 전략 방향성 (플로우차트)")

    create_strategy1_progress(prs)
    print("✓ 페이지 4: 전략1 - 순간유실 Zero (프로그레스바)")

    create_strategy2_table(prs)
    print("✓ 페이지 5: 전략2 - 불량재발 Zero (테이블)")

    create_strategy3_chart(prs)
    print("✓ 페이지 6: 전략3 - 설비CAPA (막대그래프)")

    create_strategy4_facility(prs)
    print("✓ 페이지 7: 전략4 - 설비관리 혁신 (매트릭스)")

    create_efficiency_targets(prs)
    print("✓ 페이지 8: 평가가동 효율 목표")

    create_roadmap(prs)
    print("✓ 페이지 9: Q1-Q4 로드맵 (간트차트)")

    create_kpi_dashboard(prs)
    print("✓ 페이지 10: KPI 대시보드")

    output_file = '2026전략_최종완성본_v2.pptx'
    prs.save(output_file)

    print(f"\n{'='*60}")
    print(f"✅ PPT 생성 완료: {output_file}")
    print(f"📄 총 10페이지")
    print(f"🎨 특징:")
    print(f"   - 설비관리 혁신방안 추가 (4대 전략)")
    print(f"   - 과거 이력 연속성 강화")
    print(f"   - 평가가동 효율 목표: SMD 91%, RADIAL 85%, AXIAL 85%")
    print(f"   - 다양한 시각화: 타임라인, 플로우차트, 프로그레스바,")
    print(f"     테이블, 막대그래프, 매트릭스, 간트차트, 원형게이지")
    print(f"{'='*60}")

if __name__ == "__main__":
    main()
