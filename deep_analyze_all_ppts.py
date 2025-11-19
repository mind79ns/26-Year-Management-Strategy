#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ì „ëµì°¸ê³ ìë£Œ í´ë”ì˜ ëª¨ë“  PPT ì‹¬ì¸µ ë¶„ì„
- ëª¨ë“  ìŠ¬ë¼ì´ë“œì˜ ëª¨ë“  í…ìŠ¤íŠ¸ ì¶”ì¶œ
- í˜ì´ì§€ë³„ ì£¼ì œ ë¶„ë¥˜
- í•µì‹¬ KPI ë° ëª©í‘œ ì¶”ì¶œ
- ê³¼ì œ ë° ì‹¤í–‰ ê³„íš ì¶”ì¶œ
"""

from pptx import Presentation
import os
import json
import re
from collections import defaultdict, Counter

def extract_detailed_content(ppt_path):
    """PPT íŒŒì¼ì—ì„œ ìƒì„¸ ë‚´ìš© ì¶”ì¶œ"""
    try:
        prs = Presentation(ppt_path)
        file_info = {
            "file_name": os.path.basename(ppt_path),
            "file_path": ppt_path,
            "total_slides": len(prs.slides),
            "slides": []
        }

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_info = {
                "slide_number": slide_num,
                "texts": [],
                "all_text": "",
                "bullets": [],
                "numbers": [],
                "keywords": []
            }

            # ëª¨ë“  í…ìŠ¤íŠ¸ ìˆ˜ì§‘
            all_texts = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    all_texts.append(text)
                    slide_info["texts"].append(text)

                    # í…ìŠ¤íŠ¸ í”„ë ˆì„ì´ ìˆëŠ” ê²½ìš° ë‹¨ë½ë³„ ë¶„ì„
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = paragraph.text.strip()
                            if para_text:
                                # ë¶ˆë¦¿ í¬ì¸íŠ¸ ê°ì§€
                                if para_text.startswith(('â€¢', '-', 'Â·', 'â€»', 'â†’', 'â–¶', 'â–¡', 'â– ', 'â—‹', 'â—')):
                                    slide_info["bullets"].append(para_text)

            # ì „ì²´ í…ìŠ¤íŠ¸ í†µí•©
            slide_info["all_text"] = " ".join(all_texts)

            # ìˆ«ì íŒ¨í„´ ì¶”ì¶œ (ëª©í‘œ, KPI ë“±)
            numbers = re.findall(r'\d+(?:\.\d+)?(?:%|ê°œ|ëª…|ì‹œê°„|ë¶„|ì´ˆ|ê±´|íšŒ|ëŒ€|ì–µ|ë§Œ|ì²œ)', slide_info["all_text"])
            slide_info["numbers"] = numbers

            # í‚¤ì›Œë“œ ì¶”ì¶œ (í•œê¸€ 2ê¸€ì ì´ìƒ)
            korean_words = re.findall(r'[ê°€-í£]{2,}', slide_info["all_text"])
            slide_info["keywords"] = korean_words

            file_info["slides"].append(slide_info)

        return file_info

    except Exception as e:
        print(f"Error processing {ppt_path}: {e}")
        return None

def analyze_themes(file_info):
    """ìŠ¬ë¼ì´ë“œë³„ ì£¼ì œ ë¶„ë¥˜"""
    theme_keywords = {
        "ëª©í‘œ": ["ëª©í‘œ", "Target", "ì „ëµ", "ë°©í–¥", "ë¹„ì „"],
        "ìœ ì‹¤ì‹œê°„": ["ìœ ì‹¤", "Loss", "ì •ì§€", "ë¹„ê°€ë™", "ê°€ë™ë¥ "],
        "ë¶ˆëŸ‰": ["ë¶ˆëŸ‰", "í’ˆì§ˆ", "Defect", "PPM", "ì¬ë°œ"],
        "ì„¤ë¹„": ["ì„¤ë¹„", "Equipment", "CAPA", "ëŠ¥ë ¥", "ìƒì‚°ë ¥"],
        "ì¸ë ¥": ["ì¸ë ¥", "ì¸ì›", "ì‘ì—…ì", "êµìœ¡", "í›ˆë ¨"],
        "ì›ê°€": ["ì›ê°€", "ë¹„ìš©", "Cost", "ì ˆê°", "ê°€ê³µë¹„"],
        "ìƒì‚°ì„±": ["ìƒì‚°ì„±", "íš¨ìœ¨", "Productivity", "ê°œì„ "],
        "ìë™í™”": ["ìë™í™”", "Automation", "ìŠ¤ë§ˆíŠ¸", "MES", "ì‹œìŠ¤í…œ"],
        "ì‹¤í–‰ê³„íš": ["ì‹¤í–‰", "ê³„íš", "ì¼ì •", "ë¡œë“œë§µ", "ì¶”ì§„"],
        "ì„±ê³¼": ["ì„±ê³¼", "ê²°ê³¼", "ë‹¬ì„±", "ì‹¤ì "]
    }

    for slide in file_info["slides"]:
        slide["themes"] = []
        text = slide["all_text"].lower()

        for theme, keywords in theme_keywords.items():
            for keyword in keywords:
                if keyword.lower() in text:
                    if theme not in slide["themes"]:
                        slide["themes"].append(theme)
                    break

def extract_tasks_and_kpis(file_info):
    """ê³¼ì œ ë° KPI ì¶”ì¶œ"""
    tasks = []
    kpis = []

    for slide in file_info["slides"]:
        # ê³¼ì œ íŒ¨í„´ ì¶”ì¶œ
        for bullet in slide["bullets"]:
            # ê³¼ì œë¡œ ë³´ì´ëŠ” íŒ¨í„´
            if any(keyword in bullet for keyword in ["ê°œì„ ", "êµ¬ì¶•", "ë„ì…", "ì¶”ì§„", "ì‹¤í–‰", "ìˆ˜ë¦½", "ê°•í™”", "í™•ëŒ€"]):
                tasks.append({
                    "slide": slide["slide_number"],
                    "task": bullet,
                    "themes": slide["themes"]
                })

        # KPI íŒ¨í„´ ì¶”ì¶œ (ìˆ«ì í¬í•¨)
        for text in slide["texts"]:
            if any(char in text for char in ['%', 'â†’', 'â–¶']) and any(char.isdigit() for char in text):
                kpis.append({
                    "slide": slide["slide_number"],
                    "kpi": text,
                    "numbers": slide["numbers"]
                })

    return tasks, kpis

def main():
    """ë©”ì¸ ë¶„ì„ ì‹¤í–‰"""
    ppt_folder = "ì „ëµìë£Œì°¸ê³ "

    # PPT íŒŒì¼ ëª©ë¡
    ppt_files = [
        "21ë…„smdì „ëµ.pptx",
        "í•˜ë…¸ì´ ë²•ì¸ 21ë…„ ê²½ì˜ ì „ëµ 20201217.pptx",
        "í•˜ë…¸ì´ ë²•ì¸ 21ë…„ í•˜ë°˜ê¸° ê²½ì˜ ì „ëµ_ ì œì¡°1_REV3.pptx",
        "22ë…„ ì œì¡°1 ê²½ì˜ì „ëµ R2.pptx",
        "í•˜ë…¸ì´ ë²•ì¸ 22ë…„ í•˜ë°˜ê¸° ê²½ì˜ ì „ëµ ì œì¡°1íŒ€ R3.pptx"
    ]

    all_analysis = {
        "total_files": len(ppt_files),
        "total_slides": 0,
        "files": [],
        "all_tasks": [],
        "all_kpis": [],
        "keyword_frequency": {},
        "theme_distribution": defaultdict(int),
        "yearly_summary": {}
    }

    print("=" * 70)
    print("ì „ëµì°¸ê³ ìë£Œ ì‹¬ì¸µ ë¶„ì„ ì‹œì‘")
    print("=" * 70)

    for ppt_file in ppt_files:
        ppt_path = os.path.join(ppt_folder, ppt_file)

        if not os.path.exists(ppt_path):
            print(f"âš  íŒŒì¼ ì—†ìŒ: {ppt_file}")
            continue

        print(f"\nğŸ“„ ë¶„ì„ ì¤‘: {ppt_file}")

        # ìƒì„¸ ë‚´ìš© ì¶”ì¶œ
        file_info = extract_detailed_content(ppt_path)

        if file_info:
            # ì£¼ì œ ë¶„ì„
            analyze_themes(file_info)

            # ê³¼ì œ ë° KPI ì¶”ì¶œ
            tasks, kpis = extract_tasks_and_kpis(file_info)

            file_info["tasks"] = tasks
            file_info["kpis"] = kpis

            all_analysis["files"].append(file_info)
            all_analysis["total_slides"] += file_info["total_slides"]
            all_analysis["all_tasks"].extend(tasks)
            all_analysis["all_kpis"].extend(kpis)

            # ì£¼ì œ ë¶„í¬ ì§‘ê³„
            for slide in file_info["slides"]:
                for theme in slide["themes"]:
                    all_analysis["theme_distribution"][theme] += 1

            # í‚¤ì›Œë“œ ë¹ˆë„ ì§‘ê³„
            for slide in file_info["slides"]:
                for keyword in slide["keywords"]:
                    if len(keyword) >= 2:
                        all_analysis["keyword_frequency"][keyword] = \
                            all_analysis["keyword_frequency"].get(keyword, 0) + 1

            print(f"   âœ“ {file_info['total_slides']}ê°œ ìŠ¬ë¼ì´ë“œ ë¶„ì„ ì™„ë£Œ")
            print(f"   âœ“ {len(tasks)}ê°œ ê³¼ì œ ì¶”ì¶œ")
            print(f"   âœ“ {len(kpis)}ê°œ KPI ì¶”ì¶œ")

    # í‚¤ì›Œë“œ ë¹ˆë„ ìƒìœ„ ì •ë ¬
    sorted_keywords = sorted(
        all_analysis["keyword_frequency"].items(),
        key=lambda x: x[1],
        reverse=True
    )
    all_analysis["top_keywords"] = sorted_keywords[:100]

    # ì£¼ì œ ë¶„í¬ ì •ë ¬
    all_analysis["theme_distribution"] = dict(
        sorted(all_analysis["theme_distribution"].items(),
               key=lambda x: x[1],
               reverse=True)
    )

    # JSON ì €ì¥
    output_file = "ì „ëµìë£Œ_ì‹¬ì¸µë¶„ì„_ê²°ê³¼.json"
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(all_analysis, f, ensure_ascii=False, indent=2)

    print("\n" + "=" * 70)
    print("âœ… ë¶„ì„ ì™„ë£Œ")
    print("=" * 70)
    print(f"ì´ íŒŒì¼: {all_analysis['total_files']}ê°œ")
    print(f"ì´ ìŠ¬ë¼ì´ë“œ: {all_analysis['total_slides']}ê°œ")
    print(f"ì´ ê³¼ì œ: {len(all_analysis['all_tasks'])}ê°œ")
    print(f"ì´ KPI: {len(all_analysis['all_kpis'])}ê°œ")
    print(f"ê³ ìœ  í‚¤ì›Œë“œ: {len(all_analysis['keyword_frequency'])}ê°œ")
    print(f"\nê²°ê³¼ ì €ì¥: {output_file}")
    print("=" * 70)

    # ì£¼ìš” í†µê³„ ì¶œë ¥
    print("\nğŸ“Š ì£¼ì œ ë¶„í¬ (ìƒìœ„ 10ê°œ):")
    for i, (theme, count) in enumerate(list(all_analysis["theme_distribution"].items())[:10], 1):
        print(f"   {i}. {theme}: {count}íšŒ")

    print("\nğŸ”‘ ê³ ë¹ˆë„ í‚¤ì›Œë“œ (ìƒìœ„ 20ê°œ):")
    for i, (keyword, count) in enumerate(all_analysis["top_keywords"][:20], 1):
        print(f"   {i}. {keyword}: {count}íšŒ")

    return all_analysis

if __name__ == "__main__":
    analysis_result = main()
