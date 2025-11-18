#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
옵션 2: 데이터 중심 디자인
대형 숫자, 차트, 미니멀, 화이트 스페이스
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

def create_data_driven_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 미니멀 색상 팔레트
    BLACK = RGBColor(33, 33, 33)
    GRAY = RGBColor(117, 117, 117)
    LIGHT_GRAY = RGBColor(245, 245, 245)
    WHITE = RGBColor(255, 255, 255)
    ACCENT = RGBColor(0, 122, 255)  # iOS Blue
    RED = RGBColor(255, 59, 48)
    GREEN = RGBColor(52, 199, 89)

    # ========== 슬라이드 1: 표지 - 미니멀 ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = WHITE

    # 작은 년도 표시
    year_box = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.5))
    year_frame = year_box.text_frame
    year_para = year_frame.paragraphs[0]
    year_para.text = "2026"
    year_para.font.size = Pt(120)
    year_para.font.bold = True
    year_para.font.color.rgb = LIGHT_GRAY
    year_para.alignment = PP_ALIGN.CENTER

    # 메인 타이틀
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "경영전략"
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = BLACK
    title_para.alignment = PP_ALIGN.CENTER

    # 서브타이틀
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.9), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = "제조1팀 자동화 제조라인 스마트화"
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER

    # 하단 심플 라인
    line = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4), Inches(5.5), Inches(2), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLACK
    line.line.fill.background()

    # ========== 슬라이드 2: 대형 숫자 KPI ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = WHITE

    # 작은 제목
    title_box = slide2.shapes.add_textbox(Inches(1), Inches(0.6), Inches(8), Inches(0.4))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "핵심 목표"
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.font.color.rgb = BLACK
    title_para.alignment = PP_ALIGN.LEFT

    # 3개 대형 KPI 박스
    kpis = [
        {"value": "10%", "label": "가공비 절감", "color": ACCENT, "y": 1.4},
        {"value": "10%", "label": "품질 개선", "color": GREEN, "y": 3.2},
        {"value": "5%", "label": "손실시간 감소", "color": RED, "y": 5.0}
    ]

    for kpi in kpis:
        # 대형 숫자
        num_box = slide2.shapes.add_textbox(Inches(1), Inches(kpi["y"]), Inches(3), Inches(1.2))
        num_frame = num_box.text_frame
        num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        num_para = num_frame.paragraphs[0]
        num_para.text = f"-{kpi['value']}"
        num_para.font.size = Pt(96)
        num_para.font.bold = True
        num_para.font.color.rgb = kpi["color"]
        num_para.alignment = PP_ALIGN.LEFT

        # 라벨
        label_box = slide2.shapes.add_textbox(Inches(4.5), Inches(kpi["y"] + 0.3), Inches(4.5), Inches(0.6))
        label_frame = label_box.text_frame
        label_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        label_para = label_frame.paragraphs[0]
        label_para.text = kpi["label"]
        label_para.font.size = Pt(32)
        label_para.font.color.rgb = BLACK
        label_para.alignment = PP_ALIGN.LEFT

        # 얇은 라인
        line = slide2.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(kpi["y"] + 1.35), Inches(8), Inches(0.01)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = LIGHT_GRAY
        line.line.fill.background()

    # 페이지 번호
    page_box = slide2.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page_frame = page_box.text_frame
    page_frame.text = "02"
    page_frame.paragraphs[0].font.size = Pt(12)
    page_frame.paragraphs[0].font.color.rgb = GRAY
    page_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # ========== 슬라이드 3: 현황 데이터 (차트) ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = WHITE

    # 제목
    title_box = slide3.shapes.add_textbox(Inches(1), Inches(0.6), Inches(8), Inches(0.4))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "현황 분석"
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.font.color.rgb = BLACK
    title_para.alignment = PP_ALIGN.LEFT

    # 좌측: 문제점 (미니멀 리스트)
    problems_box = slide3.shapes.add_textbox(Inches(1), Inches(1.4), Inches(4), Inches(5))
    problems_frame = problems_box.text_frame
    problems_frame.word_wrap = True

    problem_items = [
        ("01", "느린 대응", "설비 이상 DATA 집계 대응으로\n기회 손실 발생"),
        ("02", "반복 작업", "수동적 원인 분석으로\n업무 비효율 심화"),
        ("03", "품질 문제", "불량 원인 대응 지연 및\n재발 방지 체계 미흡")
    ]

    for i, (num, title, desc) in enumerate(problem_items):
        # 번호
        if i == 0:
            num_para = problems_frame.paragraphs[0]
        else:
            num_para = problems_frame.add_paragraph()
        num_para.text = num
        num_para.font.size = Pt(48)
        num_para.font.bold = True
        num_para.font.color.rgb = LIGHT_GRAY
        num_para.space_after = Pt(5)

        # 제목
        title_para = problems_frame.add_paragraph()
        title_para.text = title
        title_para.font.size = Pt(20)
        title_para.font.bold = True
        title_para.font.color.rgb = BLACK
        title_para.space_after = Pt(8)

        # 설명
        desc_para = problems_frame.add_paragraph()
        desc_para.text = desc
        desc_para.font.size = Pt(13)
        desc_para.font.color.rgb = GRAY
        desc_para.space_after = Pt(25)

    # 우측: 간단한 막대 차트 효과 (숫자로 표현)
    metrics_y = Inches(1.4)
    metrics = [
        {"label": "기회손실", "current": "HIGH", "target": "LOW", "color": RED},
        {"label": "설비효율", "current": "85%", "target": "95%", "color": ACCENT},
        {"label": "재발불량", "current": "30건", "target": "0건", "color": GREEN}
    ]

    for i, metric in enumerate(metrics):
        y_pos = metrics_y + (i * Inches(1.5))

        # 라벨
        label_box = slide3.shapes.add_textbox(Inches(5.5), y_pos, Inches(1.5), Inches(0.3))
        label_frame = label_box.text_frame
        label_para = label_frame.paragraphs[0]
        label_para.text = metric["label"]
        label_para.font.size = Pt(14)
        label_para.font.color.rgb = GRAY

        # Current
        curr_box = slide3.shapes.add_textbox(Inches(5.5), y_pos + Inches(0.35), Inches(1.5), Inches(0.6))
        curr_frame = curr_box.text_frame
        curr_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        curr_para = curr_frame.paragraphs[0]
        curr_para.text = metric["current"]
        curr_para.font.size = Pt(36)
        curr_para.font.bold = True
        curr_para.font.color.rgb = LIGHT_GRAY

        # 화살표
        arrow_box = slide3.shapes.add_textbox(Inches(7.2), y_pos + Inches(0.45), Inches(0.5), Inches(0.4))
        arrow_frame = arrow_box.text_frame
        arrow_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        arrow_para = arrow_frame.paragraphs[0]
        arrow_para.text = "→"
        arrow_para.font.size = Pt(32)
        arrow_para.font.color.rgb = metric["color"]
        arrow_para.alignment = PP_ALIGN.CENTER

        # Target
        target_box = slide3.shapes.add_textbox(Inches(7.8), y_pos + Inches(0.35), Inches(1.5), Inches(0.6))
        target_frame = target_box.text_frame
        target_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        target_para = target_frame.paragraphs[0]
        target_para.text = metric["target"]
        target_para.font.size = Pt(36)
        target_para.font.bold = True
        target_para.font.color.rgb = metric["color"]

    # 페이지 번호
    page_box = slide3.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page_frame = page_box.text_frame
    page_frame.text = "03"
    page_frame.paragraphs[0].font.size = Pt(12)
    page_frame.paragraphs[0].font.color.rgb = GRAY
    page_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # ========== 슬라이드 4: 전략 1 - 데이터 중심 ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    slide4.background.fill.solid()
    slide4.background.fill.fore_color.rgb = WHITE

    # 제목
    title_box = slide4.shapes.add_textbox(Inches(1), Inches(0.6), Inches(8), Inches(0.4))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "전략 1"
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.font.color.rgb = BLACK
    title_para.alignment = PP_ALIGN.LEFT

    # 대형 전략명
    strategy_box = slide4.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(1))
    strategy_frame = strategy_box.text_frame
    strategy_para = strategy_frame.paragraphs[0]
    strategy_para.text = "손실 시간 제로화"
    strategy_para.font.size = Pt(52)
    strategy_para.font.bold = True
    strategy_para.font.color.rgb = ACCENT
    strategy_para.alignment = PP_ALIGN.LEFT

    # 핵심 지표
    kpi_box = slide4.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(0.8))
    kpi_frame = kpi_box.text_frame

    kpi_label = kpi_frame.paragraphs[0]
    kpi_label.text = "목표"
    kpi_label.font.size = Pt(14)
    kpi_label.font.color.rgb = GRAY
    kpi_label.space_after = Pt(5)

    kpi_value = kpi_frame.add_paragraph()
    kpi_value.text = "손실시간 5% 감소 → 점당 가공비 직접 절감"
    kpi_value.font.size = Pt(20)
    kpi_value.font.bold = True
    kpi_value.font.color.rgb = BLACK

    # 3개 액션 (미니멀 카드)
    actions_y = Inches(3.7)
    actions = [
        {"num": "1", "title": "자동분석 TOOL", "desc": "MES DATA 활용\n이상 감지 및 C/T 모니터링"},
        {"num": "2", "title": "순간유실 가시화", "desc": "초 단위 손실 기록\n모바일 앱 원터치 입력"},
        {"num": "3", "title": "TOP 10 개선", "desc": "손실 항목 순위화\n집중 타격 전략"}
    ]

    action_width = Inches(2.5)
    action_spacing = Inches(0.4)
    action_start = Inches(1)

    for i, action in enumerate(actions):
        x_pos = action_start + (i * (action_width + action_spacing))

        # 번호 (큰 숫자)
        num_box = slide4.shapes.add_textbox(x_pos, actions_y, action_width, Inches(0.8))
        num_frame = num_box.text_frame
        num_para = num_frame.paragraphs[0]
        num_para.text = action["num"]
        num_para.font.size = Pt(72)
        num_para.font.bold = True
        num_para.font.color.rgb = LIGHT_GRAY
        num_para.alignment = PP_ALIGN.LEFT

        # 제목
        title_box = slide4.shapes.add_textbox(x_pos, actions_y + Inches(0.9), action_width, Inches(0.4))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = action["title"]
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = BLACK

        # 설명
        desc_box = slide4.shapes.add_textbox(x_pos, actions_y + Inches(1.4), action_width, Inches(1))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_para = desc_frame.paragraphs[0]
        desc_para.text = action["desc"]
        desc_para.font.size = Pt(12)
        desc_para.font.color.rgb = GRAY

        # 하단 라인
        line = slide4.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x_pos, actions_y + Inches(2.5), action_width, Inches(0.02)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT
        line.line.fill.background()

    # 페이지 번호
    page_box = slide4.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page_frame = page_box.text_frame
    page_frame.text = "04"
    page_frame.paragraphs[0].font.size = Pt(12)
    page_frame.paragraphs[0].font.color.rgb = GRAY
    page_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # ========== 슬라이드 5: 전략 2 & 3 - 콤팩트 ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    slide5.background.fill.solid()
    slide5.background.fill.fore_color.rgb = WHITE

    # 전략 2
    s2_title = slide5.shapes.add_textbox(Inches(1), Inches(0.6), Inches(4), Inches(0.4))
    s2_title_frame = s2_title.text_frame
    s2_title_para = s2_title_frame.paragraphs[0]
    s2_title_para.text = "전략 2"
    s2_title_para.font.size = Pt(18)
    s2_title_para.font.bold = True
    s2_title_para.font.color.rgb = BLACK

    s2_name = slide5.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(0.6))
    s2_name_frame = s2_name.text_frame
    s2_name_para = s2_name_frame.paragraphs[0]
    s2_name_para.text = "불량 재발 Zero"
    s2_name_para.font.size = Pt(36)
    s2_name_para.font.bold = True
    s2_name_para.font.color.rgb = GREEN

    s2_kpi = slide5.shapes.add_textbox(Inches(1), Inches(1.7), Inches(4), Inches(0.4))
    s2_kpi_frame = s2_kpi.text_frame
    s2_kpi_para = s2_kpi_frame.paragraphs[0]
    s2_kpi_para.text = "재발 불량 30% ↓ / 품질 불량 10% ↓"
    s2_kpi_para.font.size = Pt(14)
    s2_kpi_para.font.color.rgb = GRAY

    s2_actions = [
        "• 불량 사진 + MES 이력 즉시 공유",
        "• 전 조 자동 알람 시스템",
        "• 재발 불량 자동 추적 DB"
    ]

    s2_content = slide5.shapes.add_textbox(Inches(1), Inches(2.3), Inches(4), Inches(1.5))
    s2_content_frame = s2_content.text_frame
    s2_content_frame.word_wrap = True
    for i, action in enumerate(s2_actions):
        if i == 0:
            para = s2_content_frame.paragraphs[0]
        else:
            para = s2_content_frame.add_paragraph()
        para.text = action
        para.font.size = Pt(13)
        para.font.color.rgb = BLACK
        para.space_after = Pt(8)

    # 전략 3
    s3_title = slide5.shapes.add_textbox(Inches(5.5), Inches(0.6), Inches(4), Inches(0.4))
    s3_title_frame = s3_title.text_frame
    s3_title_para = s3_title_frame.paragraphs[0]
    s3_title_para.text = "전략 3"
    s3_title_para.font.size = Pt(18)
    s3_title_para.font.bold = True
    s3_title_para.font.color.rgb = BLACK

    s3_name = slide5.shapes.add_textbox(Inches(5.5), Inches(1), Inches(4), Inches(0.6))
    s3_name_frame = s3_name.text_frame
    s3_name_para = s3_name_frame.paragraphs[0]
    s3_name_para.text = "설비 CAPA 증가"
    s3_name_para.font.size = Pt(36)
    s3_name_para.font.bold = True
    s3_name_para.font.color.rgb = RED

    s3_kpi = slide5.shapes.add_textbox(Inches(5.5), Inches(1.7), Inches(4), Inches(0.4))
    s3_kpi_frame = s3_kpi.text_frame
    s3_kpi_para = s3_kpi_frame.paragraphs[0]
    s3_kpi_para.text = "CAPA 증가 / OVERTIME 감소"
    s3_kpi_para.font.size = Pt(14)
    s3_kpi_para.font.color.rgb = GRAY

    s3_actions = [
        "• C/T 단축 DEEP 분석",
        "• 설비 효율 성능/PM 강화",
        "• 공정 최적화 (RADIAL2, SMD)"
    ]

    s3_content = slide5.shapes.add_textbox(Inches(5.5), Inches(2.3), Inches(4), Inches(1.5))
    s3_content_frame = s3_content.text_frame
    s3_content_frame.word_wrap = True
    for i, action in enumerate(s3_actions):
        if i == 0:
            para = s3_content_frame.paragraphs[0]
        else:
            para = s3_content_frame.add_paragraph()
        para.text = action
        para.font.size = Pt(13)
        para.font.color.rgb = BLACK
        para.space_after = Pt(8)

    # 하단: 종합 효과
    effect_y = Inches(4.3)

    effect_title = slide5.shapes.add_textbox(Inches(1), effect_y, Inches(8), Inches(0.3))
    effect_title_frame = effect_title.text_frame
    effect_title_para = effect_title_frame.paragraphs[0]
    effect_title_para.text = "종합 효과"
    effect_title_para.font.size = Pt(18)
    effect_title_para.font.bold = True
    effect_title_para.font.color.rgb = BLACK

    # 대형 숫자 박스 3개
    final_kpis = [
        {"value": "10%", "label": "가공비 절감", "color": ACCENT},
        {"value": "10%", "label": "품질 개선", "color": GREEN},
        {"value": "5%", "label": "손실시간", "color": RED}
    ]

    kpi_width = Inches(2.5)
    kpi_spacing = Inches(0.4)
    kpi_start = Inches(1)

    for i, kpi in enumerate(final_kpis):
        x_pos = kpi_start + (i * (kpi_width + kpi_spacing))

        # 값
        val_box = slide5.shapes.add_textbox(x_pos, effect_y + Inches(0.5), kpi_width, Inches(0.8))
        val_frame = val_box.text_frame
        val_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        val_para = val_frame.paragraphs[0]
        val_para.text = f"-{kpi['value']}"
        val_para.font.size = Pt(64)
        val_para.font.bold = True
        val_para.font.color.rgb = kpi["color"]
        val_para.alignment = PP_ALIGN.CENTER

        # 라벨
        label_box = slide5.shapes.add_textbox(x_pos, effect_y + Inches(1.4), kpi_width, Inches(0.3))
        label_frame = label_box.text_frame
        label_para = label_frame.paragraphs[0]
        label_para.text = kpi["label"]
        label_para.font.size = Pt(14)
        label_para.font.color.rgb = GRAY
        label_para.alignment = PP_ALIGN.CENTER

    # 페이지 번호
    page_box = slide5.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page_frame = page_box.text_frame
    page_frame.text = "05"
    page_frame.paragraphs[0].font.size = Pt(12)
    page_frame.paragraphs[0].font.color.rgb = GRAY
    page_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # ========== 슬라이드 6: 로드맵 (타임라인 숫자) ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    slide6.background.fill.solid()
    slide6.background.fill.fore_color.rgb = WHITE

    # 제목
    title_box = slide6.shapes.add_textbox(Inches(1), Inches(0.6), Inches(8), Inches(0.4))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "실행 로드맵"
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.font.color.rgb = BLACK
    title_para.alignment = PP_ALIGN.LEFT

    # 4분기 타임라인
    quarters = [
        {"q": "Q1", "label": "개발/구축", "tasks": ["TOOL 개발", "모바일 앱", "시스템 구축"]},
        {"q": "Q2", "label": "시범 운영", "tasks": ["파일럿 테스트", "피드백 반영", "개선"]},
        {"q": "Q3", "label": "전사 확대", "tasks": ["전체 적용", "교육 실시", "모니터링"]},
        {"q": "Q4", "label": "목표 달성", "tasks": ["성과 점검", "최적화", "차년도 계획"]}
    ]

    quarter_width = Inches(2)
    quarter_spacing = Inches(0.25)
    quarter_start = Inches(1)
    quarter_y = Inches(1.5)

    for i, quarter in enumerate(quarters):
        x_pos = quarter_start + (i * (quarter_width + quarter_spacing))

        # Q 번호 (대형)
        q_box = slide6.shapes.add_textbox(x_pos, quarter_y, quarter_width, Inches(0.7))
        q_frame = q_box.text_frame
        q_para = q_frame.paragraphs[0]
        q_para.text = quarter["q"]
        q_para.font.size = Pt(56)
        q_para.font.bold = True
        q_para.font.color.rgb = LIGHT_GRAY
        q_para.alignment = PP_ALIGN.LEFT

        # 라벨
        label_box = slide6.shapes.add_textbox(x_pos, quarter_y + Inches(0.7), quarter_width, Inches(0.3))
        label_frame = label_box.text_frame
        label_para = label_frame.paragraphs[0]
        label_para.text = quarter["label"]
        label_para.font.size = Pt(16)
        label_para.font.bold = True
        label_para.font.color.rgb = BLACK

        # 태스크
        task_box = slide6.shapes.add_textbox(x_pos, quarter_y + Inches(1.1), quarter_width, Inches(1.2))
        task_frame = task_box.text_frame
        task_frame.word_wrap = True
        for j, task in enumerate(quarter["tasks"]):
            if j == 0:
                para = task_frame.paragraphs[0]
            else:
                para = task_frame.add_paragraph()
            para.text = f"• {task}"
            para.font.size = Pt(11)
            para.font.color.rgb = GRAY
            para.space_after = Pt(4)

        # 연결 라인 (마지막 제외)
        if i < 3:
            line = slide6.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x_pos + quarter_width + Inches(0.05), quarter_y + Inches(0.3),
                Inches(0.15), Inches(0.02)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = LIGHT_GRAY
            line.line.fill.background()

    # 하단 박스: 최종 목표
    final_y = Inches(4.5)

    final_title = slide6.shapes.add_textbox(Inches(1), final_y, Inches(8), Inches(0.3))
    final_title_frame = final_title.text_frame
    final_title_para = final_title_frame.paragraphs[0]
    final_title_para.text = "2026 최종 목표"
    final_title_para.font.size = Pt(18)
    final_title_para.font.bold = True
    final_title_para.font.color.rgb = BLACK

    final_goals = [
        "기회손실 최소화",
        "점당 가공비 직접 절감",
        "OH 달성",
        "OVERTIME 감소"
    ]

    final_content = slide6.shapes.add_textbox(Inches(1), final_y + Inches(0.5), Inches(8), Inches(1.5))
    final_content_frame = final_content.text_frame
    final_content_frame.word_wrap = True

    for i, goal in enumerate(final_goals):
        if i == 0:
            para = final_content_frame.paragraphs[0]
        else:
            para = final_content_frame.add_paragraph()
        para.text = f"▪ {goal}"
        para.font.size = Pt(16)
        para.font.color.rgb = BLACK
        para.space_after = Pt(10)

    # 페이지 번호
    page_box = slide6.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page_frame = page_box.text_frame
    page_frame.text = "06"
    page_frame.paragraphs[0].font.size = Pt(12)
    page_frame.paragraphs[0].font.color.rgb = GRAY
    page_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # 저장
    output_file = '옵션2_데이터중심_2026전략.pptx'
    prs.save(output_file)
    print(f"✅ 옵션 2 (데이터 중심) PPT 생성 완료: {output_file}")
    return output_file

if __name__ == "__main__":
    create_data_driven_presentation()
