#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
전략 PPT 파일 분석 스크립트
모든 슬라이드의 텍스트 내용을 추출하여 분석
"""

from pptx import Presentation
import os
import json

def extract_ppt_content(ppt_path):
    """PPT 파일에서 모든 텍스트 추출"""
    try:
        prs = Presentation(ppt_path)

        content = {
            "file_name": os.path.basename(ppt_path),
            "total_slides": len(prs.slides),
            "slides": []
        }

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = {
                "slide_number": slide_num,
                "texts": []
            }

            # 모든 도형에서 텍스트 추출
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content["texts"].append(shape.text.strip())

                # 테이블 체크
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                slide_content["texts"].append(f"[TABLE] {cell.text.strip()}")

            content["slides"].append(slide_content)

        return content
    except Exception as e:
        return {"error": str(e), "file_name": os.path.basename(ppt_path)}

def analyze_all_ppts(folder_path):
    """폴더 내 모든 PPT 분석"""
    ppt_files = [f for f in os.listdir(folder_path) if f.endswith('.pptx')]

    all_content = []

    for ppt_file in sorted(ppt_files):
        ppt_path = os.path.join(folder_path, ppt_file)
        print(f"\n{'='*80}")
        print(f"📊 분석 중: {ppt_file}")
        print(f"{'='*80}")

        content = extract_ppt_content(ppt_path)

        if "error" in content:
            print(f"❌ 오류: {content['error']}")
        else:
            print(f"✅ 총 슬라이드: {content['total_slides']}개")

            # 주요 내용 출력
            for slide in content["slides"][:3]:  # 처음 3개 슬라이드만
                print(f"\n--- 슬라이드 {slide['slide_number']} ---")
                for text in slide["texts"][:5]:  # 각 슬라이드에서 최대 5개 텍스트
                    print(f"  • {text[:100]}...")  # 최대 100자

        all_content.append(content)

    # JSON으로 저장
    output_file = "전략PPT_분석결과.json"
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(all_content, f, ensure_ascii=False, indent=2)

    print(f"\n\n{'='*80}")
    print(f"✅ 분석 완료! 결과 저장: {output_file}")
    print(f"{'='*80}")

    return all_content

if __name__ == "__main__":
    folder_path = "전략자료참고"
    analyze_all_ppts(folder_path)
