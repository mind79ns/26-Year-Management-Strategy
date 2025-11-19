#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
2026년 제조1팀 경영전략 - 최종 통합 버전
과거 분석 + 차별화 포인트 + 참고디자인 스타일 + 프리미엄 요소 통합
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_progress_bar(slide, x, y, width, height, percentage, color, bg_color):
    """프로그레스 바 생성"""
    bg_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
    bg_bar.fill.solid()
    bg_bar.fill.fore_color.rgb = bg_color
    bg_bar.line.fill.background()

    if percentage > 0:
        progress_width = width * (percentage / 100)
        if progress_width > Inches(0.1):
            progress_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, progress_width, height)
            progress_bar.fill.solid()
            progress_bar.fill.fore_color.rgb = color
            progress_bar.line.fill.background()

    circle_size = Inches(0.5)
    circle_x = x + (width * (percentage / 100)) - circle_size / 2
    circle_y = y - Inches(0.1)

    if percentage > 0:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, circle_x, circle_y, circle_size, circle_size)
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        percent_box = slide.shapes.add_textbox(circle_x, circle_y, circle_size, circle_size)
        percent_frame = percent_box.text_frame
        percent_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        percent_para = percent_frame.paragraphs[0]
        percent_para.text = f"{int(percentage)}%"
        percent_para.font.size = Pt(11)
        percent_para.font.bold = True
        percent_para.font.color.rgb = RGBColor(255, 255, 255)
        percent_para.alignment = PP_ALIGN.CENTER

def create_circular_gauge(slide, center_x, center_y, radius, percentage, color):
    """원형 게이지"""
    bg_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, center_x - radius, center_y - radius, radius * 2, radius * 2)
    bg_circle.fill.solid()
    bg_circle.fill.fore_color.rgb = RGBColor(240, 240, 240)
    bg_circle.line.color.rgb = RGBColor(220, 220, 220)
    bg_circle.line.width = Pt(2)

    inner_radius = radius * 0.7
    inner_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, center_x - inner_radius, center_y - inner_radius, inner_radius * 2, inner_radius * 2)
    inner_circle.fill.solid()
    inner_circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
    inner_circle.line.fill.background()

    percent_box = slide.shapes.add_textbox(center_x - inner_radius, center_y - inner_radius, inner_radius * 2, inner_radius * 2)
    percent_frame = percent_box.text_frame
    percent_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    percent_para = percent_frame.paragraphs[0]
    percent_para.text = f"{int(percentage)}%"
    percent_para.font.size = Pt(36)
    percent_para.font.bold = True
    percent_para.font.color.rgb = color
    percent_para.alignment = PP_ALIGN.CENTER

def create_premium_box(slide, x, y, width, height, color, shadow=True):
    """프리미엄 박스"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()

    if shadow:
        box.shadow.inherit = False
        box.shadow.visible = True
        box.shadow.distance = Pt(3)
        box.shadow.angle = 45
        box.shadow.blur_radius = Pt(8)

    return box

def create_final_integrated_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 색상 팔레트
    PRIMARY_BLUE = RGBColor(41, 128, 185)
    LIGHT_BLUE = RGBColor(93, 173, 226)
    DARK_BLUE = RGBColor(21, 67, 96)
    NAVY = RGBColor(15, 32, 56)
    GOLD = RGBColor(212, 175, 55)
    GREEN = RGBColor(39, 174, 96)
    ORANGE = RGBColor(230, 126, 34)
    RED = RGBColor(231, 76, 60)
    GRAY = RGBColor(149, 165, 166)
    LIGHT_GRAY = RGBColor(220, 220, 220)
    BG_GRAY = RGBColor(245, 245, 245)
    WHITE = RGBColor(255, 255, 255)
    DARK_TEXT = RGBColor(33, 33, 33)

    # ========== 슬라이드 1: 표지 (프리미엄 + 임팩트) ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = NAVY

    # 좌측 골드 액센트
    accent = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    # 메인 타이틀
    title_box = slide1.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(1.2))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "2026년 경영전략"
    title_para.font.size = Pt(66)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.LEFT

    # 서브타이틀
    subtitle_box = slide1.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(7), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = "디지털 전환을 통한 스마트 팩토리 실현"
    subtitle_para.font.size = Pt(22)
    subtitle_para.font.color.rgb = LIGHT_BLUE
    subtitle_para.alignment = PP_ALIGN.LEFT

    # 골드 라인
    gold_line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.5), Inches(4), Inches(0.03))
    gold_line.fill.solid()
    gold_line.fill.fore_color.rgb = GOLD
    gold_line.line.fill.background()

    # 3개 KPI 카드
    kpis = [
        {"label": "가공비 절감", "value": "10%", "color": PRIMARY_BLUE},
        {"label": "품질 개선", "value": "10%", "color": GREEN},
        {"label": "손실시간", "value": "5%", "color": ORANGE}
    ]

    card_y = Inches(5.3)
    card_width = Inches(2)
    card_spacing = Inches(0.25)
    card_start = Inches(1.5)

    for i, kpi in enumerate(kpis):
        x_pos = card_start + i * (card_width + card_spacing)

        card = create_premium_box(slide1, x_pos, card_y, card_width, Inches(1.4), RGBColor(30, 55, 85), shadow=True)

        label_box = slide1.shapes.add_textbox(x_pos + Inches(0.2), card_y + Inches(0.2), card_width - Inches(0.4), Inches(0.3))
        label_frame = label_box.text_frame
        label_para = label_frame.paragraphs[0]
        label_para.text = kpi["label"]
        label_para.font.size = Pt(14)
        label_para.font.color.rgb = LIGHT_GRAY
        label_para.alignment = PP_ALIGN.LEFT

        value_box = slide1.shapes.add_textbox(x_pos + Inches(0.2), card_y + Inches(0.55), card_width - Inches(0.4), Inches(0.6))
        value_frame = value_box.text_frame
        value_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        value_para = value_frame.paragraphs[0]
        value_para.text = f"-{kpi['value']}"
        value_para.font.size = Pt(36)
        value_para.font.bold = True
        value_para.font.color.rgb = kpi["color"]
        value_para.alignment = PP_ALIGN.LEFT

    # 팀 정보
    team_box = slide1.shapes.add_textbox(Inches(1.5), Inches(6.8), Inches(7), Inches(0.4))
    team_frame = team_box.text_frame
    team_para = team_frame.paragraphs[0]
    team_para.text = "제조1팀  |  Manufacturing Team 1  |  2026 Strategy"
    team_para.font.size = Pt(14)
    team_para.font.color.rgb = LIGHT_BLUE
    team_para.alignment = PP_ALIGN.LEFT

    # ========== 슬라이드 2: 과거 전략 회고 (신규) ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = WHITE

    # 상단 바
    top_bar = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PRIMARY_BLUE
    top_bar.line.fill.background()

    # 제목
    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title2_frame = title2.text_frame
    title2_para = title2_frame.paragraphs[0]
    title2_para.text = "01. 과거 전략 회고 (2021~2022)"
    title2_para.font.size = Pt(28)
    title2_para.font.bold = True
    title2_para.font.color.rgb = DARK_BLUE

    # 좌측: 과거 성과
    past_title = slide2.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.4), Inches(0.4))
    pt_frame = past_title.text_frame
    pt_para = pt_frame.paragraphs[0]
    pt_para.text = "✓ 주요 성과"
    pt_para.font.size = Pt(20)
    pt_para.font.bold = True
    pt_para.font.color.rgb = PRIMARY_BLUE

    achievements = [
        "MES System 구축 및 정착",
        "WORST LINE/MODEL 집중 개선",
        "SMD 설비 유실 개선 추진",
        "평가가동율 지속 향상"
    ]

    ach_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.4), Inches(2))
    ach_frame = ach_box.text_frame
    ach_frame.word_wrap = True

    for i, ach in enumerate(achievements):
        if i == 0:
            para = ach_frame.paragraphs[0]
        else:
            para = ach_frame.add_paragraph()
        para.text = f"• {ach}"
        para.font.size = Pt(14)
        para.font.color.rgb = DARK_TEXT
        para.space_after = Pt(10)

    # 지속 과제
    cont_title = slide2.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(4.4), Inches(0.4))
    ct_frame = cont_title.text_frame
    ct_para = ct_frame.paragraphs[0]
    ct_para.text = "⚠ 지속 과제 (미해결)"
    ct_para.font.size = Pt(20)
    ct_para.font.bold = True
    ct_para.font.color.rgb = RED

    challenges = [
        "느린 대응 속도 (주간 단위)",
        "수동적 원인 분석",
        "재발 불량 반복",
        "순간유실 미관리"
    ]

    chal_box = slide2.shapes.add_textbox(Inches(0.5), Inches(4.7), Inches(4.4), Inches(1.8))
    chal_frame = chal_box.text_frame
    chal_frame.word_wrap = True

    for i, chal in enumerate(challenges):
        if i == 0:
            para = chal_frame.paragraphs[0]
        else:
            para = chal_frame.add_paragraph()
        para.text = f"• {chal}"
        para.font.size = Pt(14)
        para.font.color.rgb = DARK_TEXT
        para.space_after = Pt(10)

    # 우측: 핵심 키워드 (과거 분석 결과)
    keyword_title = slide2.shapes.add_textbox(Inches(5.1), Inches(1.3), Inches(4.4), Inches(0.4))
    kt_frame = keyword_title.text_frame
    kt_para = kt_frame.paragraphs[0]
    kt_para.text = "📊 과거 전략 핵심 키워드 TOP 5"
    kt_para.font.size = Pt(18)
    kt_para.font.bold = True
    kt_para.font.color.rgb = DARK_BLUE

    keywords = [
        ("개선", 202, PRIMARY_BLUE),
        ("유실", 115, ORANGE),
        ("설비", 112, GREEN),
        ("관리", 106, GRAY),
        ("목표", 92, RED)
    ]

    keyword_y = Inches(1.9)
    for i, (word, count, color) in enumerate(keywords):
        y_pos = keyword_y + i * Inches(0.8)

        # 키워드
        word_box = slide2.shapes.add_textbox(Inches(5.1), y_pos, Inches(1.5), Inches(0.3))
        w_frame = word_box.text_frame
        w_para = w_frame.paragraphs[0]
        w_para.text = word
        w_para.font.size = Pt(16)
        w_para.font.bold = True
        w_para.font.color.rgb = DARK_TEXT

        # 프로그레스 바 (횟수를 %로 변환)
        percentage = (count / 202) * 100
        create_progress_bar(slide2, Inches(5.1), y_pos + Inches(0.35), Inches(3.5), Inches(0.2), percentage, color, LIGHT_GRAY)

        # 횟수
        count_box = slide2.shapes.add_textbox(Inches(8.7), y_pos + Inches(0.05), Inches(0.8), Inches(0.5))
        c_frame = count_box.text_frame
        c_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        c_para = c_frame.paragraphs[0]
        c_para.text = f"{count}회"
        c_para.font.size = Pt(12)
        c_para.font.color.rgb = GRAY
        c_para.alignment = PP_ALIGN.RIGHT

    # 하단: 분석 출처
    source_box = slide2.shapes.add_textbox(Inches(5.1), Inches(6.2), Inches(4.4), Inches(0.8))
    s_frame = source_box.text_frame
    s_para = s_frame.paragraphs[0]
    s_para.text = "📁 분석 출처"
    s_para.font.size = Pt(12)
    s_para.font.bold = True
    s_para.font.color.rgb = GRAY
    s_para.space_after = Pt(5)

    s_detail = s_frame.add_paragraph()
    s_detail.text = "5개 과거 전략 PPT / 72개 슬라이드\n6,547줄 데이터 분석 완료"
    s_detail.font.size = Pt(10)
    s_detail.font.color.rgb = GRAY

    # 페이지 번호
    page2 = slide2.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page2.text_frame.text = "02"
    page2.text_frame.paragraphs[0].font.size = Pt(11)
    page2.text_frame.paragraphs[0].font.color.rgb = GRAY
    page2.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # ========== 슬라이드 3: 2026 전략 배경 및 차별화 ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = WHITE

    top_bar3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    top_bar3.fill.solid()
    top_bar3.fill.fore_color.rgb = GOLD
    top_bar3.line.fill.background()

    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title3_frame = title3.text_frame
    title3_para = title3_frame.paragraphs[0]
    title3_para.text = "02. 2026 전략 목표 및 차별화 포인트"
    title3_para.font.size = Pt(28)
    title3_para.font.bold = True
    title3_para.font.color.rgb = DARK_BLUE

    # 중앙: 차별화 포인트 (대형)
    diff_box = create_premium_box(slide3, Inches(1), Inches(1.5), Inches(8), Inches(2.5), PRIMARY_BLUE, shadow=True)

    df_frame = diff_box.text_frame
    df_frame.margin_left = Inches(0.4)
    df_frame.margin_top = Inches(0.3)

    df_title = df_frame.paragraphs[0]
    df_title.text = "🚀 과거 대비 혁신 포인트"
    df_title.font.size = Pt(24)
    df_title.font.bold = True
    df_title.font.color.rgb = WHITE
    df_title.space_after = Pt(15)

    innovations = [
        "실시간 대응: 주간 단위 → 초 단위 (1000배 빠름)",
        "순간유실 가시화: 완전 신규 개념 (모바일 앱 원터치)",
        "불량 재발 Zero: 사진 공유 + 자동 알람 시스템",
        "데이터 기반: 자동 분석 TOOL → 순위화 → TOP 10 집중"
    ]

    for innov in innovations:
        para = df_frame.add_paragraph()
        para.text = f"▪ {innov}"
        para.font.size = Pt(15)
        para.font.color.rgb = WHITE
        para.space_after = Pt(10)

    # 하단: Before / After 비교
    comp_title = slide3.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(0.3))
    comp_frame = comp_title.text_frame
    comp_para = comp_frame.paragraphs[0]
    comp_para.text = "Before / After 비교"
    comp_para.font.size = Pt(20)
    comp_para.font.bold = True
    comp_para.font.color.rgb = DARK_BLUE
    comp_para.alignment = PP_ALIGN.CENTER

    comparisons = [
        {"item": "대응 속도", "before": "주간 단위", "after": "실시간 (초)", "improve": "1000배"},
        {"item": "데이터 입력", "before": "수기 기록", "after": "원터치", "improve": "10배 간편"},
        {"item": "불량 공유", "before": "구두 전달", "after": "자동 알람", "improve": "100%"},
        {"item": "재발 방지", "before": "수동 추적", "after": "자동 표시", "improve": "Zero 목표"}
    ]

    comp_y = Inches(4.8)
    comp_width = Inches(2)
    comp_spacing = Inches(0.15)
    comp_start = Inches(0.5)

    # 헤더
    headers = ["구분", "과거 (21~22)", "2026", "개선도"]
    for i, header in enumerate(headers):
        x_pos = comp_start + i * (comp_width + comp_spacing)
        h_box = slide3.shapes.add_textbox(x_pos, comp_y, comp_width, Inches(0.35))
        h_frame = h_box.text_frame
        h_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        h_para = h_frame.paragraphs[0]
        h_para.text = header
        h_para.font.size = Pt(12)
        h_para.font.bold = True
        h_para.font.color.rgb = WHITE
        h_para.alignment = PP_ALIGN.CENTER

        # 헤더 배경
        h_bg = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, comp_y, comp_width, Inches(0.35))
        h_bg.fill.solid()
        h_bg.fill.fore_color.rgb = DARK_BLUE
        h_bg.line.fill.background()
        h_bg.z_order = 1
        h_box.z_order = 2

    # 데이터 행
    for j, comp in enumerate(comparisons):
        row_y = comp_y + Inches(0.45) + j * Inches(0.4)

        values = [comp["item"], comp["before"], comp["after"], comp["improve"]]
        for i, value in enumerate(values):
            x_pos = comp_start + i * (comp_width + comp_spacing)

            # 배경 (교차 색상)
            bg_color = BG_GRAY if j % 2 == 0 else WHITE
            bg = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, row_y, comp_width, Inches(0.35))
            bg.fill.solid()
            bg.fill.fore_color.rgb = bg_color
            bg.line.color.rgb = LIGHT_GRAY
            bg.line.width = Pt(0.5)

            # 텍스트
            v_box = slide3.shapes.add_textbox(x_pos + Inches(0.1), row_y, comp_width - Inches(0.2), Inches(0.35))
            v_frame = v_box.text_frame
            v_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            v_para = v_frame.paragraphs[0]
            v_para.text = value
            v_para.font.size = Pt(11)

            if i == 0:
                v_para.font.bold = True
                v_para.font.color.rgb = DARK_BLUE
            elif i == 3:
                v_para.font.bold = True
                v_para.font.color.rgb = RED
            else:
                v_para.font.color.rgb = DARK_TEXT

            v_para.alignment = PP_ALIGN.CENTER

    page3 = slide3.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page3.text_frame.text = "03"
    page3.text_frame.paragraphs[0].font.size = Pt(11)
    page3.text_frame.paragraphs[0].font.color.rgb = GRAY
    page3.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # ========== 슬라이드 4-6: 3대 전략 (프로그레스 바 스타일) ==========
    strategies = [
        {
            "num": "03",
            "title": "전략 1: 손실 시간 제로화 프로젝트",
            "color": PRIMARY_BLUE,
            "innovation": "💡 혁신: 순간유실 가시화 (완전 신규 개념)",
            "actions": [
                {"title": "자동분석 TOOL", "progress": 100, "items": ["MES DATA 활용", "이상 감지", "C/T 모니터링"]},
                {"title": "순간유실 가시화", "progress": 100, "items": ["초 단위 기록", "모바일 원터치", "LINE별 분석"]},
                {"title": "TOP 10 집중", "progress": 100, "items": ["순위화", "WORST 타격", "리포트 자동화"]}
            ],
            "kpi": "손실시간 5% 감소 → 점당 가공비 직접 절감"
        },
        {
            "num": "04",
            "title": "전략 2: 불량 재발 Zero 챌린지",
            "color": GREEN,
            "innovation": "💡 혁신: 사진 공유 시스템 (완전 신규)",
            "actions": [
                {"title": "즉시 FEEDBACK", "progress": 100, "items": ["사진+MES 연동", "자동 기록", "즉시 업로드"]},
                {"title": "전 조 자동 알람", "progress": 100, "items": ["태그 분류", "메모 공유", "자동 전달"]},
                {"title": "재발 불량 추적", "progress": 100, "items": ["재발 표시", "월별 추적", "Zero KPI화"]}
            ],
            "kpi": "재발 불량 30% 감소 / 품질 불량 10% 감소"
        },
        {
            "num": "05",
            "title": "전략 3: 설비 CAPA 증가",
            "color": ORANGE,
            "innovation": "✓ 계승: 과거 성공 요소 강화",
            "actions": [
                {"title": "C/T 단축", "progress": 100, "items": ["DEEP 분석", "최단거리", "1초 줄이기"]},
                {"title": "BASE 강화", "progress": 100, "items": ["효율 유지", "PM 강화", "모니터링"]},
                {"title": "공정 최적화", "progress": 100, "items": ["RADIAL2", "SMD LAY OUT", "배치 효율"]}
            ],
            "kpi": "CAPA 증가 / OVERTIME 감소 / OH 달성"
        }
    ]

    for strat_idx, strategy in enumerate(strategies):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = WHITE

        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = strategy["color"]
        top_bar.line.fill.background()

        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
        t_frame = title.text_frame
        t_para = t_frame.paragraphs[0]
        t_para.text = f"{strategy['num']}. {strategy['title']}"
        t_para.font.size = Pt(28)
        t_para.font.bold = True
        t_para.font.color.rgb = DARK_BLUE

        # 혁신 포인트
        innov_box = create_premium_box(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(0.6), strategy["color"], shadow=False)
        i_frame = innov_box.text_frame
        i_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        i_frame.margin_left = Inches(0.3)
        i_para = i_frame.paragraphs[0]
        i_para.text = strategy["innovation"]
        i_para.font.size = Pt(18)
        i_para.font.bold = True
        i_para.font.color.rgb = WHITE

        # 3개 액션 카드
        action_y = Inches(2.1)
        action_width = Inches(2.8)
        action_spacing = Inches(0.3)
        action_start = Inches(0.5)

        for i, action in enumerate(strategy["actions"]):
            x_pos = action_start + i * (action_width + action_spacing)

            card = create_premium_box(slide, x_pos, action_y, action_width, Inches(3.6), BG_GRAY, shadow=True)

            # 타이틀
            title_box = slide.shapes.add_textbox(x_pos + Inches(0.2), action_y + Inches(0.2), action_width - Inches(0.4), Inches(0.5))
            tf = title_box.text_frame
            tf.word_wrap = True
            tp = tf.paragraphs[0]
            tp.text = action["title"]
            tp.font.size = Pt(16)
            tp.font.bold = True
            tp.font.color.rgb = DARK_BLUE
            tp.alignment = PP_ALIGN.CENTER

            # 프로그레스 바
            prog_y = action_y + Inches(0.8)
            create_progress_bar(slide, x_pos + Inches(0.3), prog_y, action_width - Inches(0.6), Inches(0.25), action["progress"], strategy["color"], LIGHT_GRAY)

            # 항목들
            items_box = slide.shapes.add_textbox(x_pos + Inches(0.2), action_y + Inches(1.3), action_width - Inches(0.4), Inches(2))
            if_frame = items_box.text_frame
            if_frame.word_wrap = True

            for j, item in enumerate(action["items"]):
                if j == 0:
                    para = if_frame.paragraphs[0]
                else:
                    para = if_frame.add_paragraph()
                para.text = f"• {item}"
                para.font.size = Pt(12)
                para.font.color.rgb = DARK_TEXT
                para.space_after = Pt(8)

        # KPI 박스
        kpi_box = create_premium_box(slide, Inches(0.5), Inches(6), Inches(9), Inches(0.7), RED, shadow=False)
        kpi_frame = kpi_box.text_frame
        kpi_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        kpi_para = kpi_frame.paragraphs[0]
        kpi_para.text = f"🎯 목표 KPI: {strategy['kpi']}"
        kpi_para.font.size = Pt(18)
        kpi_para.font.bold = True
        kpi_para.font.color.rgb = WHITE
        kpi_para.alignment = PP_ALIGN.CENTER

        # 페이지 번호
        page = slide.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
        page.text_frame.text = f"0{strat_idx + 4}"
        page.text_frame.paragraphs[0].font.size = Pt(11)
        page.text_frame.paragraphs[0].font.color.rgb = GRAY
        page.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # ========== 슬라이드 7: 실행 로드맵 ==========
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    slide7.background.fill.solid()
    slide7.background.fill.fore_color.rgb = WHITE

    top_bar7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    top_bar7.fill.solid()
    top_bar7.fill.fore_color.rgb = NAVY
    top_bar7.line.fill.background()

    title7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title7_frame = title7.text_frame
    title7_para = title7_frame.paragraphs[0]
    title7_para.text = "06. 2026 실행 로드맵 (Q1~Q4)"
    title7_para.font.size = Pt(28)
    title7_para.font.bold = True
    title7_para.font.color.rgb = DARK_BLUE

    # Q1~Q4 타임라인
    quarters = [
        {"q": "Q1", "color": PRIMARY_BLUE, "progress": 100, "tasks": "시스템 개발 (모바일 앱, 사진 공유, TOOL)"},
        {"q": "Q2", "color": GREEN, "progress": 75, "tasks": "파일럿 운영 (테스트 LINE, 피드백 반영)"},
        {"q": "Q3", "color": ORANGE, "progress": 50, "tasks": "전사 확대 (전체 LINE 적용, 교육 실시)"},
        {"q": "Q4", "color": RED, "progress": 25, "tasks": "성과 평가 (목표 달성 점검, 고도화)"}
    ]

    q_y = Inches(1.5)
    for i, qt in enumerate(quarters):
        current_y = q_y + i * Inches(1.2)

        # Q 라벨
        q_label = slide7.shapes.add_textbox(Inches(0.5), current_y, Inches(0.8), Inches(0.5))
        ql_frame = q_label.text_frame
        ql_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        ql_para = ql_frame.paragraphs[0]
        ql_para.text = qt["q"]
        ql_para.font.size = Pt(24)
        ql_para.font.bold = True
        ql_para.font.color.rgb = qt["color"]
        ql_para.alignment = PP_ALIGN.CENTER

        # 프로그레스 바
        create_progress_bar(slide7, Inches(1.5), current_y + Inches(0.125), Inches(5), Inches(0.25), qt["progress"], qt["color"], LIGHT_GRAY)

        # 태스크
        task_box = slide7.shapes.add_textbox(Inches(6.7), current_y, Inches(2.8), Inches(0.5))
        t_frame = task_box.text_frame
        t_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        t_frame.word_wrap = True
        t_para = t_frame.paragraphs[0]
        t_para.text = qt["tasks"]
        t_para.font.size = Pt(12)
        t_para.font.color.rgb = DARK_TEXT

    page7 = slide7.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page7.text_frame.text = "07"
    page7.text_frame.paragraphs[0].font.size = Pt(11)
    page7.text_frame.paragraphs[0].font.color.rgb = GRAY
    page7.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # ========== 슬라이드 8: 종합 KPI 및 기대효과 ==========
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    slide8.background.fill.solid()
    slide8.background.fill.fore_color.rgb = WHITE

    top_bar8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    top_bar8.fill.solid()
    top_bar8.fill.fore_color.rgb = GOLD
    top_bar8.line.fill.background()

    title8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    title8_frame = title8.text_frame
    title8_para = title8_frame.paragraphs[0]
    title8_para.text = "07. 종합 KPI 및 기대효과"
    title8_para.font.size = Pt(28)
    title8_para.font.bold = True
    title8_para.font.color.rgb = DARK_BLUE

    # 3개 원형 게이지
    kpis_final = [
        {"label": "가공비 절감", "past": "미달", "target": 10, "color": PRIMARY_BLUE},
        {"label": "품질 개선", "past": "부분 달성", "target": 10, "color": GREEN},
        {"label": "손실시간 감소", "past": "미관리", "target": 5, "color": ORANGE}
    ]

    gauge_y = Inches(1.5)
    gauge_spacing = Inches(3.1)
    gauge_start = Inches(0.8)

    for i, kpi in enumerate(kpis_final):
        x_pos = gauge_start + i * gauge_spacing

        # 라벨
        label_box = slide8.shapes.add_textbox(x_pos, gauge_y, Inches(2.5), Inches(0.3))
        l_frame = label_box.text_frame
        l_para = l_frame.paragraphs[0]
        l_para.text = kpi["label"]
        l_para.font.size = Pt(16)
        l_para.font.bold = True
        l_para.font.color.rgb = DARK_BLUE
        l_para.alignment = PP_ALIGN.CENTER

        # 과거 실적
        past_box = slide8.shapes.add_textbox(x_pos, gauge_y + Inches(0.4), Inches(2.5), Inches(0.25))
        p_frame = past_box.text_frame
        p_para = p_frame.paragraphs[0]
        p_para.text = f"과거: {kpi['past']}"
        p_para.font.size = Pt(11)
        p_para.font.color.rgb = GRAY
        p_para.alignment = PP_ALIGN.CENTER

        # 원형 게이지
        gauge_x = x_pos + Inches(1.25)
        gauge_cy = gauge_y + Inches(1.4)
        create_circular_gauge(slide8, gauge_x, gauge_cy, Inches(0.8), kpi["target"] * 10, kpi["color"])

        # 2026 목표
        target_box = slide8.shapes.add_textbox(x_pos, gauge_y + Inches(2.4), Inches(2.5), Inches(0.3))
        tg_frame = target_box.text_frame
        tg_para = tg_frame.paragraphs[0]
        tg_para.text = f"2026 목표: {kpi['target']}%"
        tg_para.font.size = Pt(14)
        tg_para.font.bold = True
        tg_para.font.color.rgb = kpi["color"]
        tg_para.alignment = PP_ALIGN.CENTER

    # 하단: 최종 기대효과
    effect_y = Inches(4.5)

    effect_title = slide8.shapes.add_textbox(Inches(0.5), effect_y, Inches(9), Inches(0.4))
    et_frame = effect_title.text_frame
    et_para = et_frame.paragraphs[0]
    et_para.text = "🎯 최종 기대효과"
    et_para.font.size = Pt(22)
    et_para.font.bold = True
    et_para.font.color.rgb = DARK_BLUE
    et_para.alignment = PP_ALIGN.CENTER

    effect_box = create_premium_box(slide8, Inches(0.5), effect_y + Inches(0.6), Inches(9), Inches(2), NAVY, shadow=True)

    ef_frame = effect_box.text_frame
    ef_frame.margin_left = Inches(0.4)
    ef_frame.margin_top = Inches(0.3)

    effects = [
        "기회손실 최소화 (실시간 대응으로 손실 즉시 차단)",
        "설비 CAPA 증가 (C/T 단축 및 효율 향상)",
        "OVERTIME 감소 (자동화로 인력 부담 감소)",
        "점당 가공비 직접 절감 (종합 효과로 원가 절감)",
        "OH(간접비) 목표 달성 (전사 목표 기여)"
    ]

    for i, eff in enumerate(effects):
        if i == 0:
            para = ef_frame.paragraphs[0]
        else:
            para = ef_frame.add_paragraph()
        para.text = f"▪ {eff}"
        para.font.size = Pt(15)
        para.font.color.rgb = WHITE
        para.space_after = Pt(10)

    page8 = slide8.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.5), Inches(0.3))
    page8.text_frame.text = "08"
    page8.text_frame.paragraphs[0].font.size = Pt(11)
    page8.text_frame.paragraphs[0].font.color.rgb = GRAY
    page8.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # 저장
    output_file = '최종_통합_2026전략_완결판.pptx'
    prs.save(output_file)
    print(f"✅ 최종 통합 PPT 생성 완료: {output_file}")
    print(f"📄 총 8페이지")
    print(f"🎨 특징:")
    print(f"   - 과거 전략 회고 포함 (신규)")
    print(f"   - 차별화 포인트 강조")
    print(f"   - Before/After 비교 시각화")
    print(f"   - 프로그레스 바 + 원형 게이지")
    print(f"   - 프리미엄 디자인 + 참고디자인 스타일 통합")
    return output_file

if __name__ == "__main__":
    create_final_integrated_presentation()
